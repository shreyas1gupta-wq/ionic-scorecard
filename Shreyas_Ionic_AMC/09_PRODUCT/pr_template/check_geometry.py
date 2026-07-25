# -*- coding: utf-8 -*-
"""check_geometry.py — deterministic layout QA for the v9 decks.
Flags: (1) TEXT-vs-TEXT box collisions (>18% of smaller box), (2) shapes past the right/bottom
edge, (3) text below the footer line (y>7.05in) that isn't the footer itself, (4) probable
text overflow of its own box (char-area heuristic), (5) picture-over-text collisions.
Usage: python check_geometry.py out/DECK.pptx [--json out.json]
"""
import sys, json
from pptx import Presentation
from pptx.util import Emu

IN = 914400.0
SW, SH = 13.333, 7.5


def box(sh):
    try:
        return (sh.left / IN, sh.top / IN, (sh.left + sh.width) / IN, (sh.top + sh.height) / IN)
    except Exception:
        return None


def inter(a, b):
    x0, y0 = max(a[0], b[0]), max(a[1], b[1])
    x1, y1 = min(a[2], b[2]), min(a[3], b[3])
    if x1 <= x0 or y1 <= y0: return 0.0
    return (x1 - x0) * (y1 - y0)


def area(a): return max(0.0, (a[2] - a[0]) * (a[3] - a[1]))


def txt_of(sh):
    try:
        return sh.text_frame.text.strip() if sh.has_text_frame else ""
    except Exception:
        return ""


def est_text_height(sh):
    """Rough: sum over paragraphs of ceil(chars*char_w/box_w)*line_h."""
    b = box(sh)
    if not b: return 0
    w = max(b[2] - b[0], 0.1)
    h = 0.0
    for p in sh.text_frame.paragraphs:
        chars = sum(len(r.text or "") for r in p.runs)
        if chars == 0: h += 0.05; continue
        sz = max((r.font.size.pt if r.font.size else 10) for r in p.runs)
        char_w = sz * 0.0075  # ~avg char width inches for Bahnschrift/Georgia
        line_h = sz / 72.0 * 1.18
        import math
        lines = max(1, math.ceil(chars * char_w / w))
        h += lines * line_h
    return h


def check(path):
    prs = Presentation(path)
    findings = []
    for i, s in enumerate(prs.slides, 1):
        shapes = [sh for sh in s.shapes]
        texts = [(sh, box(sh), txt_of(sh)) for sh in shapes if sh.has_text_frame and txt_of(sh)]
        pics = [(sh, box(sh)) for sh in shapes if sh.shape_type == 13]
        # 2/3: bounds
        for sh, b, t in texts:
            if not b: continue
            if b[2] > SW + 0.05 or b[3] > SH + 0.05:
                findings.append({"slide": i, "kind": "off-slide", "text": t[:60], "box": [round(v, 2) for v in b]})
            elif b[1] > 7.02 and "portfolio review" not in t.lower() and len(t) > 30:
                findings.append({"slide": i, "kind": "below-footer", "text": t[:60], "box": [round(v, 2) for v in b]})
        # 1: text-text collisions (ignore tiny labels over big canvases; both must be >12 chars)
        for a in range(len(texts)):
            for b_ in range(a + 1, len(texts)):
                sa, ba, ta = texts[a]; sb, bb, tb = texts[b_]
                if not ba or not bb or len(ta) < 13 or len(tb) < 13: continue
                ov = inter(ba, bb); small = min(area(ba), area(bb))
                if small > 0 and ov / small > 0.18:
                    findings.append({"slide": i, "kind": "text-collision", "pct": round(ov / small * 100),
                                     "a": ta[:42], "b": tb[:42]})
        # 5: picture over text (>25% of the text box)
        for psh, pb in pics:
            if not pb: continue
            for sh, tb_, t in texts:
                if not tb_ or len(t) < 13: continue
                ov = inter(pb, tb_)
                if area(tb_) > 0 and ov / area(tb_) > 0.25:
                    findings.append({"slide": i, "kind": "pic-over-text", "text": t[:42],
                                     "pct": round(ov / area(tb_) * 100)})
        # 4: text overflow of own box (estimate 25% slack)
        for sh, b, t in texts:
            if not b or len(t) < 60: continue
            est = est_text_height(sh); bh = b[3] - b[1]
            if est > bh * 1.30 and bh > 0.15:
                findings.append({"slide": i, "kind": "possible-overflow", "text": t[:50],
                                 "box_h": round(bh, 2), "est_h": round(est, 2)})
    return findings


if __name__ == "__main__":
    path = sys.argv[1]
    fs = check(path)
    by = {}
    for f in fs: by.setdefault(f["kind"], []).append(f)
    print(f"{path}: {len(fs)} findings")
    for k, v in sorted(by.items()):
        print(f"  {k}: {len(v)}")
        for f in v[:14]:
            print("   ", {kk: vv for kk, vv in f.items() if kk != 'kind'})
    if "--json" in sys.argv:
        out = sys.argv[sys.argv.index("--json") + 1]
        json.dump(fs, open(out, "w", encoding="utf-8"), indent=1)
        print("saved ->", out)
