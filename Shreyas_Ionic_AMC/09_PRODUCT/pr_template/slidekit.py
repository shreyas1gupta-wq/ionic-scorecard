# -*- coding: utf-8 -*-
"""slidekit.py, importable Ionic Wealth pptx toolkit for the v9 template (ported from build_pr_full.py).
A Deck wraps one Presentation and exposes the house primitives as methods. Modules call:
    def render(deck, ctx, tier): s = deck.content(...); deck.txt(s, ...); ...
Run styling identical to the v8 deck (indigo #1B27A3 / orange #F2A93C, Bahnschrift head / Georgia body).

Text 'paras' format: a paragraph is a list of runs; a run is a tuple
    (text, font_name, size_pt, color, bold[, italic[, letter_spacing]]).
Pass a single paragraph as a list of runs; pass multiple as a list of paragraphs.
"""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
try:
    from PIL import Image as _PILImage
except Exception:
    _PILImage = None

# ---- brand ----
NAVYD = RGBColor(0x10, 0x19, 0x7A); NAVY = RGBColor(0x1B, 0x27, 0xA3)
NT1 = RGBColor(0x4A, 0x57, 0xC4); NT2 = RGBColor(0x8C, 0x95, 0xDE); NT3 = RGBColor(0xC9, 0xCE, 0xF0)
GOLD = RGBColor(0xF2, 0xA9, 0x3C); ORANGE = GOLD
INK = RGBColor(0x16, 0x23, 0x3B); SLATE = RGBColor(0x6B, 0x72, 0x80)
HAIR = RGBColor(0xE5, 0xE7, 0xEB); TRACK = RGBColor(0xEE, 0xEF, 0xF7)
PANEL = RGBColor(0xF5, 0xF6, 0xFC); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SELL = RGBColor(0xE0, 0x40, 0x2F); SELLBG = RGBColor(0xFB, 0xE3, 0xE0)
HOLD = RGBColor(0x1E, 0x9E, 0x6A); HOLDBG = RGBColor(0xE0, 0xF2, 0xEA)
AMBER = RGBColor(0x92, 0x40, 0x0E); AMBERBG = RGBColor(0xFB, 0xEF, 0xDC)
SERIF = "Georgia"; SANS = "Bahnschrift"
CW, CH = 13.333, 7.5
ML, MR = 0.92, 0.92
RX = CW - MR; UW = RX - ML

# hollow intensifiers never reach a client slide, wherever they sit in the sentence
# (mid-word matches like 'ingenuine' are protected by the leading whitespace requirement)
_TELL_RE = re.compile(r"\s+(?:genuinely|genuine|truly)(?=[\s,.;:…)]|$)")
# internal epistemic tags (D-035 keeps them in research FILES; they never render client-side)
_TAG_RE = re.compile(r"\[\s*(?:OPINION|INFERENCE|DATA|ESTIMATE)\b[^\]]*\]\s*")


def short_name(name, n=28):
    """Shorten a scheme/holding name without ever cutting mid-word: strip common
    suffixes, then drop trailing words until it fits — a clean shorter name beats
    'LIC MF Balanced Advanta…' (declutter pass, 2026-07-25)."""
    name = (name or "").replace(" Fund", "").replace(" (Regular)", " (Reg)").replace(" (Direct)", " (Dir)")
    if len(name) <= n:
        return name
    words = name.split(" ")
    while len(words) > 2 and len(" ".join(words)) > n:
        words.pop()
    out = " ".join(words)
    return out if len(out) <= n else out[:n - 1].rsplit(" ", 1)[0]


def clip_clause(txt, n):
    """Clip long data text for a table cell WITHOUT the broken look. Rules learned the
    hard way (visual QA 2026-07-25): only a real sentence/semicolon boundary may end
    with a period (a comma-cut fakes completeness: 'India's best-capitalised.');
    never leave an unbalanced '('; otherwise word boundary + ellipsis."""
    txt = (txt or "").strip()
    if len(txt) <= n:
        return txt
    cut = txt[:n]
    best = max(cut.rfind(". "), cut.rfind("; "))
    if best >= n * 0.45:
        return cut[:best].rstrip(" ,;") + "."
    if cut.count("(") > cut.count(")"):        # cut landed inside a parenthetical
        cut = cut[:cut.rfind("(")]
    sp = cut.rfind(" ")
    if sp > n * 0.5:
        cut = cut[:sp]
    return cut.rstrip(" ,.;:-") + "…"


def clip_sentences(txt, n):
    """Trim analyst prose to whole SENTENCES within the budget — the text always ends
    with its own full stop, never a mid-clause '…' or doubled punctuation. Falls back
    to clip_clause only when even the first sentence exceeds the budget."""
    txt = (txt or "").strip()
    if len(txt) <= n:
        return txt
    # split at whitespace AFTER sentence punctuation — decimals ('1.5x') never split,
    # and nothing is silently skipped (the [^.]*\. approach dropped text before the
    # first decimal point: a card once rendered starting mid-sentence at '5x across…')
    out = ""
    for part in re.split(r"(?<=[.!?])\s+", txt):
        cand = (out + " " + part).strip()
        if len(cand) > n:
            break
        out = cand
    if out and len(out) >= n * 0.45 and out.endswith((".", "!", "?")):
        return out
    return clip_clause(txt, n)

REC_STYLE = {"Sell": (SELLBG, SELL), "Exit": (SELLBG, SELL), "Redeem-to-Direct": (AMBERBG, AMBER),
             "Redeem": (AMBERBG, AMBER), "Switch": (AMBERBG, AMBER), "Trim": (AMBERBG, AMBER),
             "Hold": (HOLDBG, HOLD), "Aligned": (HOLDBG, HOLD), "Watch": (PANEL, SLATE),
             "Gap": (SELLBG, SELL), "Breach": (SELLBG, SELL)}


class Deck:
    def __init__(self, logo_path=None):
        self.prs = Presentation()
        self.prs.slide_width = Inches(CW); self.prs.slide_height = Inches(CH)
        self.BLANK = self.prs.slide_layouts[6]
        self.logo_path = logo_path if (logo_path and os.path.exists(logo_path)) else None
        self.folio = 0
        self._anchors = {}    # key -> (prio, slide, folio)
        self._links = []      # (shape, key)
        self._pagerefs = []   # (run, key)

    def save(self, path):
        self.resolve_links()
        self.prs.save(path); return path

    # ---------- internal links (annexure cross-references, real-deck 'see p.NN') ----------
    def anchor(self, key, s, prio=0):
        """Register slide s as the jump target for key. Higher prio wins (a Sell card
        outranks the all-holdings page the same name also appears on)."""
        cur = self._anchors.get(key)
        if cur is None or prio >= cur[0]:
            self._anchors[key] = (prio, s, self.folio)

    def link(self, shape, key, own=None):
        """Make an existing shape jump to anchor(key) once resolve_links() runs.
        own: the shape's slide — a link resolving to its own slide is dropped."""
        self._links.append((shape, key, own))

    def hotspot(self, s, x, y, w, h, key):
        """Invisible click area (0%-alpha fill so the whole area takes the click)."""
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(max(w, 0.001)), Inches(max(h, 0.001)))
        shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
        try:
            from pptx.oxml.ns import qn
            clr = shp.fill._xPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
            clr.append(clr.makeelement(qn('a:alpha'), {'val': '0'}))
        except Exception:
            pass
        shp.line.fill.background(); self._nosh(shp)
        self._links.append((shp, key, s))
        return shp

    def pageref(self, s, x, y, key, w=0.55, color=None, align=PP_ALIGN.RIGHT, label=None):
        """Small 'p.NN' cross-reference (optionally 'LABEL · p.NN'); the number is
        patched at resolve time and the whole box clicks through to the anchor."""
        runs = [(label + "  ·  ", SANS, 7.5, SLATE, True, False, 60)] if label else []
        runs.append(("p.00", SANS, 8, color or NT2, True))
        tb = self.txt(s, x, y, w, 0.2, [runs], align=align, wrap=False)
        self._pagerefs.append((tb, key))
        self._links.append((tb, key, None))
        return tb

    def resolve_links(self):
        """Bind every registered link/pageref to its anchor slide. Called by save()."""
        for tb, key in self._pagerefs:
            tgt = self._anchors.get(key)
            runs = tb.text_frame.paragraphs[0].runs
            if tgt:
                runs[-1].text = f"p.{tgt[2]:02d}"
            else:                       # dead ref: blank the whole 'LABEL · p.NN' box
                for r in runs:
                    r.text = ""
        n = 0
        for shape, key, own in self._links:
            tgt = self._anchors.get(key)
            if tgt is not None and tgt[1] is not own:
                try:
                    shape.click_action.target_slide = tgt[1]; n += 1
                except Exception:
                    pass
        return n

    # ---------- primitives ----------
    @staticmethod
    def _nosh(sh):
        try: sh.shadow.inherit = False
        except Exception: pass

    def slide(self, bg=WHITE):
        s = self.prs.slides.add_slide(self.BLANK)
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); self._nosh(r)
        return s

    def rect(self, s, x, y, w, h, fill=None, line=None, lw=0.75, round_=0.0):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(max(w, 0.001)), Inches(max(h, 0.001)))
        if fill is None: shp.fill.background()
        else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
        if line is None: shp.line.fill.background()
        else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
        self._nosh(shp)
        if round_:
            try: shp.adjustments[0] = round_
            except Exception: pass
        return shp

    def oval(self, s, x, y, d, fill):
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        o.fill.solid(); o.fill.fore_color.rgb = fill; o.line.fill.background(); self._nosh(o)
        return o

    def rule(self, s, x, y, w, color=HAIR, h=0.01): return self.rect(s, x, y, w, h, fill=color)
    def vrule(self, s, x, y, h, color=HAIR, w=0.008): return self.rect(s, x, y, w, h, fill=color)

    def txt(self, s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, ls=None):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"): setattr(tf, m, 0)
        if paras and isinstance(paras[0], tuple): paras = [paras]
        for pi, para in enumerate(paras):
            p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
            p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(0)
            if ls: p.line_spacing = ls
            for run in para:
                t, fn, sz, col, bold = run[0], run[1], run[2], run[3], run[4]
                ital = run[5] if len(run) > 5 else False
                spc = run[6] if len(run) > 6 else None
                r = p.add_run()
                # house no-AI-tell rule (detell-lite): dashes and hollow intensifiers never
                # reach a client slide, even when they arrive inside analyst-data strings
                if isinstance(t, str):
                    t = t.replace(" — ", ", ").replace("—", ", ").replace(" -- ", ", ")
                    # 'genuine(ly)' swaps to a plain word (bare strip dangles articles:
                    # 'a genuine, company-acknowledged' -> 'a,'); 'truly' just goes
                    t = t.replace(" genuinely", " clearly").replace(" genuine", " clear")
                    t = _TELL_RE.sub("", t)
                    t = _TAG_RE.sub("", t)
                    # demo-name consistency + jargon softening + data-engine narration
                    # never reaches a client slide (leak audit 2026-07-26; widened after the
                    # CEO sweep caught snake_case fields + data-QA vocabulary in rationale text)
                    t = (t.replace("AZBY", "ABXY")
                          .replace("Forensic / governance flag", "Governance concern")
                          .replace("our own PIT data", "our data").replace("PIT data", "our data")
                          .replace("quant data cutoff", "close of our scoring window")
                          .replace("quant data cut", "close of our scoring window")
                          .replace("the quant feed", "our screening data")
                          .replace("our data feed", "our screening data")
                          .replace("quant snapshot", "screening snapshot")
                          .replace("data snapshot", "screening snapshot")
                          .replace("fcf_yield", "FCF yield"))
                    t = re.sub(r"\s*\(\d+\s*rows?\b[^)]{0,60}\)", "", t)
                    # glyphs Bahnschrift lacks (render as tofu in charts/PDF): never ship
                    t = (t.replace(" → ", " to ").replace("→", "to")
                          .replace("≤ ", "max ").replace("≤", "max ")
                          .replace("≥ ", "min ").replace("≥", "min "))
                r.text = t
                r.font.name = fn; r.font.size = Pt(sz); r.font.bold = bold
                r.font.italic = ital; r.font.color.rgb = col
                if spc is not None: r.font._rPr.set('spc', str(int(spc)))
        return tb

    def pic(self, s, path, x, y, boxw, boxh, valign="middle", halign="center"):
        """Fit image into box preserving aspect."""
        if not path or not os.path.exists(path):
            return self.rect(s, x, y, boxw, boxh, fill=PANEL, line=HAIR)
        try:
            iw, ih = _PILImage.open(path).size; ar = iw / ih
        except Exception:
            return s.shapes.add_picture(path, Inches(x), Inches(y), Inches(boxw), Inches(boxh))
        if boxw / boxh > ar:
            h = boxh; w = boxh * ar
        else:
            w = boxw; h = boxw / ar
        ox = x + {"center": (boxw - w) / 2, "left": 0, "right": boxw - w}.get(halign, 0)
        oy = y + {"middle": (boxh - h) / 2, "top": 0, "bottom": boxh - h}.get(valign, 0)
        return s.shapes.add_picture(path, Inches(ox), Inches(oy), Inches(w), Inches(h))

    def pill(self, s, x, y, text, w=0.9, kind=None):
        fill, tc = REC_STYLE.get(kind or text, (PANEL, SLATE))
        self.rect(s, x, y, w, 0.24, fill=fill, line=tc, lw=0.75, round_=0.5)
        self.txt(s, x, y - 0.006, w, 0.25, [(text.upper(), SANS, 8, tc, True, False, 60)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---------- chrome ----------
    def logo(self, s):
        if self.logo_path: self.pic(s, self.logo_path, RX - 1.85, 0.40, 1.8, 0.31, halign="right")
        else: self.txt(s, RX - 2.3, 0.42, 2.3, 0.3, [("IONIC WEALTH", SANS, 12, NAVY, True, False, 60)], align=PP_ALIGN.RIGHT)

    def classified(self, s, dark=False):
        # client deliverable ⇒ 'Private & Confidential', never an internal security marking
        # (leak audit 2026-07-26)
        self.txt(s, CW / 2 - 1.7, 0.14, 3.4, 0.2, [("Private & Confidential", SANS, 7.5, (NT2 if dark else SLATE), False)], align=PP_ALIGN.CENTER)

    def footer(self, s, dark=False):
        self.txt(s, RX - 2.4, 7.14, 2.4, 0.2, [(f"Portfolio Review  ·  {self.folio:02d}", SANS, 7.5, (NT2 if dark else SLATE), False)], align=PP_ALIGN.RIGHT)

    def source(self, s, text):
        self.txt(s, ML, 6.66, UW, 0.24, [(text, SANS, 7, SLATE, False)])

    def marker(self, s, section_no, section_name):
        if not section_name: return
        self.txt(s, 8.6, 0.30, RX - 8.6 - 1.9, 0.22, [(f"{section_no:02d} · {section_name.upper()}", SANS, 7.5, NT2, True, False, 120)], align=PP_ALIGN.RIGHT)
        # section progress ticks (v8 pattern), clear of the logo
        for k in range(5):
            col = NAVY if (k + 1) == section_no else HAIR
            self.rect(s, 9.35 + 0.20 * k, 0.56, 0.15, 0.022, fill=col)

    @staticmethod
    def _fit(text, base_pt, box_in, char_w_per_pt=0.0079, floor=14):
        """Shrink a header font so the line NEVER wraps into content below it."""
        need = len(text) * base_pt * char_w_per_pt
        if need <= box_in:
            return base_pt
        return max(floor, int(base_pt * box_in / need))

    def content(self, section_no, section_name, eyebrow, title, standfirst=None):
        """Standard content-slide header. Returns the slide. Increments folio.
        standfirst: one Georgia-italic thesis line under the rule (v8 pattern).
        Eyebrow/title fonts auto-shrink to one line — headers can never wrap into the body."""
        self.folio += 1
        s = self.slide(WHITE)
        self.logo(s); self.classified(s)
        if section_name: self.marker(s, section_no, section_name)
        self.txt(s, ML, 0.46, 7.5, 0.55, [(eyebrow, SANS, self._fit(eyebrow, 26, 7.4), NAVY, True)])
        self.txt(s, ML, 1.04, 10.6, 0.45, [(title, SANS, self._fit(title, 17.5, 10.5), ORANGE, True)])
        self.rule(s, ML, 1.54, UW, NAVY, 0.024)
        if standfirst:
            self.txt(s, ML, 1.60, UW, 0.24, [(standfirst, SERIF, 11, SLATE, False, True)])
        self.footer(s)
        return s

    def section_divider(self, num, title, subtitle="", pages=None):
        """v8 editorial divider: NAVYD ground, Georgia ghost numeral, gold kicker, serif
        title — plus a low-alpha flow-art field on the right so the page reads designed,
        not empty (Principal 2026-07-25)."""
        self.folio += 1
        GHOST = RGBColor(0x24, 0x2F, 0x8E)
        s = self.slide(NAVYD)
        try:
            import art as _art
            self.pic(s, _art.flow_art(f"div_art_{num}", w=6.4, h=7.5, seed=31 + num,
                                      gold=False, alpha=0.5, transparent=True),
                     CW - 6.4, 0, 6.4, 7.5, valign="top", halign="right")
        except Exception:
            pass
        self.classified(s, dark=True)
        self.txt(s, 7.4, 0.95, RX - 7.4, 2.6, [(f"{num:02d}", SERIF, 170, GHOST, True)], align=PP_ALIGN.RIGHT)
        self.txt(s, ML, 2.58, 6.5, 0.3, [(f"SECTION {num:02d}", SANS, 10, GOLD, True, False, 250)])
        self.txt(s, ML, 2.96, 8.6, 0.75, [(title, SERIF, 30, WHITE, True)])
        self.rule(s, ML, 3.74, 2.2, GOLD, 0.028)
        if subtitle:
            self.txt(s, ML, 3.98, 8.6, 0.6, [(subtitle, SERIF, 13, NT3, False, True)])
        if pages:
            y = 5.45
            for label in pages[:5]:
                self.txt(s, ML, y, 8.0, 0.24, [(label, SANS, 9.5, NT3, False, False, 40)])
                y += 0.30
        return s

    def pullquote(self, s, x, y, quote, attribution, w=5.95):
        """v8 spotlight pull-quote: gold bar + Georgia-italic quote + spaced-caps attribution."""
        self.rect(s, x, y, 0.03, 2.35, fill=GOLD)
        self.txt(s, x + 0.24, y, w, 2.0, [(quote, SERIF, 18, INK, False, True)], ls=1.16)
        self.txt(s, x + 0.24, y + 2.12, w, 0.24, [(attribution.upper(), SANS, 8.5, SLATE, True, False, 150)])

    # ---------- components ----------
    def kpi_strip(self, s, stats, y=1.82, x=ML, w=UW, h=0.95):
        n = len(stats); cw = w / n
        for i, st in enumerate(stats):
            cx = x + i * cw
            val, lab = st[0], st[1]; sub = st[2] if len(st) > 2 else None
            col = st[3] if len(st) > 3 else INK
            self.txt(s, cx, y, cw - 0.15, 0.5, [(val, SANS, 27, col, False)])
            self.txt(s, cx, y + 0.52, cw - 0.15, 0.19, [(lab.upper(), SANS, 8, SLATE, True, False, 200)], ls=1.0)
            if sub: self.txt(s, cx, y + 0.72, cw - 0.15, 0.2, [(sub, SANS, 8, NT2, False)])
            if i: self.vrule(s, cx - 0.06, y + 0.05, 0.62, HAIR, 0.008)

    def score_bar(self, s, x, y, score, w=0.65):
        # v8 law: neutral NT1 fill + ink 40-tick — the Call pill carries the colour, not the bar
        self.rect(s, x, y, w, 0.07, fill=TRACK)
        if score is not None:
            self.rect(s, x, y, w * max(score, 0) / 100.0, 0.07, fill=NT1)
        self.rect(s, x + w * 0.40, y - 0.03, 0.012, 0.13, fill=INK)
        self.txt(s, x + w + 0.07, y - 0.085, 0.45, 0.24, [((f"{score:.0f}" if score is not None else "-"), SANS, 9, INK, False)], anchor=MSO_ANCHOR.MIDDLE)

    def callout(self, s, x, y, w, h, title, body, kind="note"):
        """Boxed callout. kind: note|warn|good|human."""
        bg, bar, tc = {"note": (PANEL, NAVY, NAVY), "warn": (SELLBG, SELL, SELL),
                       "good": (HOLDBG, HOLD, HOLD), "human": (AMBERBG, GOLD, AMBER)}.get(kind, (PANEL, NAVY, NAVY))
        self.rect(s, x, y, w, h, fill=bg, round_=0.04)
        self.rect(s, x, y, 0.06, h, fill=bar)
        self.txt(s, x + 0.22, y + 0.14, w - 0.4, 0.3, [(title.upper(), SANS, 9.5, tc, True, False, 60)])
        self.txt(s, x + 0.22, y + 0.44, w - 0.4, h - 0.5, [(body, SERIF, 10.5, INK, False)], ls=1.06)

    def callout_h(self, w, body, min_h=1.0, max_h=2.6):
        """Height that hugs the text for a callout of width w — a box sized to its
        worst case renders 40-60% empty tint on short copy (declutter, 2026-07-25)."""
        import math as _m
        cpl = max(10, int((w - 0.44) / (0.0102 * 10.5)))
        lines = max(1, _m.ceil(len(body or "") / cpl))
        return min(max_h, max(min_h, 0.62 + lines * 0.185))

    def table(self, s, x, y, w, cols, rows, rowh=0.34, fs=10, hfs=8, header=True, zebra=False,
              maxrows=None, totals=None):
        """cols = [(label, width_frac, align 'l'/'c'/'r')]. Each cell:
           str (Georgia serif, v8 register) | ('b', text[, color]) bold sans | ('c', text, color[, bold])
           | ('pill', text, kind) | ('bar', score) | ('flags', [str,...]).
           totals: optional footer row (same cell forms) under a navy rule (v8 utbl pattern).
           Returns y after the table."""
        tot = sum(c[1] for c in cols); xs = []; ws = []; cx = x
        for (_, cwf, _a) in cols:
            xs.append(cx); ws.append(w * cwf / tot); cx += w * cwf / tot
        PAD = 0.08
        def al(i): return PP_ALIGN.RIGHT if cols[i][2] == "r" else (PP_ALIGN.CENTER if cols[i][2] == "c" else PP_ALIGN.LEFT)
        def cell_out(cell, i, ry, fsz, force_bold=False):
            cx = xs[i] + PAD; cwid = ws[i] - 2 * PAD
            if isinstance(cell, tuple):
                k = cell[0]
                if k == "pill":
                    self.pill(s, cx, ry + rowh / 2 - 0.13, cell[1], w=min(cwid, 1.4), kind=cell[2] if len(cell) > 2 else cell[1])
                elif k == "bar":
                    self.score_bar(s, cx, ry + rowh / 2 - 0.02, cell[1], w=min(cwid - 0.5, 0.7))
                elif k == "flags":
                    fx = cx
                    for fl in cell[1][:3]:
                        self.pill(s, fx, ry + rowh / 2 - 0.13, fl[:9], w=0.82,
                                  kind="Sell" if fl in ("NEG_ALPHA", "CLOSET_INDEX", "DEEP_DD", "CAPACITY",
                                                        "TRAILS", "INDEX HUG", "DEEP FALL", "TOO LARGE") else "Trim")
                        fx += 0.88
                elif k in ("b", "c"):
                    col = cell[2] if (k == "c" and len(cell) > 2) else INK
                    bold = (k == "b") or (len(cell) > 3 and cell[3]) or force_bold
                    self.txt(s, cx, ry, cwid, rowh, [(str(cell[1]), SANS, fsz, col, bold)], align=al(i), anchor=MSO_ANCHOR.MIDDLE)
            else:
                # plain cells in Georgia — the private-bank serif/sans register (v8)
                self.txt(s, cx, ry, cwid, rowh, [(str(cell), SERIF, fsz, INK, force_bold)], align=al(i), anchor=MSO_ANCHOR.MIDDLE)
        if header:
            for i, (lab, _c, _a) in enumerate(cols):
                self.txt(s, xs[i] + PAD, y, ws[i] - 2 * PAD, 0.24, [(lab.upper(), SANS, hfs, SLATE, True, False, 200)], align=al(i), anchor=MSO_ANCHOR.MIDDLE)
            self.rule(s, x, y + 0.28, w, NAVY, 0.015); ry = y + 0.33
        else:
            ry = y
        rows = rows[:maxrows] if maxrows else rows
        for ri, row in enumerate(rows):
            if zebra and ri % 2 == 1:
                self.rect(s, x, ry - 0.02, w, rowh, fill=PANEL)
            for i, cell in enumerate(row):
                cell_out(cell, i, ry, fs)
            if not zebra:
                self.rule(s, x, ry + rowh - 0.02, w, HAIR, 0.006)
            ry += rowh
        if totals:
            self.rule(s, x, ry, w, NAVY, 0.014); ry += 0.05
            for i, cell in enumerate(totals):
                cell_out(cell, i, ry, fs + 0.5, force_bold=True)
            ry += rowh
        return ry

    def scope_tag(self, s, text, x=None, y=1.62):
        """CMP-DATASCOPE tag, always states data scope + as-of date. One line that must
        never spill: over-long scopes drop whole ' · ' segments (keeping the as-of tail)
        rather than truncating mid-word."""
        x = ML if x is None else x
        budget = int((RX - x - 0.55) / (0.0102 * 8.5)) - 8
        if len(text) > budget and " · " in text:
            parts = text.split(" · ")
            while len(parts) > 2 and len(" · ".join(parts)) > budget:
                parts.pop(-2)          # drop detail segments, keep the first + the as-of tail
            text = " · ".join(parts)
        if len(text) > budget:
            text = text[:budget - 1].rstrip(" ·,;") + "…"
        self.rect(s, x, y, 0.14, 0.14, fill=NT2, round_=0.3)
        self.txt(s, x + 0.20, y - 0.03, RX - x - 0.35, 0.22,
                 [("SCOPE  ", SANS, 7.5, SLATE, True, False, 80), (text, SERIF, 8.5, INK, False, True)], wrap=False)

    _BAND_VARIANTS = (
        "The Ionic Score flags candidates; the Portfolio Review team confirms every call.",
        "Scores flag, people decide: each call here carries a named reviewer's sign-off.",
        "The score is the starting point; the final word on every name belongs to the desk.",
        "Numbers shortlist, judgment decides: every call passed a human review.",
    )

    def score_band(self, s, y=6.9):
        """The score-positioning line (F13), rotated across four phrasings so the refrain
        reads authored, not templated. Attaches to score-bearing slides only."""
        line = self._BAND_VARIANTS[self.folio % len(self._BAND_VARIANTS)]
        self.txt(s, ML, y, UW, 0.18, [(line, SERIF, 8.5, SLATE, False, True)])


def new_deck():
    logo = r"C:\Users\SHREYA~1.1GU\AppData\Local\Temp\claude\c--Users-Shreyas-1Gupta-OneDrive---Angel-Broking-Limited-Desktop-Backup-NIFTY-500\5ec2bf16-8c38-4f40-9e4f-8e07be6545fd\scratchpad\assets\logo_clean.png"
    return Deck(logo_path=logo)
