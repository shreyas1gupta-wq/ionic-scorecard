# -*- coding: utf-8 -*-
"""Do the five-signal dots carry COLOUR in the built deck, or are they all hollow?

This is the check that the whole build pipeline cannot do for itself. A failed universe join produces
a structurally perfect page whose every dot is an unfilled grey ring, exit code 0. Run this against
the .pptx, because the .pptx is the thing the client sees.
"""
import glob
import os
import sys
from collections import Counter

from pptx import Presentation
from pptx.util import Emu

MAX_DOT = Emu(int(0.30 * 914400))          # a signal dot is ~0.19-0.24in; anything bigger is art
HOLLOW = "hollow/no-fill"


def fills_on(slide):
    out = Counter()
    for sh in slide.shapes:
        # add_shape(MSO_SHAPE.OVAL) reports shape_type == AUTO_SHAPE; the oval identity is in
        # auto_shape_type. Matching on shape_type finds nothing and looks like "no dots on the page".
        try:
            ast = str(sh.auto_shape_type or "")
        except Exception:
            continue
        if "OVAL" not in ast or sh.width is None or sh.width > MAX_DOT:
            continue
        try:
            out[str(sh.fill.fore_color.rgb)] += 1
        except Exception:
            out[HOLLOW] += 1
    return out


def main():
    files = sorted(glob.glob(os.path.join("out", "*.pptx")))
    if not files:
        print("no decks in out/ -- build first")
        return 1
    bad = 0
    for f in files:
        prs = Presentation(f)
        best = None
        for i, sl in enumerate(prs.slides):
            fl = fills_on(sl)
            n = sum(fl.values())
            if n >= 15 and (best is None or n > sum(best[1].values())):
                best = (i + 1, fl)
        name = os.path.basename(f)
        if best is None:
            print(f"{name:40s}  no dot page (tier may not include book_scored)")
            continue
        pg, fl = best
        n = sum(fl.values())
        hollow = fl.get(HOLLOW, 0)
        distinct = len([k for k in fl if k != HOLLOW])
        # A FEW hollow rings are correct -- a genuinely unscored signal must show as unscored, and a
        # thin-history name legitimately has some. The failure this gate exists to catch is the page
        # where the join produced NOTHING, so demanding zero hollow would cry wolf on a healthy deck.
        if hollow > n * 0.5 or distinct < 2:
            verdict, why = "*** FAIL ***", " <- join looks empty/degenerate"
            bad += 1
        else:
            verdict, why = "OK", ""
        print(f"{name:40s}  p.{pg:<4d} {n:3d} dots  {distinct} colours  "
              f"{hollow} hollow ({hollow / n * 100:.0f}%)   {verdict}{why}")
        for k, v in fl.most_common():
            print(f"{'':40s}     {k:18s} x{v}")
    print()
    print("PASS -- every deck's signal dots carry colour" if not bad else
          f"FAIL -- {bad} deck(s) look like a failed universe join")
    return 1 if bad else 0


if __name__ == "__main__":
    sys.exit(main())
