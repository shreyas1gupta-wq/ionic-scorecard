# -*- coding: utf-8 -*-
"""render_preview.py — approximate raster preview of a python-pptx deck, so slides can be
visually reviewed without PowerPoint. Draws, in z-order: solid-fill autoshapes (rect/round-
rect/oval), pictures (actual embedded blobs), and text frames with the real Bahnschrift/
Georgia fonts, per-run size/bold/italic/color, word-wrap, alignment and vertical anchor.
Faithful enough to judge clutter, density and hierarchy; not a pixel-perfect renderer.

Usage: python render_preview.py <deck.pptx> <outdir> [scale_px_per_inch=140]
"""
import io
import os
import sys
import shutil
import tempfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EMU = 914400.0

_FONT_FILES = {
    ("Bahnschrift", False, False): "bahnschrift.ttf",
    ("Bahnschrift", True, False): "bahnschrift.ttf",   # variable font; PIL picks default weight
    ("Bahnschrift", False, True): "bahnschrift.ttf",
    ("Bahnschrift", True, True): "bahnschrift.ttf",
    ("Georgia", False, False): "georgia.ttf",
    ("Georgia", True, False): "georgiab.ttf",
    ("Georgia", False, True): "georgiai.ttf",
    ("Georgia", True, True): "georgiaz.ttf",
}
_FONT_DIR = r"C:\Windows\Fonts"
_font_cache = {}


def _font(name, size_px, bold, italic):
    key = (name, bold, italic, size_px)
    if key in _font_cache:
        return _font_cache[key]
    fn = _FONT_FILES.get((name, bold, italic)) or _FONT_FILES.get((name, False, False))
    if fn is None:
        fn = "arialbd.ttf" if bold else "arial.ttf"
    try:
        f = ImageFont.truetype(os.path.join(_FONT_DIR, fn), max(size_px, 4))
    except Exception:
        f = ImageFont.load_default()
    # crude bold emulation for Bahnschrift variable font
    _font_cache[key] = f
    return f


def _rgb(color_obj, default=(22, 35, 59)):
    try:
        rgb = color_obj.rgb
        if rgb is None:
            return default
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return default


def _draw_text_frame(draw, sh, S):
    tf = sh.text_frame
    x0 = sh.left / EMU * S
    y0 = sh.top / EMU * S
    w = sh.width / EMU * S
    h = sh.height / EMU * S
    wrap = tf.word_wrap is not False
    anchor = tf.vertical_anchor

    # first pass: build wrapped lines [(runs=[(text,font,color,bold)], line_h)]
    lines = []
    for p in tf.paragraphs:
        align = p.alignment
        ls = p.line_spacing or 1.0
        runs = []
        for r in p.runs:
            t = r.text or ""
            if not t:
                continue
            sz = (r.font.size.pt if r.font.size else 10)
            px = int(sz * S / 72.0)
            f = _font(r.font.name or "Georgia", px, bool(r.font.bold), bool(r.font.italic))
            col = _rgb(r.font.color)
            runs.append([t, f, col, px])
        if not runs:
            lines.append(([], 6, align, 1.0))
            continue
        # wrap: split into words across runs
        words = []
        for t, f, col, px in runs:
            for wd in t.replace("\u000b", "\n ").split(" "):
                words.append((wd, f, col, px))
        cur, cur_w = [], 0.0
        max_px = max(px for _, _, _, px in runs)
        line_h = max_px * 1.22 * ls
        for wd, f, col, px in words:
            piece = wd + " "
            pw = draw.textlength(piece, font=f)
            if wrap and cur and cur_w + pw > w:
                lines.append((cur, line_h, align, ls))
                cur, cur_w = [], 0.0
            cur.append((piece, f, col))
            cur_w += pw
        if cur:
            lines.append((cur, line_h, align, ls))

    total_h = sum(lh for _, lh, _, _ in lines)
    if anchor == MSO_ANCHOR.MIDDLE:
        ty = y0 + max(0, (h - total_h) / 2)
    elif anchor == MSO_ANCHOR.BOTTOM:
        ty = y0 + max(0, h - total_h)
    else:
        ty = y0

    for runs, lh, align, _ls in lines:
        if not runs:
            ty += lh
            continue
        lw = sum(draw.textlength(t, font=f) for t, f, _ in runs)
        if align == PP_ALIGN.CENTER:
            tx = x0 + (w - lw) / 2
        elif align == PP_ALIGN.RIGHT:
            tx = x0 + w - lw
        else:
            tx = x0
        for t, f, col in runs:
            draw.text((tx, ty), t, font=f, fill=col)
            tx += draw.textlength(t, font=f)
        ty += lh


def render(path, outdir, S=140):
    os.makedirs(outdir, exist_ok=True)
    tmp = os.path.join(tempfile.gettempdir(), "_preview_src.pptx")
    shutil.copy2(path, tmp)
    prs = Presentation(tmp)
    W = int(prs.slide_width / EMU * S)
    H = int(prs.slide_height / EMU * S)
    n = 0
    for i, s in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), (255, 255, 255))
        draw = ImageDraw.Draw(img)
        for sh in s.shapes:
            try:
                st = sh.shape_type
                if st == 13:  # picture
                    try:
                        blob = sh.image.blob
                        pic = Image.open(io.BytesIO(blob)).convert("RGBA")
                        bw = max(1, int(sh.width / EMU * S))
                        bh = max(1, int(sh.height / EMU * S))
                        pic = pic.resize((bw, bh))
                        img.paste(pic, (int(sh.left / EMU * S), int(sh.top / EMU * S)), pic)
                    except Exception:
                        pass
                    continue
                if sh.has_text_frame and sh.text_frame.text.strip():
                    # autoshapes may carry both fill and text: fill first
                    pass
                # fills (autoshapes incl. rounded rect / oval)
                if st in (1, 9) or str(st).startswith("AUTO_SHAPE"):
                    try:
                        if sh.fill.type is not None and sh.fill.type == 1:  # solid
                            # skip fully-transparent fills (link hotspots carry alpha=0)
                            if 'alpha val="0"' in sh._element.xml:
                                col = None
                            else:
                                col = _rgb(sh.fill.fore_color, None)
                            if col is not None:
                                x0 = sh.left / EMU * S; y0 = sh.top / EMU * S
                                x1 = x0 + sh.width / EMU * S; y1 = y0 + sh.height / EMU * S
                                if getattr(sh, "shape_type", None) == 9 or "OVAL" in str(getattr(sh, "auto_shape_type", "")):
                                    draw.ellipse([x0, y0, x1, y1], fill=col)
                                else:
                                    draw.rectangle([x0, y0, x1, y1], fill=col)
                    except Exception:
                        pass
                if sh.has_text_frame and sh.text_frame.text.strip():
                    _draw_text_frame(draw, sh, S)
            except Exception:
                continue
        img.save(os.path.join(outdir, f"slide_{i:02d}.png"))
        n += 1
    print(f"rendered {n} slides -> {outdir}")


if __name__ == "__main__":
    src = sys.argv[1]
    out = sys.argv[2]
    scale = int(sys.argv[3]) if len(sys.argv) > 3 else 140
    render(src, out, scale)
