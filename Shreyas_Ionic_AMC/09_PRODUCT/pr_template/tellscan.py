# -*- coding: utf-8 -*-
"""tellscan.py — deterministic client-safe-copy scan for rendered NDPMS decks (standing gate,
alongside check_geometry.py/check_geometry2.py). Codifies the QA LAW #3 tell-scan that was
previously re-derived from memory each session (2026-07-27 permanent fix).

Catches, by bucket: AI writing tells, "Buy" recommendation language, internal jargon/codenames,
data-QA vocabulary, source/analyst citations, raw snake_case field names, and "synthetic/demo"
language outside the annexure. See BUCKETS below for the exact term lists.

Usage:
    python tellscan.py out/DECK.pptx [--json out.json]
    python tellscan.py data/<client>.py               # scan raw ctx source too — the render-time
                                                       # scrub in slidekit.txt() can't rescue whole
                                                       # sentences of internal audit trail, so the
                                                       # DATA file must be clean at the source.
"""
import ast
import json
import re
import sys

from pptx import Presentation

# ---------------------------------------------------------------------------
# Term buckets (SKILL.md QA LAW #3 + the internal-jargon/mislabeling classes
# actually caught during the Anand Reddy HNI_DEEP build, 2026-07-27).
# ---------------------------------------------------------------------------
BUCKETS = {
    "AI_TELL": [
        "--", "genuinely", "genuine", "truly", "robust", "seamless", "holistic",
        "delve", "boasts",
    ],
    "RECOMMENDATION_LANGUAGE": [
        r"\bBuy\b",
    ],
    "INTERNAL_JARGON": [
        "SENTINEL", "QFRA", "MERIT", "pf_qual", "AZBY",
        "quant-only", "analyst view", "Ratified Sell", "Ratified Hold",
        "One-time review", "House decision", "Quant-head research",
    ],
    "DATA_QA_VOCAB": [
        "stale", "does not reconcile", "data feed", "data cut", "quant snapshot",
        "data snapshot", "Data Office", "CoPilot", "DATA_GAPS",
    ],
    "SOURCE_CITATIONS": [
        "screener.in", "MF Dashboard", "NSE bhavcopy", "INDmoney", "Groww",
        "Paytm Money", "Advisorkhoj", "AMFI NAV",
    ],
    "SYNTHETIC_DEMO_LEAK": [
        "synthetic demo", "synthetic book", "synthetic funds", "illustrative for this demo",
    ],
    "GLYPH_HYGIENE": [
        "→", "≤", "≥",  # -> <= >=  (Bahnschrift can't render these)
    ],
    "LITERAL_NONE": [
        "None-year", "built not yet", "None%", "+0.0%",
    ],
}

# em-dash is handled at render time by slidekit.txt() for RENDERED pptx text, but a raw ctx
# .py source file never passes through that scrub -- check it there too.
_EMDASH = "—"

_SNAKE_RE = re.compile(r"\b[a-z]+(?:_[a-z0-9]+)+\b")
# allowlist: legitimate snake_case that isn't an internal field-name leak (file paths, real
# ticker-ish tokens, etc.) -- extend as false positives turn up
_SNAKE_ALLOW = {"pf_qual"}  # already caught by INTERNAL_JARGON; don't double-report


def _bucket_hits(text):
    hits = []
    low = text.lower()
    for bucket, terms in BUCKETS.items():
        for term in terms:
            pat = term if term.startswith(r"\b") else re.escape(term)
            if re.search(pat, text if bucket == "RECOMMENDATION_LANGUAGE" else low,
                         0 if bucket == "RECOMMENDATION_LANGUAGE" else re.IGNORECASE):
                hits.append((bucket, term))
    if _EMDASH in text:
        hits.append(("AI_TELL", "em-dash"))
    for m in _SNAKE_RE.finditer(text):
        tok = m.group(0)
        if tok not in _SNAKE_ALLOW and "_" in tok:
            hits.append(("SNAKE_CASE_FIELD", tok))
    return hits


def extract_texts_pptx(path):
    """Yield (location, text) for every non-empty text_frame in a rendered deck."""
    pres = Presentation(path)
    for i, slide in enumerate(pres.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if t.strip():
                    yield (f"slide {i}", t)


def extract_texts_pyfile(path):
    """Yield (location, text) for every string literal in a ctx-builder .py source file."""
    src = open(path, encoding="utf-8").read()
    tree = ast.parse(src, filename=path)
    for node in ast.walk(tree):
        if isinstance(node, ast.Constant) and isinstance(node.value, str) and len(node.value) > 3:
            yield (f"{path}:{node.lineno}", node.value)


def scan(pairs):
    findings = []
    for loc, text in pairs:
        for bucket, term in _bucket_hits(text):
            findings.append({"loc": loc, "bucket": bucket, "term": term, "text": text[:100]})
    return findings


def main():
    path = sys.argv[1]
    if path.endswith(".py"):
        pairs = list(extract_texts_pyfile(path))
    else:
        pairs = list(extract_texts_pptx(path))
    findings = scan(pairs)

    if "--json" in sys.argv:
        out = sys.argv[sys.argv.index("--json") + 1]
        json.dump(findings, open(out, "w", encoding="utf-8"), indent=2)

    if not findings:
        print(f"{path}: 0 findings")
        return 0

    by_bucket = {}
    for f in findings:
        by_bucket.setdefault(f["bucket"], []).append(f)
    print(f"{path}: {len(findings)} findings")
    for bucket, items in sorted(by_bucket.items(), key=lambda kv: -len(kv[1])):
        print(f"  [{bucket}] x{len(items)}")
        for it in items[:5]:
            print(f"    {it['loc']}: {it['term']!r} in {it['text']!r}")
        if len(items) > 5:
            print(f"    ... and {len(items) - 5} more")
    return 1


if __name__ == "__main__":
    sys.exit(main())
