"""
build_methodology_deck.py - Ionic Wealth Stock Scorecard: methodology & workflow deck.
A team-facing PowerPoint (.pptx) documenting the full process, every formula, and the
step-by-step workflow of the STOCK_SCORECARD_750 engine (FROZEN_METHODOLOGY.md v6.3).
Premium Ionic theme; all prose passes ionic_style.detell() as a de-AI backstop.
Output: 09_PRODUCT/reports/IONIC_SCORECARD_METHODOLOGY_DECK.pptx
Usage: python build_methodology_deck.py
"""
import os, sys
os.environ["PYTHONIOENCODING"] = "utf-8"
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
try:
    from ionic_style import detell
except Exception:
    def detell(t):
        return t

# ---- palette (Ionic premium theme) ----
NAVY   = RGBColor(0x1E, 0x3A, 0x8A)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
BLUEDK = RGBColor(0x1D, 0x4E, 0xD8)
INK    = RGBColor(0x11, 0x18, 0x27)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
CARD   = RGBColor(0xEF, 0xF6, 0xFF)
BAND   = RGBColor(0xF3, 0xF4, 0xF6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SELL   = RGBColor(0xB9, 0x1C, 0x1C)
SELLBG = RGBColor(0xFE, 0xE2, 0xE2)
TRIM   = RGBColor(0x92, 0x40, 0x0E)
TRIMBG = RGBColor(0xFE, 0xF3, 0xC7)
HOLD   = RGBColor(0x15, 0x80, 0x3D)
HOLDBG = RGBColor(0xDC, 0xFC, 0xE7)
GOLD   = RGBColor(0xB0, 0x8D, 0x57)

HEAD = "Bahnschrift"
BODY = "Georgia"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
PAGE = {"n": 0}


def _tf(shape):
    tf = shape.text_frame
    tf.word_wrap = True
    return tf


def _set(run, text, size, color, font=BODY, bold=False, italic=False):
    run.text = detell(text) if isinstance(text, str) else text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic


def box(slide, l, t, w, h, fill=None, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def rbox(slide, l, t, w, h, fill, radius=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = _tf(tb)
    tf.vertical_anchor = anchor
    for i, r in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.space_after = Pt(space)
        run = p.add_run()
        _set(run, r["t"], r["s"], r["c"], r.get("f", BODY), r.get("b", False), r.get("i", False))
    return tb


def bullets(slide, l, t, w, h, items, size=14, gap=8):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = _tf(tb)
    for i, it in enumerate(items):
        lvl = it[2] if len(it) > 2 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = lvl
        mark = "-  " if lvl else "•  "
        r1 = p.add_run(); _set(r1, mark, size, BLUE if lvl == 0 else MUTED, BODY, bold=True)
        r2 = p.add_run(); _set(r2, it[0], size, INK if lvl == 0 else RGBColor(0x37,0x41,0x51))
        if it[1]:
            r3 = p.add_run(); _set(r3, "  " + it[1], size - 1, MUTED, BODY, italic=True)
    return tb


def formula(slide, l, t, w, lines, size=13, h=None):
    hh = h if h else 0.34 * len(lines) + 0.3
    card = rbox(slide, l, t, w, hh, CARD, radius=0.06)
    card.line.color.rgb = RGBColor(0xBF, 0xDB, 0xFE); card.line.width = Pt(0.75)
    tb = slide.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.12), Inches(w - 0.3), Inches(hh - 0.24))
    tf = _tf(tb)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run(); _set(run, ln, size, NAVY, MONO)
    return card


def header(slide, title, kicker=None):
    box(slide, 0, 0, 13.333/1, 0.14, fill=BLUE)
    text(slide, 0.6, 0.34, 12.1, 0.7, [{"t": title, "s": 25, "c": NAVY, "f": HEAD, "b": True}])
    if kicker:
        text(slide, 0.62, 0.92, 12.1, 0.3, [{"t": kicker, "s": 11.5, "c": MUTED, "f": BODY, "i": True}])
    box(slide, 0.6, 1.28 if kicker else 1.12, 12.13, 0.028, fill=BLUE)
    return 1.55 if kicker else 1.4


def footer(slide):
    PAGE["n"] += 1
    box(slide, 0, 7.28, 13.333, 0.22, fill=BAND)
    text(slide, 0.6, 7.28, 9, 0.22,
         [{"t": "Ionic Wealth   ·   Stock Scorecard: Methodology & Workflow   ·   Internal, for team review",
           "s": 8, "c": MUTED, "f": BODY}], anchor=MSO_ANCHOR.MIDDLE)
    text(slide, 11.5, 7.28, 1.2, 0.22, [{"t": str(PAGE["n"]), "s": 8, "c": MUTED, "f": BODY}],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def new(kicker=None, title=None):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, 13.333, 7.5, fill=WHITE)
    if title:
        y = header(s, title, kicker)
        footer(s)
        return s, y
    return s


def table(slide, l, t, w, col_w, rows, header_row=True, size=12, row_h=0.34,
          fills=None, fonts=None, aligns=None):
    """Borderless table: navy header + subtle banding. rows = list of list[str]."""
    nrows, ncols = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(row_h * nrows))
    tbl = gt.table
    # strip default style -> no theme borders/banding
    tblPr = tbl._tbl.tblPr
    for ch in list(tblPr):
        tblPr.remove(ch)
    tblPr.set("firstRow", "0"); tblPr.set("bandRow", "0")
    tot = sum(col_w)
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Emu(int(w * 914400 * cw / tot))
    for i in range(nrows):
        tbl.rows[i].height = Inches(row_h)
        for j in range(ncols):
            cell = tbl.cell(i, j)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if header_row and i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif fills and fills.get((i, j)):
                cell.fill.solid(); cell.fill.fore_color.rgb = fills[(i, j)]
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if (i % 2 == 1) else BAND
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (aligns[j] if aligns else (PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))
            run = p.add_run()
            is_h = header_row and i == 0
            col = WHITE if is_h else (fonts.get((i, j), INK) if fonts else INK)
            _set(run, rows[i][j], size, col, HEAD if is_h else BODY, bold=is_h or (j == 0 and not is_h))
    return gt


# ================================================================= TITLE
s = new()
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
box(s, 0, 0, 13.333, 0.35, fill=GOLD)
box(s, 0, 5.0, 13.333, 0.03, fill=BLUE)
text(s, 0.9, 1.5, 11.5, 0.5, [{"t": "IONIC WEALTH", "s": 18, "c": RGBColor(0xDB,0xEA,0xFE), "f": HEAD, "b": True}])
text(s, 0.9, 2.25, 11.5, 2.0, [
    {"t": "The Stock Scorecard", "s": 46, "c": WHITE, "f": HEAD, "b": True},
    {"t": "Methodology and end-to-end workflow", "s": 24, "c": RGBColor(0xBF,0xDB,0xFE), "f": BODY, "i": True},
], space=8)
text(s, 0.9, 5.25, 11.5, 1.4, [
    {"t": "A quantamental engine that scores stocks 0 to 100 and turns those scores into",
     "s": 14, "c": RGBColor(0xDB,0xEA,0xFE), "f": BODY},
    {"t": "clear Sell / Trim / Hold guidance for client portfolios.",
     "s": 14, "c": RGBColor(0xDB,0xEA,0xFE), "f": BODY},
], space=3)
text(s, 0.9, 6.6, 11.5, 0.4, [
    {"t": "Frozen methodology v6.3   ·   prepared 2026-07-20   ·   internal team walkthrough",
     "s": 11, "c": RGBColor(0x93,0xC5,0xFD), "f": BODY, "i": True}])

# ================================================================= WHAT IT IS
s, y = new("What the system is", "One engine, two jobs")
bullets(s, 0.6, y, 7.4, 4.6, [
    ("A repeatable pipeline that scores any Nifty stock on fundamentals and technicals, then converts the score into portfolio advice.", "", 0),
    ("Two independent scores per stock, never blended in the engine:", "", 0),
    ("3-Year score, tilted to fundamentals (quality, growth, value).", "", 1),
    ("1-Year score, tilted to price behaviour and momentum.", "", 1),
    ("Every stock passes through the same steps, so two analysts on two days get the same discipline.", "", 0),
    ("Recommendation vocabulary is Sell, Trim or Hold. Never Buy.", "this reviews existing holdings, it does not solicit new positions (NDPMS context)", 0),
    ("Live use so far: a real 59-stock client portfolio, plus a 66-name Nifty-100 coverage build.", "", 0),
], size=14.5, gap=11)
c = rbox(s, 8.35, y + 0.1, 4.4, 4.2, CARD, radius=0.05)
c.line.color.rgb = RGBColor(0xBF,0xDB,0xFE); c.line.width = Pt(1)
text(s, 8.6, y + 0.35, 3.9, 0.4, [{"t": "THE PROMISE", "s": 12, "c": BLUE, "f": HEAD, "b": True}])
text(s, 8.6, y + 0.85, 3.9, 3.3, [
    {"t": "Fund-manager judgment,", "s": 20, "c": NAVY, "f": HEAD, "b": True},
    {"t": "delivered with the", "s": 20, "c": NAVY, "f": HEAD, "b": True},
    {"t": "consistency of a machine", "s": 20, "c": NAVY, "f": HEAD, "b": True},
    {"t": " ", "s": 8, "c": INK},
    {"t": "Systematic scoring removes the coin-flips; human review keeps the judgment.",
     "s": 13, "c": RGBColor(0x37,0x41,0x51), "f": BODY, "i": True},
], space=4)

# ================================================================= PIPELINE OVERVIEW
s, y = new("The workflow at a glance", "Six stages, from raw data to a signed client sheet")
stages = [
    ("1", "Data", "Screener feed", BLUE),
    ("2", "Quant score", "0 to 100, x2", BLUEDK),
    ("3", "AI analyst", "one per stock", NAVY),
    ("4", "AI fund mgr", "trim + sizing", NAVY),
    ("5", "Human review", "sign-off", GOLD),
    ("6", "Deliverable", "client + analyst", HOLD),
]
n = len(stages); bw = 1.86; gap = (12.13 - n * bw) / (n - 1); x = 0.6
for i, (num, ttl, sub, col) in enumerate(stages):
    bx = rbox(s, x, y + 0.15, bw, 1.35, col, radius=0.09)
    text(s, x, y + 0.32, bw, 0.5, [{"t": num, "s": 26, "c": WHITE, "f": HEAD, "b": True}], align=PP_ALIGN.CENTER)
    text(s, x, y + 0.82, bw, 0.35, [{"t": ttl, "s": 13.5, "c": WHITE, "f": HEAD, "b": True}], align=PP_ALIGN.CENTER)
    text(s, x, y + 1.16, bw, 0.3, [{"t": sub, "s": 10, "c": RGBColor(0xDB,0xEA,0xFE), "f": BODY, "i": True}], align=PP_ALIGN.CENTER)
    if i < n - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.02), Inches(y + 0.72), Inches(gap - 0.04), Inches(0.22))
        ar.fill.solid(); ar.fill.fore_color.rgb = BORDER; ar.line.fill.background(); ar.shadow.inherit = False
    x += bw + gap
bullets(s, 0.6, y + 1.95, 12.1, 3.0, [
    ("Data:", "Screener fundamentals + price history, refreshed each results season, quality-gated before use.", 0),
    ("Quant score:", "seven percentile-ranked pillars, weighted two ways (3Y and 1Y), then gates and penalties.", 0),
    ("AI analyst:", "one sector-specialist agent researches each stock, forms a forward growth view and a Sell/Hold call.", 0),
    ("AI fund manager:", "sets trim targets and portfolio-level narrative where client weights are known.", 0),
    ("Human review:", "the Principal signs off; nothing reaches a client without it.", 0),
    ("Deliverable:", "a client workbook (one score, clear actions) and a full analyst workbook (every number).", 0),
], size=13, gap=7)

# ================================================================= DIVIDER 1
s = new()
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
box(s, 0, 3.05, 13.333, 0.03, fill=GOLD)
text(s, 0.9, 2.1, 11.5, 0.6, [{"t": "PART ONE", "s": 16, "c": GOLD, "f": HEAD, "b": True}])
text(s, 0.9, 3.25, 11.5, 1.2, [{"t": "The quant engine", "s": 40, "c": WHITE, "f": HEAD, "b": True}])
text(s, 0.92, 4.5, 11.5, 0.6, [{"t": "How a raw stock becomes two scores from 0 to 100.", "s": 17, "c": RGBColor(0xBF,0xDB,0xFE), "f": BODY, "i": True}])

# ================================================================= STAGE 1 DATA
s, y = new("Stage 1  ·  The data layer", "Garbage in is the one thing a scoring engine cannot survive")
bullets(s, 0.6, y, 7.5, 4.6, [
    ("Source: Screener.in deep financials (P&L, balance sheet, cash flow) plus daily adjusted prices and index data.", "", 0),
    ("Refresh cadence is basis results, four times a year, after each earnings season closes.", "next full refresh approx 25 Aug 2026", 0),
    ("Between refreshes, any covered name that reports gets a targeted re-scrape before it feeds a client sheet.", "", 0),
    ("Quality gate (D-009) before any scoring use:", "", 0),
    ("spot-check known values, confirm no look-ahead, sanity-check the schema.", "", 1),
    ("A stale price is expected and corrected silently, it is never treated as a finding.", "", 0),
    ("Every refresh stamps a run-date and a pass/fail record, so the book is auditable.", "", 0),
], size=14, gap=10)
c = rbox(s, 8.4, y + 0.1, 4.35, 3.9, BAND, radius=0.05)
text(s, 8.65, y + 0.3, 3.9, 0.4, [{"t": "THE DATA CONTRACT", "s": 12, "c": NAVY, "f": HEAD, "b": True}])
bullets(s, 8.65, y + 0.85, 3.9, 3.0, [
    ("Same fields, every stock", "", 0),
    ("Point-in-time, no future leakage", "", 0),
    ("Refresh on results, not on whim", "", 0),
    ("Verified before it scores", "", 0),
    ("One source of truth on disk", "", 0),
], size=13, gap=13)

# ================================================================= DUAL HORIZON + WEIGHTS
s, y = new("Stage 2a  ·  Two horizons, seven pillars", "Fundamentals lead the 3-year view; price behaviour leads the 1-year view")
text(s, 0.6, y, 6.0, 0.4, [{"t": "Pillar weights (before regime tilt)", "s": 14, "c": NAVY, "f": HEAD, "b": True}])
table(s, 0.6, y + 0.5, 6.1, [2.0, 0.9, 0.9], [
    ["Pillar", "3Y", "1Y"],
    ["Quality", "20%", "16%"],
    ["Growth", "20%", "16%"],
    ["Value", "18%", "16%"],
    ["Stage / Technical", "14%", "26%"],
    ["Sector & Macro", "11%", "13%"],
    ["Ownership Flow", "9%", "8%"],
    ["Accumulation", "8%", "5%"],
], size=12.5, row_h=0.42)
bullets(s, 7.1, y + 0.1, 5.65, 4.6, [
    ("Each pillar is scored as a percentile rank against the universe, so every score is 'where this stock sits versus its peers', not a raw ratio.", "", 0),
    ("All inputs are winsorized at 2% and 98% before ranking, so a single outlier cannot claim an extreme rank.", "", 0),
    ("The 3Y book leans on Quality, Growth and Value (58% together).", "", 0),
    ("The 1Y book leans on Stage / Technical and near-term flow (39% together).", "", 0),
    ("Intrinsic value (DCF) is deliberately excluded from the machine: the mechanical version produced implausible growth figures, so it now lives as the analyst's own reverse-DCF judgment.", "", 0),
], size=13, gap=10)

# ================================================================= PILLAR FORMULAS I
s, y = new("Stage 2b  ·  Pillar formulas (1 of 2)", "Quality, Growth, Value")
text(s, 0.6, y, 12, 0.35, [{"t": "Quality", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 0.42, 12.13, [
    "Quality = mean( pctile(ROE, sector-neutral),  pctile(ROCE, sector-neutral) )",
    "ROE  = Net Profit / (Equity Capital + Reserves)      ROCE = Operating Profit / (Equity + Borrowings)",
])
text(s, 0.6, y + 1.5, 12, 0.35, [{"t": "Growth", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 1.92, 12.13, [
    "Growth(3Y) = pctile( 3-year revenue CAGR,  whole universe )",
    "Growth(1Y) = pctile( 1-year / TTM revenue growth,  whole universe )",
])
text(s, 0.6, y + 3.0, 12, 0.35, [{"t": "Value", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 3.42, 12.13, [
    "Value = 0.25 x pctile(-PE, universe) + 0.35 x pctile(-PE, sector x cap-tier)",
    "      + 0.20 x pctile(-PB, sector x cap-tier) + 0.20 x pctile(FCF yield, sector x cap-tier)",
])

# ================================================================= PILLAR FORMULAS II
s, y = new("Stage 2b  ·  Pillar formulas (2 of 2)", "Stage / Technical, Sector & Macro, Ownership Flow, Accumulation")
text(s, 0.6, y, 12, 0.35, [{"t": "Stage / Technical", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 0.42, 12.13, [
    "mean( pctile(return, universe), pctile(return, sector) )  over 12M+24M [3Y] / 3M+6M [1Y]",
    "gated by price vs 200DMA [3Y] / 50DMA [1Y];  +/-5 pt RSI nudge on 1Y only",
    "where a chart-analyst pass has run, its Long-Term Pattern Score REPLACES the mechanical 3Y score",
], size=12)
text(s, 0.6, y + 1.75, 12, 0.35, [{"t": "Sector & Macro   /   Ownership Flow   /   Accumulation", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 2.17, 12.13, [
    "Sector & Macro = pctile(sector-mean return, universe)  +  regime-cyclicality fit adjustment",
    "Ownership Flow = pctile( mean(FII qoq + DII qoq) ),  6-of-8 quarters [3Y] / 1-2 quarters [1Y]",
    "Accumulation   = pctile( OBV slope ),  6-12 months [3Y] / 1-3 months [1Y]",
], size=12)
c = rbox(s, 0.6, y + 3.5, 12.13, 0.85, HOLDBG, radius=0.05)
text(s, 0.85, y + 3.62, 11.7, 0.6, [
    {"t": "Why percentiles, not raw numbers", "s": 12.5, "c": HOLD, "f": HEAD, "b": True},
    {"t": "A 22% ROE means nothing in isolation. Ranked against peers it becomes a score you can compare, add and weight. Every pillar speaks the same 0 to 100 language.", "s": 12.5, "c": INK, "f": BODY},
], space=3)

# ================================================================= REGIME
s, y = new("Stage 2c  ·  The regime tilt", "A mild, deliberate nudge, not a market-timing bet")
bullets(s, 0.6, y, 6.1, 3.2, [
    ("Six fixed regimes describe the macro backdrop (risk-on, risk-off, domestic expansion or slowdown, value rotation, quality leadership).", "", 0),
    ("The regime shifts the pillar weights slightly and adjusts the Sector & Macro score for cyclical vs defensive names.", "", 0),
    ("It is slow-moving, re-checked monthly, and the analyst can override it per stock.", "", 0),
], size=13.5, gap=11)
text(s, 6.95, y, 5.8, 0.35, [{"t": "Weight tilt", "s": 14, "c": NAVY, "f": HEAD, "b": True}])
table(s, 6.95, y + 0.45, 5.8, [1.3, 3.6], [
    ["Tag", "Tilt applied"],
    ["Cyclical", "Value +3, Sector +3, Growth -2, Quality -2, Stage -2"],
    ["Not cyclical", "Quality +3, Growth +2, Stage -3, Sector -2"],
], size=11.5, row_h=0.55)
c = rbox(s, 6.95, y + 1.8, 5.8, 1.35, CARD, radius=0.05)
c.line.color.rgb = RGBColor(0xBF,0xDB,0xFE); c.line.width = Pt(0.75)
text(s, 7.2, y + 1.92, 5.35, 1.1, [
    {"t": "Current call (macro desk)", "s": 12, "c": BLUE, "f": HEAD, "b": True},
    {"t": "Value / Cyclical Rotation (mild) 70%  +  India Slowdown (mild) 30%.", "s": 12.5, "c": INK, "f": BODY},
    {"t": "Fit adjustment: Cyclical +4.6, Not-cyclical -2.6.", "s": 12.5, "c": INK, "f": BODY},
], space=3)
bullets(s, 0.6, y + 3.35, 12.1, 1.0, [
    ("Design intent: the tilt is about one-fifth the magnitude of a full melt-up. It shades the score, it does not drive it.", "", 0),
], size=13)

# ================================================================= GATES
s, y = new("Stage 2d  ·  Overlay gates", "Safety checks applied after the weighted score, before it is trusted")
text(s, 0.6, y, 12, 0.35, [{"t": "Balance-sheet safety", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 0.42, 12.13, [
    "RED   if  D/E > 2.5  or  Interest Coverage < 1.5   ->   score capped at 40",
    "AMBER if  D/E > 1.5  or  Interest Coverage < 3     ->   score x 0.85",
])
text(s, 0.6, y + 1.5, 12, 0.35, [{"t": "Liquidity", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 1.92, 12.13, [
    "RED if median 60-day turnover < size-tier bar (Rs 5cr / 1cr / 25L for Large / Mid / Small)  ->  cap 40",
])
c = rbox(s, 0.6, y + 2.85, 12.13, 1.5, TRIMBG, radius=0.05)
text(s, 0.85, y + 3.0, 11.7, 1.25, [
    {"t": "The financial-sector exemption (a lesson we paid for)", "s": 13, "c": TRIM, "f": HEAD, "b": True},
    {"t": "Banks, NBFCs and insurers run high debt-to-equity by design, leverage is their business model, not distress.",
     "s": 13, "c": INK, "f": BODY},
    {"t": "They are exempt from the D/E trigger and tagged 'N/A financial-sector', never silently passed. Their true health is judged on asset quality and capital adequacy in the analyst layer.",
     "s": 13, "c": INK, "f": BODY},
], space=4)

# ================================================================= PENALTY / BOOST
s, y = new("Stage 2e  ·  Penalty and boost", "Red flags compound; a clean bill earns a small lift")
text(s, 0.6, y, 12, 0.35, [{"t": "Exponential red-flag penalty", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 0.42, 7.4, [
    "penalty = -min( 10,  2^(flags) - 1 )",
    "1 flag -1   2 flags -3   3 flags -7   4+ -10",
])
table(s, 8.2, y + 0.42, 4.55, [3.0, 1.0], [
    ["Red flag counted", "Scope"],
    ["Interest coverage < 1.5", "all"],
    ["Debt / equity > 2.5", "non-fin"],
    ["Negative 1yr revenue", "all"],
    ["Growth decel > 15pp", "all"],
    ["Analyst fwd growth < 10%", "all"],
], size=11, row_h=0.36)
text(s, 0.6, y + 2.55, 12, 0.35, [{"t": "Boost", "s": 15, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 2.97, 12.13, [
    "boost = +3   if  zero red flags  AND  Quality > 60th pctile  AND  Value > 60th pctile",
    "full +10 is reserved for an analyst-confirmed clean bill of health (not yet automatic)",
], size=12)
bullets(s, 0.6, y + 4.05, 12.1, 0.7, [
    ("A forward growth estimate below 10% is penalized twice by design: once here, once as a cap on the client score. Slow growth is not rewarded.", "", 0),
], size=12.5)

# ================================================================= QUANT REC
s, y = new("Stage 2f  ·  The quant recommendation", "Conservative by construction")
formula(s, 0.6, y, 12.13, [
    "per horizon:   gate RED -> Sell     score missing -> No Recommendation     else  score >= 40 -> Hold,  < 40 -> Sell",
    "OVERALL  =  Sell  if EITHER horizon says Sell        (surface any real concern, do not average it away)",
], size=12.5)
c = rbox(s, 0.6, y + 1.5, 12.13, 1.15, CARD, radius=0.05)
c.line.color.rgb = RGBColor(0xBF,0xDB,0xFE); c.line.width = Pt(0.75)
text(s, 0.85, y + 1.62, 11.7, 0.9, [
    {"t": "The analyst always has the last word", "s": 13, "c": BLUE, "f": HEAD, "b": True},
    {"t": "Where human research exists, the analyst's Sell/Hold call OVERRIDES the quant call entirely. The machine proposes; the analyst disposes. A leveraged financial that the machine flags on debt can be held on a sound asset-quality read.",
     "s": 13, "c": INK, "f": BODY},
], space=4)
bullets(s, 0.6, y + 2.95, 12.1, 1.3, [
    ("This is why the same score can carry different calls: the number is an input to judgment, not a verdict.", "", 0),
    ("A short-term-tax-aware threshold (30 instead of 40) is designed but inactive until we hold client purchase dates.", "", 0),
], size=13, gap=9)

# ================================================================= DIVIDER 2
s = new()
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
box(s, 0, 3.05, 13.333, 0.03, fill=GOLD)
text(s, 0.9, 2.1, 11.5, 0.6, [{"t": "PART TWO", "s": 16, "c": GOLD, "f": HEAD, "b": True}])
text(s, 0.9, 3.25, 11.5, 1.2, [{"t": "AI research and the Ionic Score", "s": 38, "c": WHITE, "f": HEAD, "b": True}])
text(s, 0.92, 4.5, 11.5, 0.6, [{"t": "From a quant score to the one number a client sees.", "s": 17, "c": RGBColor(0xBF,0xDB,0xFE), "f": BODY, "i": True}])

# ================================================================= AI ANALYST
s, y = new("Stage 3  ·  The AI analyst pass", "One sector specialist per stock, in a single disciplined pass")
bullets(s, 0.6, y, 7.5, 4.6, [
    ("One agent per stock, routed to the right sector desk:", "", 0),
    ("industrials, financials, consumer, IT and telecom, pharma and chemicals.", "", 1),
    ("Each pass, in order:", "", 0),
    ("deep research on the business and its latest results,", "", 1),
    ("an earnings-quality check for one-off items flattering the numbers,", "", 1),
    ("a reverse-DCF judgment: what growth does the price imply, can the business clear it,", "", 1),
    ("a genuine forward 3 to 5 year growth estimate,", "", 1),
    ("a seven-point self-review before it finalizes.", "", 1),
    ("Output is one structured record per stock, saved the moment it is done.", "", 0),
], size=13.5, gap=8)
c = rbox(s, 8.4, y + 0.1, 4.35, 4.15, BAND, radius=0.05)
text(s, 8.65, y + 0.28, 3.9, 0.4, [{"t": "ESCALATION, KEPT NARROW", "s": 12, "c": NAVY, "f": HEAD, "b": True}])
bullets(s, 8.65, y + 0.8, 3.95, 3.3, [
    ("Escalate only a genuine Hold-versus-Sell coin-flip, or a methodology gap that likely hits other stocks.", "", 0),
    ("Never escalate a stale price or ordinary uncertainty.", "", 0),
    ("Escalations are logged for the Principal to rule on, not resolved by the desk.", "", 0),
], size=12.5, gap=12)

# ================================================================= IONIC SCORE
s, y = new("Stage 4  ·  The Ionic Score", "The single 0 to 100 number a client sees, forward-looking by design")
formula(s, 0.6, y, 12.13, [
    "base = 0.60 x final(3Y score)  +  0.40 x final(1Y score)",
    "ionic_score = clamp( base + forward_adjustment,  0,  100 )",
])
text(s, 0.6, y + 1.25, 12, 0.35, [{"t": "Forward adjustment = growth leg + conviction leg  (total clamped to +/- 20)", "s": 14, "c": NAVY, "f": HEAD, "b": True}])
bullets(s, 0.6, y + 1.75, 12.1, 1.4, [
    ("Growth leg:", "scaled from the analyst's forward 3 to 5 year growth estimate (see next slide).", 0),
    ("Conviction leg:", "analyst says Sell -6; analyst holds a name the machine wanted to sell (a rescue) +6; agreement 0.", 0),
], size=13.5, gap=9)
c = rbox(s, 0.6, y + 3.15, 12.13, 1.2, TRIMBG, radius=0.05)
text(s, 0.85, y + 3.28, 11.7, 0.95, [
    {"t": "Two coherence caps keep the score honest", "s": 13, "c": TRIM, "f": HEAD, "b": True},
    {"t": "If forward growth is below 10%, the adjustment can only be zero or negative, never a boost.", "s": 12.5, "c": INK, "f": BODY},
    {"t": "If the analyst says Sell, the adjustment can only be zero or negative. A name we are exiting is never dressed up.", "s": 12.5, "c": INK, "f": BODY},
], space=3)

# ================================================================= GROWTH TABLE
s, y = new("Stage 4  ·  The growth leg", "We give real weight to the future, not the past")
table(s, 0.6, y + 0.1, 7.2, [3.4, 1.5], [
    ["Forward 3 to 5 year growth", "Points"],
    ["Below 5%  (stagnant / declining)", "-15"],
    ["5 to 10%  (below nominal GDP)", "-5"],
    ["10 to 15%  (steady compounder)", "0"],
    ["15 to 20%  (strong)", "+5"],
    ["20 to 25%  (exceptional momentum)", "+10"],
    ["25% and above  (hypergrowth)", "+15"],
    ["Exceptional tier (see right)", "+20"],
], size=12.5, row_h=0.44,
   fonts={(1,1):SELL,(2,1):SELL,(4,1):HOLD,(5,1):HOLD,(6,1):HOLD,(7,1):HOLD},
   fills={(7,0):CARD,(7,1):CARD})
c = rbox(s, 8.15, y + 0.1, 4.6, 2.5, CARD, radius=0.05)
c.line.color.rgb = RGBColor(0xBF,0xDB,0xFE); c.line.width = Pt(0.75)
text(s, 8.4, y + 0.28, 4.15, 2.25, [
    {"t": "The +20 exceptional tier", "s": 13, "c": BLUE, "f": HEAD, "b": True},
    {"t": "Reserved for the rare compounder that clears all three:", "s": 12.5, "c": INK, "f": BODY},
    {"t": "growth at or above 25%,", "s": 12.5, "c": NAVY, "f": BODY, "b": True},
    {"t": "return on equity at or above 20%,", "s": 12.5, "c": NAVY, "f": BODY, "b": True},
    {"t": "little or no share dilution.", "s": 12.5, "c": NAVY, "f": BODY, "b": True},
    {"t": "The fund manager confirms the dilution test by hand.", "s": 11.5, "c": MUTED, "f": BODY, "i": True},
], space=5)
text(s, 8.15, y + 2.8, 4.6, 1.4, [
    {"t": "Growth is the analyst's genuine forward view, never a copy of trailing CAGR. The penalty for slow growth exists on purpose.",
     "s": 12.5, "c": INK, "f": BODY, "i": True}])

# ================================================================= TWO-GATE CLIENT REC
s, y = new("Stage 5  ·  Sell, Trim or Hold", "Two gates turn a score into a portfolio action")
text(s, 0.6, y, 12, 0.35, [{"t": "Gate A  ·  stock quality", "s": 14, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 0.42, 12.13, [
    "analyst says Sell  ->  Sell         else  ionic_score < 40  ->  Sell",
], size=12)
text(s, 0.6, y + 1.25, 12, 0.35, [{"t": "Gate B  ·  concentration (the Trim band)", "s": 14, "c": NAVY, "f": HEAD, "b": True}])
formula(s, 0.6, y + 1.67, 12.13, [
    "ionic_score 40 to 50  AND  position > 2.5% of book  ->  Trim to a stated target weight",
], size=12)
text(s, 0.6, y + 2.5, 6.0, 0.35, [{"t": "Concentration guidance (not hard caps)", "s": 14, "c": NAVY, "f": HEAD, "b": True}])
table(s, 0.6, y + 2.95, 7.0, [2.2, 3.6], [
    ["Single-stock weight", "Guidance"],
    ["5 to 10%", "Fine; reduce if forward growth is modest"],
    ["Above 10%", "A little high; trim advice unless conviction is high"],
    ["Above 20%", "Extreme; strong trim advice"],
], size=11.5, row_h=0.48)
bullets(s, 7.9, y + 2.6, 4.85, 2.2, [
    ("Trim targets are set by the fund-manager pass, by judgment, not a formula.", "", 0),
    ("Buying price, the client's mandate, and large versus small cap all matter. None is a rigid rule.", "", 0),
    ("Freed cash is shown as cash. An existing-holdings review never auto-redeploys it.", "", 0),
], size=12.5, gap=10)

# ================================================================= ANALYTICS LAYER
s, y = new("Stage 6  ·  The analytics layer", "Portfolio-level risk and return, in plain words, honestly framed")
bullets(s, 0.6, y, 6.3, 4.5, [
    ("A simulated backcast holds today's mix at constant weights and compares it to the Nifty 50 (dividends added back).", "", 0),
    ("Measured over 3 years and 1 year:", "", 0),
    ("return, volatility, Sharpe and Sortino, deepest fall,", "", 1),
    ("beta, alpha, tracking error, up and down capture,", "", 1),
    ("a four-factor regression (market, size, value, momentum),", "", 1),
    ("a correlation heatmap of the largest holdings,", "", 1),
    ("index valuation context (P/E percentile vs history).", "", 1),
], size=13.5, gap=8)
c = rbox(s, 7.15, y + 0.1, 5.6, 4.3, TRIMBG, radius=0.05)
text(s, 7.4, y + 0.3, 5.15, 0.4, [{"t": "THE HONESTY RULES (BINDING)", "s": 12, "c": TRIM, "f": HEAD, "b": True}])
bullets(s, 7.4, y + 0.85, 5.15, 3.4, [
    ("The simulation is a portrait of today's mix, it is not the client's realized return, and it is not a forecast.", "", 0),
    ("The selection-bias caveat is printed on the client sheet itself.", "", 0),
    ("The risk-free rate assumption is stated (6.5%).", "", 0),
    ("Any expected-alpha number stays analyst-side and labelled an estimate, never shown to the client as a forward return.", "", 0),
], size=12.5, gap=11)

# ================================================================= DIVIDER 3
s = new()
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
box(s, 0, 3.05, 13.333, 0.03, fill=GOLD)
text(s, 0.9, 2.1, 11.5, 0.6, [{"t": "PART THREE", "s": 16, "c": GOLD, "f": HEAD, "b": True}])
text(s, 0.9, 3.25, 11.5, 1.2, [{"t": "Deliverables, controls, coverage", "s": 36, "c": WHITE, "f": HEAD, "b": True}])
text(s, 0.92, 4.5, 11.5, 0.6, [{"t": "What ships, how we keep it honest, and where we are.", "s": 17, "c": RGBColor(0xBF,0xDB,0xFE), "f": BODY, "i": True}])

# ================================================================= DELIVERABLES
s, y = new("The two deliverables", "One for the client, one for the desk. Both ship together.")
c1 = rbox(s, 0.6, y + 0.1, 6.0, 4.3, CARD, radius=0.05)
c1.line.color.rgb = RGBColor(0xBF,0xDB,0xFE); c1.line.width = Pt(1)
text(s, 0.85, y + 0.3, 5.5, 0.5, [{"t": "Client workbook", "s": 17, "c": NAVY, "f": HEAD, "b": True},
                                   {"t": "Ionic Wealth branded, four sheets", "s": 11.5, "c": MUTED, "f": BODY, "i": True}], space=2)
bullets(s, 0.85, y + 1.25, 5.5, 3.0, [
    ("At a Glance: KPI dashboard + growth-of-100 chart.", "", 0),
    ("Recommendations: one Ionic Score, one action, one rationale per holding.", "", 0),
    ("Portfolio Analytics: risk and return in plain words.", "", 0),
    ("Before vs After: what the actions do to the book.", "", 0),
    ("One score. Sell / Trim / Hold. No jargon.", "", 0),
], size=12.5, gap=9)
c2 = rbox(s, 6.75, y + 0.1, 6.0, 4.3, BAND, radius=0.05)
text(s, 7.0, y + 0.3, 5.5, 0.5, [{"t": "Analyst workbook", "s": 17, "c": NAVY, "f": HEAD, "b": True},
                                  {"t": "the full working paper, 46 columns", "s": 11.5, "c": MUTED, "f": BODY, "i": True}], space=2)
bullets(s, 7.0, y + 1.25, 5.5, 3.0, [
    ("Every pillar sub-score, both horizons, gates, penalty, boost.", "", 0),
    ("The full research record per stock, nothing clipped.", "", 0),
    ("A field guide: the workbook documents its own method.", "", 0),
    ("Portfolio Analytics (Full): factor betas with t-stats.", "", 0),
    ("Empty technical columns hide themselves until that pass runs.", "", 0),
], size=12.5, gap=9)

# ================================================================= QUALITY CONTROLS
s, y = new("How we keep it honest", "The controls matter as much as the formulas")
rows = [
    ("Schema validation", "every research file is checked for all fields and correct types before it counts."),
    ("Zero-tell style gate", "the builder refuses to ship a sheet containing AI writing tells or banned words."),
    ("Independent QA sweep", "a second agent audits each file: does the call match the reasoning, is the growth number calibrated, any Buy language, any pending-vs-done error."),
    ("Fact-check on claims", "completed-versus-pending legal and deal claims are verified against primary sources."),
    ("Growth re-adjudication", "a flagged growth number goes back to the same analyst desk to defend or revise."),
    ("Human sign-off", "nothing reaches a client without the Principal's approval."),
]
yy = y + 0.05
for i, (t1, t2) in enumerate(rows):
    rbox(s, 0.6, yy, 0.16, 0.66, BLUE, radius=0.3)
    text(s, 0.95, yy - 0.02, 3.1, 0.7, [{"t": t1, "s": 13.5, "c": NAVY, "f": HEAD, "b": True}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.15, yy - 0.02, 8.6, 0.7, [{"t": t2, "s": 12.5, "c": INK, "f": BODY}], anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.72
text(s, 0.95, yy + 0.02, 11.8, 0.5, [
    {"t": "In the Nifty-100 build, this layer caught three overstated legal claims and pulled three growth numbers down a full band before they touched a score.",
     "s": 12, "c": MUTED, "f": BODY, "i": True}])

# ================================================================= COVERAGE + ROADMAP
s, y = new("Where we are, and what is next", "Coverage to date and the road ahead")
# KPI cards
kpis = [("125", "stocks researched", "59 client holdings + 66 Nifty-100"),
        ("7", "pillars x 2 horizons", "one consistent scoring spine"),
        ("100%", "human signed-off", "no client sheet ships otherwise")]
x = 0.6
for num, lab, sub in kpis:
    cc = rbox(s, x, y + 0.05, 3.9, 1.55, CARD, radius=0.06)
    cc.line.color.rgb = RGBColor(0xBF,0xDB,0xFE); cc.line.width = Pt(0.75)
    text(s, x + 0.2, y + 0.18, 3.6, 0.6, [{"t": num, "s": 30, "c": NAVY, "f": HEAD, "b": True}])
    text(s, x + 0.22, y + 0.82, 3.6, 0.3, [{"t": lab.upper(), "s": 11, "c": BLUE, "f": HEAD, "b": True}])
    text(s, x + 0.22, y + 1.12, 3.6, 0.4, [{"t": sub, "s": 10.5, "c": MUTED, "f": BODY, "i": True}])
    x += 4.06
text(s, 0.6, y + 1.85, 6.0, 0.35, [{"t": "Coverage so far", "s": 14, "c": NAVY, "f": HEAD, "b": True}])
bullets(s, 0.6, y + 2.3, 6.1, 2.3, [
    ("Real 59-stock client portfolio: fully scored, researched and delivered.", "", 0),
    ("Nifty-100 build: 66 new names researched, 27 Sell / 39 Hold.", "valuation vs deliverable growth, not business quality", 0),
    ("Every call carries a rationale, a forward growth view and a source list.", "", 0),
], size=13, gap=9)
text(s, 7.0, y + 1.85, 5.7, 0.35, [{"t": "The road ahead", "s": 14, "c": NAVY, "f": HEAD, "b": True}])
bullets(s, 7.0, y + 2.3, 5.75, 2.3, [
    ("Extend the quant score to the full Nifty-100, then the Nifty-750.", "", 0),
    ("Add human proof-reading and an intern re-review for a second pair of eyes.", "", 0),
    ("A consistency check so calls never contradict our house views or PMS holdings.", "", 0),
    ("Quarterly model and format upgrades, basis results.", "", 0),
], size=13, gap=8)

# ================================================================= CLOSE
s = new()
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
box(s, 0, 0, 13.333, 0.35, fill=GOLD)
text(s, 0.9, 2.7, 11.5, 1.4, [
    {"t": "A repeatable, defensible process", "s": 34, "c": WHITE, "f": HEAD, "b": True},
    {"t": "Systematic where machines are better, human where judgment is better.", "s": 18, "c": RGBColor(0xBF,0xDB,0xFE), "f": BODY, "i": True},
], space=10)
text(s, 0.92, 5.4, 11.5, 0.5, [
    {"t": "Method frozen in FROZEN_METHODOLOGY.md v6.3   ·   Ionic Wealth   ·   internal",
     "s": 12, "c": RGBColor(0x93,0xC5,0xFD), "f": BODY, "i": True}])

out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "reports", "IONIC_SCORECARD_METHODOLOGY_DECK.pptx")
try:
    prs.save(out)
except PermissionError:
    out = out.replace(".pptx", "_v2.pptx"); prs.save(out)
print("Saved", out, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
