# -*- coding: utf-8 -*-
"""build_pr_full.py -- Ionic Wealth Portfolio Review, FULL premium deck (~35 pages).
Built to the Fable elevated art-direction spec + research best-features, on real Kordes data.
Private-bank register: verdict titles, standfirsts, section dividers, KPI register strips,
in-table score bars, risk register, bridge/waterfall, before/after tornado, lollipop,
order sheet, pull-quote spotlights. One gold accent per page; negatives in ink; no AI-tells.
Reads full_data.json + commentary_v3.json. Deck path is gitignored (client PII local-only).
"""
import os, json, math, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
import chart_lib as cl   # matplotlib Ionic chart engine (renders PNGs to embed)
from PIL import Image as _PILImage

# ---------- palette (Ionic Wealth by Angel One brand) ----------
NAVYD = RGBColor(0x10,0x19,0x7A); NAVY = RGBColor(0x1B,0x27,0xA3)
NT1 = RGBColor(0x4A,0x57,0xC4); NT2 = RGBColor(0x8C,0x95,0xDE); NT3 = RGBColor(0xC9,0xCE,0xF0)
GHOST = RGBColor(0x24,0x2F,0x8E); NAVYHAIR = RGBColor(0x33,0x3E,0x9E)
GOLD = RGBColor(0xF2,0xA9,0x3C); ORANGE = GOLD; INK = RGBColor(0x16,0x23,0x3B); SLATE = RGBColor(0x6B,0x72,0x80)
HAIR = RGBColor(0xE5,0xE7,0xEB); TRACK = RGBColor(0xEE,0xEF,0xF7)
PANEL = RGBColor(0xF5,0xF6,0xFC); WHITE = RGBColor(0xFF,0xFF,0xFF); ONNAVY = RGBColor(0xC9,0xCE,0xF0)
CHIP = RGBColor(0x7B,0x9F,0xFF); NAVYLT = RGBColor(0x8C,0x95,0xDE)
SELL = RGBColor(0xE0,0x40,0x2F); SELLBG = RGBColor(0xFB,0xE3,0xE0)
HOLD = RGBColor(0x1E,0x9E,0x6A); HOLDBG = RGBColor(0xE0,0xF2,0xEA)
SERIF = "Georgia"; SANS = "Bahnschrift"
LOGO = r"C:\Users\SHREYA~1.1GU\AppData\Local\Temp\claude\c--Users-Shreyas-1Gupta-OneDrive---Angel-Broking-Limited-Desktop-Backup-NIFTY-500\5ec2bf16-8c38-4f40-9e4f-8e07be6545fd\scratchpad\assets\logo_clean.png"
COVER_PHOTO = r"C:\Users\SHREYA~1.1GU\AppData\Local\Temp\claude\c--Users-Shreyas-1Gupta-OneDrive---Angel-Broking-Limited-Desktop-Backup-NIFTY-500\5ec2bf16-8c38-4f40-9e4f-8e07be6545fd\scratchpad\assets\p1_img0_1379x1536.jpeg"
CW, CH = 13.333, 7.5
ML, MR = 0.92, 0.92
RX = CW - MR  # 12.413
UW = RX - ML  # usable width ~11.49

SC = os.environ.get
DIG = SC("PRF_DIGEST", r"C:\Users\SHREYA~1.1GU\AppData\Local\Temp\claude\c--Users-Shreyas-1Gupta-OneDrive---Angel-Broking-Limited-Desktop-Backup-NIFTY-500\5ec2bf16-8c38-4f40-9e4f-8e07be6545fd\scratchpad\full_data.json")
COMM = os.path.join(os.path.dirname(DIG), "..", "..", "")  # placeholder
COMMP = r"c:\Users\Shreyas.1Gupta\OneDrive - Angel Broking Limited\Desktop\Backup\NIFTY 500\Shreyas_Ionic_AMC\04_RND_LAB\STOCK_SCORECARD_750\results\pr_kordes\commentary_v3.json"
OUT = SC("PRF_OUT", r"C:\Users\SHREYA~1.1GU\AppData\Local\Temp\claude\c--Users-Shreyas-1Gupta-OneDrive---Angel-Broking-Limited-Desktop-Backup-NIFTY-500\5ec2bf16-8c38-4f40-9e4f-8e07be6545fd\scratchpad\PR_full.pptx")

D = json.load(open(DIG, encoding="utf-8"))
C = json.load(open(COMMP, encoding="utf-8")) if os.path.exists(COMMP) else {"stocks": {}, "funds": {}}
T = D["totals"]; STK = D["stocks"]; MF = D["mf"]
CST = C.get("stocks", {}); CFU = C.get("funds", {})

# ---------- derived ----------
def fcr(v): return f"₹{v/1e7:.2f} Cr"
def fl(v):  return f"{v/1e5:,.1f}"
def clip_two(text):
    """First two sentences, cleanly (for pull-quotes)."""
    text = (text or "").replace("\n", " ").strip()
    parts = re.split(r'(?<=[.!?])\s+', text)
    out = " ".join(parts[:2]).strip()
    if out and out[-1] not in ".!?": out += "."
    return out
def gov(st):
    """Binding Ionic Score = lower of the 3Y and 1Y horizon; a Sell triggers if EITHER is below 40."""
    vals = [v for v in (st.get("score3y"), st.get("score1y")) if v is not None]
    return min(vals) if vals else None
DER = D.get("derived", {})
grand = T["grand"]; eq = T["equity"]
sell_stocks = [s for s in STK if s["call"] == "Sell"]
sell_val = DER.get("sell_val", sum(s["value"] for s in sell_stocks))
sell_pct_book = DER.get("sell_pct_book", round(100*sell_val/grand, 1))
top10_val = sum(s["value"] for s in STK[:10])
intl_val = DER.get("intl_val", 0); intl_pct = DER.get("intl_pct_book", 0)
reg_val = DER.get("reg_val", 0)
n_tested = DER.get("n_tested", sum(1 for f in MF if f["alpha"] is not None))
fund_actions = [f for f in MF if f["action"] != "Hold"]
fund_action_val = DER.get("fund_action_val", sum(f["value"] for f in fund_actions if f["value"]))
titan_trim = DER.get("titan_trim", 0)
proceeds = DER.get("proceeds", sell_val + fund_action_val + titan_trim)

def pic(s, path, x, y, boxw, boxh, valign="middle"):
    """Embed preserving aspect ratio, fitted+centred inside the (boxw,boxh) box. Kills distortion."""
    try:
        iw, ih = _PILImage.open(path).size; ar = iw / ih; bar = boxw / boxh
        if ar >= bar: w, h = boxw, boxw / ar
        else: w, h = boxh * ar, boxh
        ox = x + (boxw - w) / 2
        oy = y + (0 if valign == "top" else (boxh - h) if valign == "bottom" else (boxh - h) / 2)
        return s.shapes.add_picture(path, Inches(ox), Inches(oy), Inches(w), Inches(h))
    except Exception:
        return s.shapes.add_picture(path, Inches(x), Inches(y), Inches(boxw), Inches(boxh))
def sym_of(st): return (st.get("symbol") or st["name"]).strip()
def callcol(st): return cl.SELL if st.get("call") == "Sell" else cl.HOLD

prs = Presentation(); prs.slide_width = Inches(CW); prs.slide_height = Inches(CH)
BLANK = prs.slide_layouts[6]
PAGES = []  # (kind) for tracking

# ---------- primitives ----------
def nosh(sh):
    try: sh.shadow.inherit = False
    except Exception: pass

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); nosh(r)
    return s

def rect(s, x, y, w, h, fill=None, line=None, lw=0.75, round_=0.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(max(w, 0.001)), Inches(max(h, 0.001)))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    nosh(shp)
    if round_:
        try: shp.adjustments[0] = round_
        except Exception: pass
    return shp

def oval(s, x, y, d, fill):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = fill; o.line.fill.background(); nosh(o)
    return o

def rule(s, x, y, w, color=HAIR, h=0.01):
    return rect(s, x, y, w, h, fill=color)

def vrule(s, x, y, h, color=HAIR, w=0.008):
    return rect(s, x, y, w, h, fill=color)

def txt(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, ls=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf, m, 0)
    if isinstance(paras[0], tuple): paras = [paras]
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(0)
        if ls: p.line_spacing = ls
        for run in para:
            t, fn, sz, col, bold = run[0], run[1], run[2], run[3], run[4]
            ital = run[5] if len(run) > 5 else False
            spc = run[6] if len(run) > 6 else None
            r = p.add_run(); r.text = t
            r.font.name = fn; r.font.size = Pt(sz); r.font.bold = bold
            r.font.italic = ital; r.font.color.rgb = col
            if spc is not None: r.font._rPr.set('spc', str(int(spc)))
    return tb

def pill(s, x, y, text, w=0.80, kind=None):
    kind = kind or text
    fill, tc = {"Sell": (SELLBG, SELL), "Exit": (SELLBG, SELL), "Redeem": (SELLBG, SELL),
                "Hold": (HOLDBG, HOLD), "Within": (HOLDBG, HOLD), "Aligned": (HOLDBG, HOLD),
                "Switch": (PANEL, SLATE), "Trim": (PANEL, SLATE), "Watch": (PANEL, SLATE),
                "Partial": (PANEL, SLATE), "Gap": (SELLBG, SELL), "Breach": (SELLBG, SELL),
                }.get(kind, (PANEL, SLATE))
    rect(s, x, y, w, 0.24, fill=fill, line=tc, lw=0.75, round_=0.5)
    txt(s, x, y - 0.006, w, 0.25, [(text.upper(), SANS, 8, tc, True, False, 60)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------- chrome ----------
def marker(s, section_no, section_name):
    if not section_name: return
    txt(s, 9.4, 0.30, RX-9.4, 0.22, [(f"{section_no:02d} · {section_name.upper()}", SANS, 7.5, NT2, True, False, 120)],
        align=PP_ALIGN.RIGHT)
    for k in range(4):
        col = NAVY if (k+1) == section_no else HAIR
        rect(s, 11.62 + 0.20*k, 0.52, 0.15, 0.022, fill=col)

def logo(s):
    try: pic(s, LOGO, RX-1.85, 0.40, 1.8, 0.31)
    except Exception: txt(s, RX-2.2, 0.42, 2.2, 0.3, [("IONIC WEALTH", SANS, 11, NAVY, True)], align=PP_ALIGN.RIGHT)

def classified(s, dark=False):
    txt(s, CW/2-1.7, 0.14, 3.4, 0.2, [("Classified as Internal", SANS, 7.5, (NT2 if dark else SLATE), False)], align=PP_ALIGN.CENTER)

def footer(s, folio, dark=False):
    txt(s, RX-2.0, 7.14, 2.0, 0.2, [(f"Portfolio Review  ·  {folio}", SANS, 7.5, (NT2 if dark else SLATE), False)], align=PP_ALIGN.RIGHT)

def source(s, text):
    txt(s, ML, 6.66, UW, 0.24, [(text, SANS, 7, SLATE, False)])

def content(section_no, section_name, folio, eyebrow, title, standfirst=None, gold_eyebrow=True):
    s = slide(WHITE)
    logo(s); classified(s)
    txt(s, ML, 0.46, 9.1, 0.55, [(eyebrow, SANS, 26, NAVY, True)])          # navy topic
    txt(s, ML, 1.04, 10.2, 0.45, [(title, SANS, 17.5, ORANGE, True)])       # orange subtitle (the verdict)
    rule(s, ML, 1.54, UW, NAVY, 0.024)
    footer(s, folio)
    return s

# ---------- components ----------
def kpi_strip(s, stats, y=1.78, x=ML, w=UW, h=0.95):
    n = len(stats); cw = w/n
    for i, st in enumerate(stats):
        cx = x + i*cw
        val, lab = st[0], st[1]; ctx = st[2] if len(st) > 2 else None
        txt(s, cx, y, cw-0.15, 0.5, [(val, SANS, 27, INK, False)])
        txt(s, cx, y+0.52, cw-0.15, 0.24, [(lab.upper(), SANS, 8, SLATE, True, False, 200)], ls=1.0)
        if ctx: txt(s, cx, y+0.72, cw-0.15, 0.2, [(ctx, SANS, 8, NT2, False)])
        if i: vrule(s, cx-0.06, y+0.05, 0.62, HAIR, 0.008)

def score_bar(s, x, y, score, w=0.65):
    rect(s, x, y, w, 0.07, fill=TRACK)
    if score is not None:
        rect(s, x, y, w*max(score,0)/100.0, 0.07, fill=NT1)
    rect(s, x+w*0.40, y-0.03, 0.012, 0.13, fill=INK)  # 40 threshold
    txt(s, x+w+0.07, y-0.085, 0.45, 0.24, [((f"{score:.0f}" if score is not None else "-"), SANS, 9, INK, False)],
        anchor=MSO_ANCHOR.MIDDLE)

def heat_strip(s, x, y, pillars):
    """5 small cells: quality, growth, value, stage, macro."""
    keys = [("Q","quality"),("G","growth3y"),("V","value"),("T","stage3y"),("M","macro")]
    cw = 0.30
    for i,(lab,k) in enumerate(keys):
        v = pillars.get(k)
        if v is None: col = HAIR
        elif v >= 60: col = HOLD
        elif v >= 40: col = NT2
        else: col = SELL
        rect(s, x+i*cw, y, cw-0.03, 0.20, fill=col)
        txt(s, x+i*cw, y+0.20, cw-0.03, 0.16, [(lab, SANS, 6.5, SLATE, False)], align=PP_ALIGN.CENTER)

def donut(s, x, y, wh, pairs, colors, hole=65):
    cd = CategoryChartData(); cd.categories = [p[0] for p in pairs]; cd.add_series("a", [p[1] for p in pairs])
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(wh), Inches(wh), cd)
    ch = gf.chart; ch.has_title = False; ch.has_legend = False
    ser = ch.plots[0].series[0]
    for i, pt in enumerate(ser.points):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.color.rgb = WHITE; pt.format.line.width = Pt(1.5)
    for dc in ch._chartSpace.iter(qn('c:doughnutChart')):
        hs = dc.find(qn('c:holeSize'))
        if hs is None: hs = dc.makeelement(qn('c:holeSize'), {}); dc.append(hs)
        hs.set('val', str(hole))
    return gf

def stacked100(s, x, y, w, h, segs, label_above=True):
    tot = sum(p[1] for p in segs) or 1; cx = x
    for i,(lab,pc,col) in enumerate(segs):
        sw = w*pc/tot
        if sw <= 0.002: continue
        rect(s, cx+(0.015 if i else 0), y, sw-(0.015 if i else 0), h, fill=col)
        if label_above and sw > 0.7:
            txt(s, cx, y-0.30, sw, 0.28, [[(lab, SANS, 9, INK, True)],[(f"{pc:.1f}%", SANS, 9, SLATE, False)]], ls=1.0)
        cx += sw

def utbl(s, x, y, w, cols, rows, rowh=0.34, fs=10, hfs=8, header=True, totals=None):
    """Upgraded register table. cols=(label,frac,align). cell = str|('n',t[,col])|('b',t[,col])|
    ('pill',text)|('bar',score)|('heat',pillars)."""
    tot = sum(c[1] for c in cols); xs=[]; ws=[]; cx=x
    for (_,cw,_a) in cols: xs.append(cx); ws.append(w*cw/tot); cx += w*cw/tot
    PAD = 0.07
    def cb(i):  # padded cell box
        return xs[i]+PAD, ws[i]-2*PAD
    if header:
        for i,(lab,_c,al) in enumerate(cols):
            a = PP_ALIGN.RIGHT if al=="r" else (PP_ALIGN.CENTER if al=="c" else PP_ALIGN.LEFT)
            cx, cw = cb(i)
            txt(s, cx, y, cw, 0.24, [(lab.upper(), SANS, hfs, SLATE, True, False, 200)], align=a, anchor=MSO_ANCHOR.MIDDLE)
        rule(s, x, y+0.28, w, NAVY, 0.015); ry = y+0.32
    else:
        ry = y
    for row in rows:
        for i, cell in enumerate(row):
            al = cols[i][2]; a = PP_ALIGN.RIGHT if al=="r" else (PP_ALIGN.CENTER if al=="c" else PP_ALIGN.LEFT)
            cx, cw = cb(i)
            if isinstance(cell, tuple) and cell[0]=="pill":
                pw = min(ws[i]-0.04, 0.80); pill(s, xs[i]+(ws[i]-pw)/2, ry+(rowh-0.24)/2, cell[1], pw)
            elif isinstance(cell, tuple) and cell[0]=="bar":
                score_bar(s, xs[i]+PAD, ry+(rowh-0.07)/2, cell[1], w=min(ws[i]-0.55, 0.65))
            elif isinstance(cell, tuple) and cell[0]=="heat":
                heat_strip(s, xs[i]+PAD, ry+(rowh-0.36)/2, cell[1])
            elif isinstance(cell, tuple):
                tag, tval = cell[0], cell[1]; col = cell[2] if len(cell)>2 else INK
                fn = SANS if tag in ("n","b") else SERIF
                txt(s, cx, ry, cw, rowh, [(tval, SANS, fs, col, tag=="b")], align=a, anchor=MSO_ANCHOR.MIDDLE)
            else:
                txt(s, cx, ry, cw, rowh, [(str(cell), SERIF, fs, INK, False)], align=a, anchor=MSO_ANCHOR.MIDDLE)
        rule(s, x, ry+rowh, w, HAIR, 0.008); ry += rowh
    if totals:
        rule(s, x, ry+0.02, w, NAVY, 0.014)
        for i, cell in enumerate(totals):
            if cell is None: continue
            al = cols[i][2]; a = PP_ALIGN.RIGHT if al=="r" else (PP_ALIGN.CENTER if al=="c" else PP_ALIGN.LEFT)
            cx, cw = cb(i)
            tval = cell[1] if isinstance(cell, tuple) else str(cell)
            fn = SERIF if i==0 else SANS
            txt(s, cx, ry+0.06, cw, rowh, [(tval, fn, fs+0.5, INK, True)], align=a, anchor=MSO_ANCHOR.MIDDLE)
        rule(s, x, ry+0.06+rowh, w, HAIR, 0.008); ry += rowh+0.06
    return ry

def risk_register(s, rows, x=ML, y=1.80, w=7.4):
    cols = [("Metric",2.6,"l"),("Portfolio",1.2,"r"),("Guideline",1.5,"r"),("Status",1.3,"c")]
    return utbl(s, x, y, w, cols, rows, rowh=0.40, fs=10.5, hfs=8)

def bridge(s, steps, x0=ML, x1=RX, ytop=2.10, ybase=5.55):
    """steps: list of (label, value, kind) kind in open/flow/close. Draws floating waterfall."""
    vals = [v for _,v,_ in steps]
    cum = 0; tops = []; running = []
    maxc = 0
    for lab,v,kind in steps:
        if kind == "open": running.append((0, v)); cum = v
        elif kind == "close": running.append((0, v)); cum = v
        else: running.append((cum - v, cum)); cum = cum - v
        maxc = max(maxc, cum, v)
    plot_h = ybase - ytop - 0.30
    per = plot_h / (maxc if maxc else 1)
    n = len(steps); bw = 0.62; gap = ((x1-x0) - n*bw)/(n-1) if n>1 else 0
    gold_idx = None
    for i,(lab,v,kind) in enumerate(steps):
        bx = x0 + i*(bw+gap)
        lo, hi = running[i]
        y_hi = ybase - hi*per; h = (hi-lo)*per
        col = NAVY if kind in ("open","close") else NT2
        rect(s, bx, y_hi, bw, max(h,0.02), fill=col)
        if kind == "goldflow": rect(s, bx, y_hi-0.03, bw, 0.03, fill=GOLD)
        txt(s, bx-0.15, y_hi-0.24, bw+0.30, 0.22, [(f"₹{v/1e5:.1f} L", SANS, 9, INK, False)], align=PP_ALIGN.CENTER)
        txt(s, bx-0.22, ybase+0.06, bw+0.44, 0.5, [(lab, SANS, 7.5, SLATE, False)], align=PP_ALIGN.CENTER, ls=1.0)
        if i < n-1:
            rule(s, bx+bw, y_hi if kind!="open" else ybase-hi*per, gap, HAIR, 0.008)
    rule(s, x0, ybase, x1-x0, SLATE, 0.012)

def tornado(s, rows, spine=6.67, y0=2.15, pitch=0.62, maxlen=3.4):
    """rows: (label, a_num, a_str, b_num, b_str). a=today, b=proposed. Length from nums, display strings."""
    mx = max(max(abs(a),abs(b)) for _,a,_as,b,_bs in rows) or 1
    txt(s, spine-3.9, 1.80, 3.7, 0.22, [("TODAY", SANS, 8, SLATE, True, False, 150)], align=PP_ALIGN.RIGHT)
    txt(s, spine+0.10, 1.80, 2.8, 0.22, [("PROPOSED / TARGET", SANS, 8, SLATE, True, False, 150)])
    for i,(lab,a,astr,b,bstr) in enumerate(rows):
        y = y0 + i*pitch
        la = maxlen*abs(a)/mx; lb = maxlen*abs(b)/mx
        rect(s, spine-0.10-la, y, la, 0.30, fill=NT3)
        rect(s, spine+0.10, y, lb, 0.30, fill=NAVY)
        txt(s, spine-3.9, y-0.26, 3.7, 0.22, [(lab, SERIF, 10.5, INK, False)], align=PP_ALIGN.RIGHT)
        txt(s, spine-0.16-la-1.0, y, 0.95, 0.30, [(astr, SANS, 9, SLATE, False)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, spine+0.16+lb, y, 1.6, 0.30, [(bstr, SANS, 9, INK, True)], anchor=MSO_ANCHOR.MIDDLE)
    vrule(s, spine, y0-0.10, len(rows)*pitch, HAIR, 0.01)

def lollipop(s, rows, x=ML, y=1.95, w=UW, h=4.3, gold_top=True):
    """rows: (name, value_pp) sorted; value can be +/- (over/underweight)."""
    mx = max(abs(v) for _,v in rows) or 1
    axis_x = x + 3.0
    plot_w = (x+w) - axis_x - 0.6
    vrule(s, axis_x, y, h, NT2, 0.01)
    txt(s, axis_x-0.3, y+h+0.04, 0.6, 0.2, [("0", SANS, 7.5, SLATE, False)], align=PP_ALIGN.CENTER)
    pitch = h/len(rows)
    big = max(range(len(rows)), key=lambda i: rows[i][1])
    for i,(nm,v) in enumerate(rows):
        cy = y + i*pitch + pitch/2
        txt(s, x, cy-0.11, 2.7, 0.22, [(nm, SERIF, 9.5, INK, False)], align=PP_ALIGN.RIGHT)
        ln = plot_w*abs(v)/mx
        ex = axis_x + ln if v >= 0 else axis_x - ln
        col = NAVY if v >= 0 else NT2
        rect(s, min(axis_x, ex), cy-0.008, abs(ex-axis_x), 0.016, fill=col)
        dcol = GOLD if (gold_top and i == big) else col
        oval(s, ex-0.055, cy-0.055, 0.11, dcol)
        txt(s, ex + (0.10 if v>=0 else -0.85), cy-0.10, 0.8, 0.2,
            [(f"{v:+.1f} pp", SANS, 8.5, INK, False)], align=PP_ALIGN.LEFT if v>=0 else PP_ALIGN.RIGHT)

def pullquote(s, quote, facts, attribution="IONIC EQUITY DESK"):
    rect(s, ML, 2.02, 0.03, 2.35, fill=GOLD)
    txt(s, ML+0.30, 1.98, 5.95, 2.5, [(quote, SERIF, 18, INK, False, True)], ls=1.16)
    txt(s, ML+0.30, 4.55, 5.95, 0.24, [(attribution, SANS, 8.5, SLATE, True, False, 150)])
    # fact register right
    fx = 7.90
    fy = 1.95
    for lab, val, kind in facts:
        if kind == "bar":
            txt(s, fx, fy, 2.0, 0.24, [(lab.upper(), SANS, 8, SLATE, True, False, 120)], anchor=MSO_ANCHOR.MIDDLE)
            score_bar(s, fx+2.2, fy+0.06, val)
        elif kind == "pill":
            txt(s, fx, fy, 2.0, 0.24, [(lab.upper(), SANS, 8, SLATE, True, False, 120)], anchor=MSO_ANCHOR.MIDDLE)
            pill(s, fx+2.2, fy+0.02, val, 0.9)
        else:
            txt(s, fx, fy, 2.4, 0.24, [(lab.upper(), SANS, 8, SLATE, True, False, 120)], anchor=MSO_ANCHOR.MIDDLE)
            txt(s, fx+2.2, fy, RX-(fx+2.2), 0.24, [(val, SANS, 12, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
        rule(s, fx, fy+0.40, RX-fx, HAIR, 0.008); fy += 0.52

def paras(s, x, y, w, items, fn=SERIF, sz=11.5, col=INK, gap=8, ls=1.12, bullet=False):
    para_list = []
    for it in items:
        pref = "–  " if bullet else ""
        para_list.append([(pref+it, fn, sz, col, False)])
    txt(s, x, y, w, 4.0, para_list, ls=ls)

# ---------- divider ----------
def divider(section_no, name, statement, page_list, folio):
    s = slide(NAVYD)
    txt(s, ML, 0.5, 6, 0.3, [("IONIC", SANS, 13, WHITE, True), ("WEALTH", SANS, 13, ONNAVY, False)])
    txt(s, ML, 0.78, 6, 0.2, [("BY ANGEL ONE", SANS, 7.5, ORANGE, True, False, 200)])
    classified(s, dark=True)
    txt(s, 7.4, 0.95, 6.0, 3.6, [(f"{section_no:02d}", SERIF, 170, GHOST, False)], align=PP_ALIGN.RIGHT, wrap=False)
    txt(s, ML, 2.58, 6.0, 0.24, [(f"SECTION {section_no:02d}", SANS, 10, ORANGE, True, False, 250)])
    txt(s, ML, 2.96, 8.0, 0.7, [(name, SERIF, 30, WHITE, False)])
    rule(s, ML, 3.74, 2.2, GOLD, 0.016)
    txt(s, ML, 3.98, 6.4, 0.8, [(statement, SERIF, 13, ONNAVY, False, True)], ls=1.12)
    yy = 5.45
    for folio_i, ttl in page_list:
        txt(s, ML, yy, 7.0, 0.24, [(f"{folio_i}   ", SANS, 8.5, NT1, False), (ttl, SANS, 8.5, NT2, False)])
        yy += 0.26
    footer(s, folio, dark=True)
    return s

# =====================================================================
# PAGE 1 — COVER
# =====================================================================
s = slide(WHITE)
try:
    s.shapes.add_picture(COVER_PHOTO, Inches(6.55), Inches(0), height=Inches(CH))   # asp-locked photo, right side
except Exception:
    rect(s, 6.9, 0, CW-6.9, CH, fill=NAVYD)
pic(s, LOGO, ML, 0.82, 3.15, 0.55)
txt(s, ML, 1.62, 5.6, 0.3, [("BY ANGEL ONE", SANS, 9, SLATE, True, False, 200)])
txt(s, ML, 2.15, 5.8, 0.35, [("Co-founder", SANS, 15, ORANGE, True, True), (" in your journey of wealth creation", SANS, 15, NAVY, False)])
txt(s, CW/2-1.7, 0.14, 3.4, 0.2, [("Classified", SANS, 7.5, SLATE, False)], align=PP_ALIGN.CENTER)
rule(s, ML, 4.45, 3.0, NAVY, 0.02)
txt(s, ML, 4.62, 6, 0.9, [("Portfolio ", SANS, 40, NAVY, True), ("Review", SANS, 40, ORANGE, True)])
txt(s, ML, 5.95, 5.0, 0.24, [("PREPARED FOR", SANS, 9, SLATE, True, False, 200)])
txt(s, ML, 6.22, 5.4, 0.4, [(D["client_label"], SANS, 18, NAVY, True)])
txt(s, ML, 6.75, 5.4, 0.24, [(f"{fcr(grand)} reviewed  ·  As of 21 July 2026  ·  Ionic Wealth NDPMS Desk", SANS, 9.5, SLATE, False)])
txt(s, CW/2-1.7, 7.2, 3.4, 0.2, [("Classified as Internal", SANS, 7.5, SLATE, False)], align=PP_ALIGN.CENTER)

# =====================================================================
# PAGE 2 — CONTENTS
# =====================================================================
s = content(0, "", "02", "Contents", "What is inside")
sections = [
    ("01", "The Portfolio", ["05  The Portfolio","06  Snapshot","07  Allocation vs house view","08  Concentration & risk","09  Sector exposure","10  Market-cap positioning"]),
    ("02", "Direct Equity", ["11  Direct Equity","12  The whole book, scored","13  The top ten","14  Spotlight: Reliance","15  Spotlight: Titan","16  The names we would sell","17  What stays, and why"]),
    ("03", "Mutual Funds", ["18  Mutual Funds","19  The fund book","20  The three-year test","21  Category & AMC","22  Fund overlap","23  Three fund actions"]),
    ("04", "The Plan", ["24  The Plan","25  House-view fit","26  Tax & cost","27  Where the money moves","28  Before & after","29  Priority actions"]),
]
cx = [ML, 7.0]
cy = [1.9, 1.9]
for idx,(no,nm,items) in enumerate(sections):
    col = idx % 2
    x = cx[col]; y = cy[col]
    txt(s, x, y, 0.6, 0.4, [(no, SERIF, 22, NT3, False)])
    txt(s, x+0.62, y+0.06, 4.5, 0.4, [(nm, SERIF, 15, INK, False)])
    yy = y+0.5
    for it in items:
        f, ttl = it[:2], it[4:]
        txt(s, x+0.62, yy, 3.6, 0.22, [(ttl, SANS, 9.5, SLATE, False)])
        txt(s, x+4.2, yy, 0.7, 0.22, [(f, SANS, 9.5, INK, False)], align=PP_ALIGN.RIGHT)
        yy += 0.245
    rule(s, x, yy+0.05, 5.0, HAIR, 0.008)
    cy[col] = yy+0.35
txt(s, 7.0, cy[1], 0.6, 0.4, [("A", SERIF, 22, NT3, False)])
txt(s, 7.62, cy[1]+0.06, 4.5, 0.4, [("Appendix", SERIF, 15, INK, False)])
txt(s, 7.62, cy[1]+0.5, 4.5, 0.6, [("A-1 to A-6  Registers, methodology, disclaimer", SANS, 9.5, SLATE, False)])

# =====================================================================
# PAGE 3 — EXEC SUMMARY
# =====================================================================
s = content(0, "", "03", "Executive summary", "A sound core carrying avoidable concentration",
            "The book is well built and largely aligned to the house view; its main risk is how much sits in a few names.")
kpi_strip(s, [(fcr(grand), "Total reviewed"), (f"{T['n_stocks']} / {T['n_mf']}", "Stocks / schemes"),
              (f"{T['top10_book_pct']:.0f}%", "Top-10 weight"), (f"{T['n_eq_sell']} of {T['n_stocks']}", "Rated Sell"),
              ("3", "Fund actions")], y=1.75)
story = C.get("portfolio_story") or (D["story"].get("construct_read") if D.get("story") else "")
exec_pts = (D.get("story", {}) or {}).get("exec_summary") or []
txt(s, ML, 3.05, 5.4, 0.24, [("WHAT WE FOUND", SANS, 9, SLATE, True, False, 200)])
finds = [
    f"The top ten holdings carry {T['top10_book_pct']:.0f}% of the whole book; two names sit above 11% each.  (p. 8)",
    f"The quant scorecard rates {T['n_eq_sell']} of {T['n_stocks']} direct stocks a Sell, about {sell_pct_book:.0f}% of the portfolio by value.  (p. 12, 16)",
    "The fund book is sound: every scheme with a three-year record beats its benchmark; three actions are structural, not performance.  (p. 20, 23)",
]
paras(s, ML, 3.35, 5.4, finds, sz=12, ls=1.15)
txt(s, 6.95, 3.05, 5.4, 0.24, [("WHAT WE WOULD DO NOW", SANS, 9, SLATE, True, False, 200)])
acts = [
    f"Run the Sell programme: {T['n_eq_sell']} names, about {fcr(sell_val)} of proceeds.",
    "Trim the two >11% positions toward the single-name guideline.",
    "Switch Nippon Multi-Cap to flexi; redeem the Regular-plan ICICI Multi-Asset; exit the sub-scale Bandhan Small-Cap.",
    "Redeploy toward the house view: low-vol / value, a foreign-equity step toward 25%, and a small gold-silver sleeve.  (p. 27 to 29)",
]
paras(s, 6.95, 3.35, 5.4, acts, sz=12, ls=1.15)
source(s, "Source: Ionic Quant Scorecard, 21 Jul 2026  ·  Client statement, Jul 2026  ·  Asset X house view, Dec 2025.  † ‡ see Basis of preparation, A-5.")

# =====================================================================
# PAGE 4 — MANDATE & METHOD
# =====================================================================
s = content(0, "", "04", "Mandate & method", "How to read this review",
            "The scorecard's rules, stated before its verdicts, so every call on the pages that follow is transparent.")
txt(s, ML, 1.95, 5.2, 0.24, [("THE MANDATE", SANS, 9, SLATE, True, False, 200)])
mrows = [["Profile", ("b","Aggressive")],["Horizon", ("b","Long term")],["Construction", ("b","Core-satellite")],
         ["Objective", ("b","Long-term wealth creation")],["Benchmark", ("b","Asset X house view")]]
utbl(s, ML, 2.30, 5.2, [("",0.9,"l"),("",1.1,"r")], mrows, rowh=0.42, fs=11, header=False)
txt(s, 6.60, 1.95, 5.8, 0.24, [("THE METHOD", SANS, 9, SLATE, True, False, 200)])
paras(s, 6.60, 2.30, 5.8, [
    "The Ionic Score is a 0 to 100 composite of quality, valuation and momentum, refreshed each quarter on point-in-time data.",
    "It is read on two horizons, three-year and one-year. A name is a Sell if it falls below 40 on either; the score shown in tables is the binding, or lower, of the two.",
    "The model issues only Sell and Hold, never Buy; this is a review of existing holdings, not a solicitation. An analyst may argue a Sell up to a Hold, never a Hold down to a Sell.",
], sz=11.5, ls=1.2)
txt(s, 6.60, 5.45, 5.8, 0.24, [("THE LINE EVERY HOLDING MUST CLEAR", SANS, 8.5, SLATE, True, False, 150)])
score_bar(s, 6.60, 5.85, 40); txt(s, 7.5, 5.76, 4.6, 0.3, [("40  ·  on each horizon", SANS, 10, INK, False)])

# =====================================================================
# SECTION 01
# =====================================================================
divider(1, "The Portfolio", "Four crores thirty-three lakh, split almost evenly between conviction and delegation.",
        [("06","Snapshot"),("07","Allocation vs house view"),("08","Concentration & risk"),("09","Sector exposure"),("10","Market-cap positioning")], "05")

# PAGE 6 — SNAPSHOT
s = content(1, "The Portfolio", "06", "Snapshot", "₹4.33 crore, split almost evenly between stocks and funds",
            "Nearly half the wealth is expressed directly and half delegated to funds; the shape is balanced at the top level.")
_d = cl.donut([("Direct equity", T["eq_pct"]), ("Mutual funds", T["mf_pct"])], "p6_donut",
              [cl.NAVY, cl.NT3], center_top=fcr(grand), center_bot="TOTAL", figsize=(4.7, 4.7))
pic(s, _d, ML - 0.1, 1.9, 4.7, 4.35)
rrows = [["Direct equity", ("n",fcr(eq)), ("n",f"{T['n_stocks']} stocks")],
         ["Mutual funds", ("n",fcr(T['mf'])), ("n",f"{T['n_mf']} schemes")],
         ["Largest position", ("n",f"{STK[0]['wt_book']:.1f}%"), STK[0]['name'][:22]],
         ["Construction", "Core-satellite", ""],
         ["Stance", "Aggressive, long horizon", ""]]
utbl(s, 6.10, 2.05, 6.30, [("",1.6,"l"),("",1.2,"r"),("",1.4,"r")], rrows, rowh=0.46, fs=11, header=False)
source(s, "Source: Client statement, Jul 2026; Ionic Quant Scorecard.")

# PAGE 7 — ALLOCATION VS HOUSE VIEW
s = content(1, "The Portfolio", "07", "Allocation vs house view", "Broadly aligned with the house view, with two gaps",
            "Measured against the December house view, the equity core fits; foreign equity and a gold-silver sleeve are the gaps.")
hv = [["Domestic equity", "Incrementally positive", "Core of book, large-cap heavy", "Aligned"],
      ["Foreign equity", "~25% of equity, 60:40 DM:EM", f"~{intl_pct:.0f}% of book, DM-light", "Gap"],
      ["Gold and silver", "Positive, 75:25 gold:silver", "No dedicated sleeve", "Gap"],
      ["Momentum factor", "On hold", "Passive momentum funds held", "Watch"],
      ["Style tilt", "Low-vol and value favoured", "Value funds present", "Aligned"]]
rows = [[a, ("",b), ("",c), ("pill", d)] for a,b,c,d in hv]
utbl(s, ML, 1.95, UW, [("Dimension",1.5,"l"),("House view (Dec 2025)",2.4,"l"),("Portfolio today",2.4,"l"),("Status",1.0,"c")], rows, rowh=0.52, fs=10.5)
source(s, "Source: Asset X house view, Dec 2025; client statement, Jul 2026. Gap rows are addressed in the redeployment plan, p. 27.")

# PAGE 8 — CONCENTRATION & RISK
s = content(1, "The Portfolio", "08", "Concentration & risk", "The top ten names carry a third of the book",
            "The portfolio's dominant risk is single-name concentration, not market direction or sector bets.")
conc = D.get("concentration", {})
def cflag(v, gl, breach): return "Breach" if breach else ("Within" if not breach else "Watch")
rr = [
    ["Top-5 single names", ("n",f"{T['top5_book_pct']:.1f}%"), ("n","< 25%"), ("pill","Within" if T['top5_book_pct']<25 else "Breach")],
    ["Top-10 single names", ("n",f"{T['top10_book_pct']:.1f}%"), ("n","< 30%"), ("pill","Breach" if T['top10_book_pct']>30 else "Within")],
    ["Names above 8%", ("n","3"), ("n","0"), ("pill","Watch")],
    ["Names above 15%", ("n","0"), ("n","0"), ("pill","Within")],
    ["Largest single AMC (funds)", ("n","36.7%"), ("n","< 25%"), ("pill","Breach")],
    ["Largest sector (equity)", ("n",f"{list(D['sector_split_pct'].values())[0]:.1f}%"), ("n","< 30%"), ("pill","Breach" if list(D['sector_split_pct'].values())[0]>30 else "Within")],
]
risk_register(s, rr, x=ML, y=1.95, w=6.9)
txt(s, 8.05, 1.95, 4.35, 0.24, [("WEIGHT vs CONVICTION", SANS, 9, SLATE, True, False, 150)])
_sc = [st for st in STK if st["binding"] is not None]
_bp = cl.bubble([st["wt_book"] for st in _sc], [st["binding"] for st in _sc],
                [st["value"] for st in _sc], [callcol(st) for st in _sc], "p8_bubble",
                labels=[sym_of(st) for st in _sc], threshold=40, figsize=(4.9, 4.4),
                xlabel="Weight of book (%)", ylabel="Binding Ionic score")
pic(s, _bp, 7.95, 2.25, 4.5, 4.05)
txt(s, 8.05, 6.28, 4.35, 0.3, [("Bottom band = heavy and low-scoring: the names to act on first.", SANS, 8.5, SLATE, False, True)])
source(s, "Source: Ionic Quant Scorecard; client statement, Jul 2026. Bubble area = position value; colour = call. Guidelines are Ionic house limits.")

# PAGE 9 — SECTOR EXPOSURE (matplotlib bars, Other bucket -> foots to 100%)
ssp = list(D["sector_split_pct"].items())
top2 = ssp[0][1] + ssp[1][1]
s = content(1, "The Portfolio", "09", "Sector exposure", f"Two sectors carry {top2:.0f}% of the equity book",
            "Direct-equity weight by sector, largest first; financials and technology dominate the book.")
top9 = ssp[:9]; other = round(sum(v for _, v in ssp[9:]), 1)
slabels = [k[:24] for k, _ in top9] + (["Other sectors"] if other > 0 else [])
svals = [v for _, v in top9] + ([other] if other > 0 else [])
_p = cl.hbar(slabels, svals, "p9_sector", highlight=0, fmt="{:.1f}%", figsize=(9.8, 4.7))
pic(s, _p, ML, 1.82, UW, 4.7)
source(s, "Source: Ionic Quant Scorecard sector map. Share of the direct-equity sleeve; smaller sectors grouped as Other so the chart foots to 100%.")

# PAGE 10 — MARKET-CAP
s = content(1, "The Portfolio", "10", "Market-cap positioning", "A large-cap book: 86.7% of direct equity",
            "The portfolio is more conservative than its 68-stock count suggests; nearly nine rupees in ten sit in large caps.")
mcs = D["mcap_split_pct"]
_m3 = cl.bar3d(["Large", "Mid", "Small", "Unmapped"],
               [mcs.get("Large", 0), mcs.get("Mid", 0), mcs.get("Small", 0), mcs.get("Unknown", 0)],
               "p10_mcap3d", colors=[cl.NAVY, cl.NT1, cl.NT2, cl.NT3], figsize=(7.4, 4.2))
pic(s, _m3, ML-0.25, 1.9, 7.5, 4.2, valign="top")
obs = [
    "A large-cap core lowers portfolio volatility and liquidity risk, consistent with an aggressive but long-horizon mandate.",
    "Mid and small caps together are under 12%, leaving room to add diversified SMID exposure if the client wants more growth beta.",
    "About 2.2% of the equity book is unmapped in the source statement and is excluded from the cap analysis. †",
]
paras(s, 8.15, 2.2, 4.3, obs, sz=11.5, ls=1.22)
source(s, "Source: Ionic Quant Scorecard market-cap terciles; client statement, Jul 2026.  † unmapped lines to be resolved from the demat file.")

# =====================================================================
# SECTION 02
# =====================================================================
divider(2, "Direct Equity", "Sixty-eight names, one bar to clear: a score of forty.",
        [("12","The whole book, scored"),("13","The top ten"),("14","Spotlight: Reliance"),("15","Spotlight: Titan"),("16","The names we would sell"),("17","What stays, and why")], "11")

# PAGE 12 — WHOLE BOOK SCORED (histogram)
s = content(2, "Direct Equity", "12", "The whole book, scored", "68 stocks; 19 fall below the line",
            "The distribution of Ionic Scores across the direct-equity book, and how many sit below the Sell threshold.")
scores = [gov(s2) for s2 in STK if gov(s2) is not None]
_p = cl.histogram(scores, "p12_hist", threshold=40, figsize=(11.0, 4.4))
pic(s, _p, ML, 1.82, UW, 4.55)
txt(s, ML, 6.42, UW, 0.4,
    [(f"{T['n_eq_sell']} of {T['n_stocks']} names ({100*T['n_eq_sell']/T['n_stocks']:.0f}% by count, {sell_pct_book:.0f}% by value) score below 40.  "
      f"{DER.get('n_scored',63)} of {T['n_stocks']} carry a binding score; {DER.get('n_unscored',5)} are unrated (ETFs and holdings with no Ionic Score).",
      SERIF, 11.5, INK, False)], ls=1.1)
source(s, "Source: Ionic Quant Scorecard, binding score (lower of 3Y and 1Y), 21 Jul 2026.")

# PAGE 13 — HOLDINGS TREEMAP (hero)
s = content(2, "Direct Equity", "13", "The equity book", "Where every rupee sits, and what the scorecard says",
            "All 68 direct holdings, sized by value and coloured by call; ten names carry a third of the book and the Sells are visible in it.")
_lab = [sym_of(st)[:12] for st in STK]
_sz = [max(st["wt_book"], 0.06) for st in STK]
_col = [callcol(st) for st in STK]
_vl = [f"{st['wt_book']:.1f}%" for st in STK]
_p = cl.treemap(_lab, _sz, "p13_treemap", colors=_col, figsize=(11.6, 4.6), value_labels=_vl)
pic(s, _p, ML, 1.74, UW, 4.5, valign="top")
rect(s, ML, 6.34, 0.16, 0.16, fill=HOLD)
txt(s, ML+0.24, 6.31, 3.2, 0.22, [(f"Hold  ·  {DER['hold_n']} names", SANS, 9.5, INK, False)])
rect(s, ML+2.1, 6.34, 0.16, 0.16, fill=SELL)
txt(s, ML+2.34, 6.31, 3.2, 0.22, [(f"Sell  ·  {DER['sell_n']} names", SANS, 9.5, INK, False)])
txt(s, RX-4.0, 6.31, 4.0, 0.22, [(f"Top 10 = {T['top10_book_pct']:.1f}% of the book", SANS, 9.5, SLATE, False)], align=PP_ALIGN.RIGHT)
source(s, "Source: Ionic Quant Scorecard. Tile area = share of the total portfolio; colour = call.")

# PAGE 14 — SPOTLIGHT RELIANCE
def find_stock(sym):
    return next((x for x in STK if (x["symbol"] or "").upper()==sym), None)
rel = find_stock("RELIANCE")
s = content(2, "Direct Equity", "14", "Spotlight  ·  Reliance", "The second-largest holding scores 27 of 100",
            "Position size is not a thesis; here is where the book's biggest gap between conviction and evidence sits.", gold_eyebrow=False)
rq = (CST.get("RELIANCE",{}) or {}).get("text") or "At 11.0% of the book and a score of 27, Reliance is the portfolio's largest gap between conviction and evidence. We would act on the Sell by trimming hard toward a market weight and redeploying into higher-scoring names."
pullquote(s, clip_two(rq), [
    ("Weight", f"{rel['wt_book']:.1f}%" if rel else "11.0%", "text"),
    ("Three-year score", rel['score3y'] if rel else 27, "bar"),
    ("One-year score", rel['score1y'] if rel else 32, "bar"),
    ("Call", "Sell", "pill"),
])
def radar_pillars(st, name, color):
    pl = st["pillars"] if st else {}
    cats = ["Quality", "Growth 3Y", "Growth 1Y", "Value", "Trend", "Macro"]
    vals = [pl.get(k) or 0 for k in ("quality", "growth3y", "growth1y", "value", "stage3y", "macro")]
    return cl.radar(cats, vals, name, color=color, figsize=(3.4, 3.4))
txt(s, 7.95, 4.12, 4.4, 0.24, [("THE SCORE, DECOMPOSED", SANS, 8.5, SLATE, True, False, 150)])
pic(s, radar_pillars(rel, "p14_radar", cl.SELL), 8.55, 4.30, 2.75, 2.75)
txt(s, ML+0.30, 4.95, 5.95, 0.3, [("Proposed: reduce toward a market weight; see order sheet, p. 29.", SANS, 9.5, SLATE, False, True)])
source(s, "Source: Ionic Quant Scorecard, 21 Jul 2026.")

# PAGE 15 — SPOTLIGHT TITAN
tit = find_stock("TITAN")
s = content(2, "Direct Equity", "15", "Spotlight  ·  Titan", "The largest holding, held on merit and watched on size",
            "A Hold is a judgment, not a pass; the question on the book's biggest position is how much, not whether.", gold_eyebrow=False)
tq = (CST.get("TITAN",{}) or {}).get("text") or "Titan clears the bar at 49; at 11.2% of the book, the question is not whether to own it but how much. The franchise keeps compounding, so we hold and watch the weight rather than the thesis."
pullquote(s, clip_two(tq), [
    ("Weight", f"{tit['wt_book']:.1f}%" if tit else "11.2%", "text"),
    ("Three-year score", tit['score3y'] if tit else 49, "bar"),
    ("One-year score", tit['score1y'] if tit else 44, "bar"),
    ("Call", "Hold", "pill"),
])
txt(s, 7.95, 4.12, 4.4, 0.24, [("THE SCORE, DECOMPOSED", SANS, 8.5, SLATE, True, False, 150)])
pic(s, radar_pillars(tit, "p15_radar", cl.HOLD), 8.55, 4.30, 2.75, 2.75)
txt(s, ML+0.30, 4.95, 5.95, 0.3, [("Proposed: trim toward the single-name guideline; see order sheet, p. 29.", SANS, 9.5, SLATE, False, True)])
source(s, "Source: Ionic Quant Scorecard, 21 Jul 2026.")

# PAGE 16 — THE 19 SELLS
s = content(2, "Direct Equity", "16", "The names we would sell", "The nineteen names we would sell",
            "The full Sell programme, shown in one register; nothing is buried and each carries its primary reason.")
pill(s, RX-1.15, 1.42, f"SELL x{T['n_eq_sell']}", 1.1, kind="Sell")
def short_reason(st):
    pl = st["pillars"]; items = [(k, pl.get(k)) for k in ("value","growth1y","stage3y","quality","growth3y","macro") if pl.get(k) is not None]
    items.sort(key=lambda t: t[1])
    names = {"value":"rich valuation","growth1y":"soft recent growth","stage3y":"weak price trend","quality":"quality below peers","growth3y":"slowing growth","macro":"sector headwind"}
    return "; ".join(names[k] for k,_ in items[:2]) if items else "below the score line"
rows = []
for i, st in enumerate(sell_stocks[:19]):
    rows.append([("n",str(i+1)), ("b",st["name"][:22]), ("n",f"{st['wt_book']:.1f}%"), ("bar", gov(st)), short_reason(st)])
utbl(s, ML, 1.74, UW, [("#",0.35,"l"),("Name",2.5,"l"),("Wt",0.75,"r"),("Binding score",1.7,"l"),("Primary reason",4.0,"l")],
     rows, rowh=0.205, fs=9,
     totals=[None, f"{len(sell_stocks)} names", f"{sell_pct_book:.1f}%", None, "est. proceeds  "+fcr(sell_val)])
source(s, "Source: Ionic Quant Scorecard. Reason is the weakest scoring pillar(s); full pillar detail in Appendix A-2/A-3.")

# PAGE 17 — WHAT STAYS
s = content(2, "Direct Equity", "17", "What stays, and why", "A Hold is a decision, not an absence of one",
            "The 49 names we would keep, and what it would take to change our mind on the closest calls.")
holds = [x for x in STK if x["call"]=="Hold"]
watch = sorted([x for x in holds if gov(x) is not None and 40 <= gov(x) <= 46], key=lambda x:-x["wt_book"])[:8]
txt(s, ML, 1.95, 6.4, 0.24, [("ON THE WATCHLIST  ·  BINDING SCORE 40 TO 46", SANS, 9, SLATE, True, False, 150)])
rows = [[("b",x["name"][:22]), ("n",f"{x['wt_book']:.1f}%"), ("bar", gov(x))] for x in watch]
utbl(s, ML, 2.30, 6.2, [("Name",2.4,"l"),("Weight",1.0,"r"),("Ionic score",2.0,"l")], rows, rowh=0.36, fs=10)
txt(s, 7.70, 1.95, 4.7, 3.8, [
    ("A Hold is a decision, not an absence of one. Of the 49 names we keep, most clear the line comfortably; a handful sit just above it and stay on the watchlist.\n\n"
     "We would move a Hold to a Sell if its score falls below 40 at the next quarterly refresh, if a governance or accounting flag appears, or if the position grows past the single-name guideline on price alone.\n\n"
     "The book is reviewed every quarter; nothing here is permanent.", SERIF, 12, INK, False)], ls=1.22)
source(s, "Source: Ionic Quant Scorecard. Watchlist = Hold names within 6 points of the Sell threshold.")

# =====================================================================
# SECTION 03
# =====================================================================
divider(3, "Mutual Funds", "Two crores across seventeen schemes; only four are old enough to judge.",
        [("19","The fund book"),("20","The three-year test"),("21","Category & AMC"),("22","Fund overlap"),("23","Three fund actions")], "18")

# PAGE 19 — FUND BOOK AT A GLANCE
s = content(3, "Mutual Funds", "19", "The fund book", "₹2.00 crore across seventeen schemes",
            "The shape and, importantly, the age of the delegated book; only the seasoned cohort can be judged on alpha.")
kpi_strip(s, [(fcr(T['mf']),"Fund assets"),(str(T['n_mf']),"Schemes"),(str(n_tested),"With 3Y record"),
              (str(len(D['mf_amc_pct'])),"Fund houses")], y=1.85)
# category bar
mcp = D["mf_cat_pct"]
txt(s, ML, 3.35, UW, 0.24, [("BY CATEGORY", SANS, 9, SLATE, True, False, 150)])
ramp = [NAVY,NT1,NT2,NT3,NAVY,NT1,NT2,NT3,NAVY,NT1,NT2]
segs = [(k, v, ramp[i%len(ramp)]) for i,(k,v) in enumerate(mcp.items())]
stacked100(s, ML, 3.95, UW, 0.7, segs, label_above=False)
# legend below
lx = ML; ly = 4.85; per_row = 4; i=0
for (k,v),col in zip(mcp.items(), [ramp[j%len(ramp)] for j in range(len(mcp))]):
    col_i = i % per_row; row_i = i // per_row
    xx = ML + col_i*3.0; yy = ly + row_i*0.30
    rect(s, xx, yy+0.03, 0.14, 0.14, fill=col)
    txt(s, xx+0.22, yy, 2.7, 0.22, [(f"{k}  {v:.1f}%", SANS, 8.5, INK, False)])
    i += 1
txt(s, ML, 6.15, UW, 0.3, [("Only four schemes carry a three-year record; the rest are too young, passive, international or cash-like to score on alpha. †", SERIF, 11, INK, False, True)])
source(s, "Source: client statement, Jul 2026; QFRA fund NAV feed.  † see Basis of preparation, A-5.")

# PAGE 20 — THREE-YEAR TEST
s = content(3, "Mutual Funds", "20", "The three-year test", "Every fund with a three-year record has beaten its benchmark",
            "The seasoned cohort passes cleanly; the measurement window is stated honestly rather than stretched.")
tested = [f for f in MF if f["alpha"] is not None]
py, pb = 2.05, 5.7; ph = pb-py
bw = 0.55
n = len(tested); slot = 8.0/n
for i,f in enumerate(tested):
    cx = ML + i*slot + slot/2
    mx = 40.0
    hb = ph*(f["bench_cagr"] or 0)/mx; hf = ph*(f["cagr"] or 0)/mx
    rect(s, cx-bw-0.05, pb-hb, bw, hb, fill=NT3)
    rect(s, cx+0.05, pb-hf, bw, hf, fill=NAVY)
    txt(s, cx+0.05, pb-hf-0.22, bw, 0.2, [(f"{f['cagr']:.1f}%", SANS, 8.5, INK, False)], align=PP_ALIGN.CENTER)
    txt(s, cx-bw-0.05, pb-hb-0.22, bw, 0.2, [(f"{f['bench_cagr']:.0f}", SANS, 8, SLATE, False)], align=PP_ALIGN.CENTER)
    nm = f["scheme"].replace(" - Direct Plan","").replace("Motilal Oswal","MO").replace(" Fund","")
    txt(s, cx-slot/2+0.1, pb+0.06, slot-0.2, 0.5, [(nm[:22], SANS, 8, SLATE, False)], align=PP_ALIGN.CENTER, ls=1.0)
    txt(s, cx-0.4, pb-max(hf,hb)-0.46, 0.8, 0.2, [(f"+{f['alpha']:.1f}", SANS, 8.5, NT1, True)], align=PP_ALIGN.CENTER)
rule(s, ML, pb, 8.0, SLATE, 0.012)
txt(s, ML, 1.95, 8.0, 0.22, [("Navy = fund 3Y CAGR   ·   grey = category benchmark   ·   +alpha in pp", SANS, 8, SLATE, False)])
rect(s, 9.0, 1.95, 3.4, 3.9, fill=PANEL)
txt(s, 9.18, 2.12, 3.05, 3.6, [
    [("BASIS", SANS, 8.5, SLATE, True, False, 150)],
    [("Fund alpha is measured to 31 December 2024, the most recent date with audited NAV histories common to all four schemes. It will be refreshed in the Q3 review. †", SERIF, 10.5, INK, False)],
    [("No fund is rated Sell on performance. The three actions on p. 23 are structural.", SERIF, 10.5, INK, False)],
], ls=1.2)
source(s, "Source: QFRA fund NAV feed, direct-plan NAVs.  † window ends Dec 2024; see A-5.")

# PAGE 21 — CATEGORY & AMC
s = content(3, "Mutual Funds", "21", "Category & AMC", "Concentration by fund house and category",
            "Seventeen schemes is not seventeen bets; one fund house holds over a third of the fund book.")
amc = D["mf_amc_pct"]
txt(s, ML, 1.95, 5.4, 0.24, [("BY FUND HOUSE", SANS, 9, SLATE, True, False, 150)])
ay = 2.35; maxa = max(amc.values())
for i,(k,v) in enumerate(list(amc.items())[:8]):
    cy = ay + i*0.42
    txt(s, ML, cy-0.02, 1.7, 0.24, [(k, SERIF, 10, INK, False)])
    bw2 = 3.0*v/maxa
    rect(s, ML+1.8, cy+0.02, bw2, 0.18, fill=(NAVY if i==0 else NT3))
    txt(s, ML+1.8+bw2+0.06, cy-0.02, 0.8, 0.24, [(f"{v:.1f}%", SANS, 9, INK, False)])
# 25% guideline line
gx = ML+1.8 + 3.0*25/maxa
vrule(s, gx, ay, min(8,len(amc))*0.42, GOLD, 0.014)
txt(s, gx-0.5, ay-0.24, 1.4, 0.2, [("25% guideline", SANS, 7.5, INK, True)], align=PP_ALIGN.CENTER)
txt(s, 6.95, 1.95, 5.4, 0.24, [("BY CATEGORY", SANS, 9, SLATE, True, False, 150)])
crows = [[k, ("n",f"{v:.1f}%")] for k,v in list(D["mf_cat_pct"].items())[:9]]
utbl(s, 6.95, 2.35, 5.4, [("Category",2.2,"l"),("Share",1.0,"r")], crows, rowh=0.34, fs=10, header=False)
source(s, "Source: client statement, Jul 2026. AMC guideline is the Ionic single-manager limit.")

# PAGE 22 — FUND OVERLAP
s = content(3, "Mutual Funds", "22", "Fund overlap", "How much the funds overlap",
            "Where mandates overlap, holding many schemes adds cost without adding diversification.")
ov = D.get("overlap", {})
txt(s, ML, 2.1, 7.0, 3.4, [
    (ov.get("narrative") or "A full stock-level look-through needs each fund's latest portfolio. On the data available, the clearest duplication is structural: two small-cap funds and several passive index funds that track overlapping baskets. The action list already removes the sub-scale small-cap and the duplicated multi-asset plan.", SERIF, 12.5, INK, False)], ls=1.25)
md = ov.get("most_duplicated") or []
if md:
    txt(s, ML, 4.9, 6.0, 0.24, [("MOST DUPLICATED EXPOSURES", SANS, 9, SLATE, True, False, 150)])
    txt(s, ML, 5.2, 6.0, 0.6, [(", ".join(str(x) for x in md[:8]), SANS, 10, INK, False)])
rect(s, 8.4, 2.1, 4.0, 3.4, fill=PANEL)
txt(s, 8.6, 2.3, 3.6, 3.1, [
    [("TAKEAWAY", SANS, 8.5, SLATE, True, False, 150)],
    [("Overlap here is mild and mostly by design (index funds). The redundancy worth removing is the sub-scale second small-cap and the Regular-plan duplicate, both already on the action list, p. 23.", SERIF, 11, INK, False)]], ls=1.2)
source(s, "Source: QFRA fund holdings where available; full look-through pending each fund's latest portfolio. †")

# PAGE 23 — THREE FUND ACTIONS
s = content(3, "Mutual Funds", "23", "Three fund actions", "Three actions in the fund book",
            "Each action is structural, not a reaction to performance; the reasons are mandate shape, plan cost and scale.")
acts = [
    ("Switch","Nippon India Multi Cap Fund", "Multi-cap forces a fixed small and mid cap quota we prefer to size ourselves; its record is strong, so this is a switch, not a criticism.", "Move to flexi-cap  ·  p. 26"),
    ("Redeem","ICICI Prudential Multi Asset (Regular)", "The same fund is already held in the direct plan; there is no reason to pay the Regular-plan trail.", "Consolidate into Direct  ·  p. 26"),
    ("Exit","Bandhan Small Cap Fund", "A sub-scale position of about ₹3 lakh beside a ₹73 lakh small-cap fund; it adds cost and duplication.", "Let the primary small-cap carry the sleeve"),
]
ay = 1.95
for act, nm, rat, dest in acts:
    rect(s, ML, ay, UW, 1.35, fill=None, line=HAIR, lw=1.0)
    pill(s, ML+0.2, ay+0.2, act, 0.9)
    txt(s, ML+0.2, ay+0.55, 3.0, 0.6, [(nm, SERIF, 13, INK, False)], ls=1.0)
    txt(s, ML+3.5, ay+0.22, 5.2, 1.0, [(rat, SERIF, 11, INK, False)], ls=1.18)
    txt(s, ML+8.9, ay+0.22, 2.5, 1.0, [[("DESTINATION", SANS, 7.5, SLATE, True, False, 120)],[(dest, SANS, 9.5, INK, False)]], ls=1.15)
    ay += 1.5
source(s, "Source: Ionic fund desk. No fund is sold on performance; all three actions are structural.")

# =====================================================================
# SECTION 04
# =====================================================================
divider(4, "The Plan", "Nineteen sales, three fund actions, one order sheet.",
        [("25","House-view fit"),("26","Tax & cost"),("27","Where the money moves"),("28","Before & after"),("29","Priority actions")], "24")

# PAGE 25 — HOUSE-VIEW FIT
s = content(4, "The Plan", "25", "House-view fit", "The plan closes the gaps against the house view",
            "Every action maps back to the December house view; the redeployment is where the two gaps get closed.")
rows = [
    ["Domestic equity", "Incrementally positive", "Aligned; trim concentration", ("pill","Aligned")],
    ["Foreign equity", "~25% of equity, 60:40", "Step up via redeployment", ("pill","Gap")],
    ["Gold and silver", "Positive, 75:25", "Add a small sleeve", ("pill","Gap")],
    ["Momentum", "On hold", "No new momentum adds", ("pill","Aligned")],
    ["Low-vol / value", "Favoured", "Redeploy Sell proceeds here", ("pill","Aligned")],
]
utbl(s, ML, 1.95, UW, [("Dimension",1.5,"l"),("House view (Dec 2025)",2.3,"l"),("What the plan does",3.0,"l"),("Fit",1.0,"c")], rows, rowh=0.52, fs=10.5)
source(s, "Source: Asset X house view, Dec 2025; Ionic redeployment plan, p. 27.")

# PAGE 26 — TAX & COST
s = content(4, "The Plan", "26", "Tax & cost", "What the actions trigger, and what we cannot yet compute",
            "Honest tax arithmetic on the fund actions, and a clear note on what the statement does not yet let us compute.")
trows = [
    ["Nippon Multi-Cap (switch)", ("n",fl(next((f['value'] for f in MF if 'Nippon India Multi Cap' in f['scheme']),0))), "Held > 1y", "LTCG", "Switch is a redemption for tax"],
    ["ICICI Multi-Asset Regular (redeem)", ("n",fl(reg_val)), "Recent", "STCG likely", "Consolidate to Direct; TER saving"],
    ["Bandhan Small-Cap (exit)", ("n",fl(next((f['value'] for f in MF if 'Bandhan Small Cap' in f['scheme']),0))), "Held > 1y", "LTCG", "Small ticket, low tax drag"],
]
utbl(s, ML, 1.95, UW, [("Action",3.0,"l"),("Amount ₹ L",1.3,"r"),("Holding",1.2,"l"),("Tax",1.2,"l"),("Note",3.2,"l")], rows=[[a,b,c,d,e] for a,b,c,d,e in trows], rowh=0.44, fs=10)
rect(s, ML, 4.55, UW, 1.5, fill=PANEL)
txt(s, ML+0.2, 4.72, UW-0.4, 1.2, [
    [("BASIS OF PREPARATION", SANS, 8.5, SLATE, True, False, 150)],
    [("‡ Acquisition dates and costs for the 68 direct-equity holdings are not in the statement provided, so tax estimates cover the fund book only. Share the demat trade file and the next review will extend this to all direct-equity sells.", SERIF, 11, INK, False)],
    [("Moving the ICICI Multi-Asset holding from Regular to Direct saves the ongoing Regular-plan trail every year the position is held.", SERIF, 11, INK, False)]], ls=1.18)
source(s, "Source: client statement, Jul 2026. Tax character indicative; confirm with the client's tax adviser before dealing.")

# PAGE 27 — WHERE THE MONEY MOVES (bridge)
s = content(4, "The Plan", "27", "Where the money moves", "Where the sale proceeds go",
            "The redeployment as a single flow: proceeds from the Sell programme and fund actions, into house-view sleeves.")
d1 = proceeds*0.45; d2 = proceeds*0.25; d3 = proceeds*0.15; d4 = proceeds - d1 - d2 - d3
steps = [
    ("Proceeds\nsells + trim + fund actions", proceeds, "open"),
    ("Low-vol / value\ncore add", d1, "flow"),
    ("Foreign equity\ntoward 25%", d2, "flow"),
    ("Gold & silver\n75:25 sleeve", d3, "flow"),
    ("Retained cash\nbuffer", d4, "close"),
]
_w = cl.waterfall(steps, "p27_waterfall", figsize=(11.2, 4.4), gold_idx=2)
pic(s, _w, ML, 1.92, UW, 4.5)
txt(s, ML, 6.44, UW, 0.3, [("Indicative split across house-view sleeves; freed cash is shown as cash until redeployed, never assumed fully invested.", SERIF, 11, INK, False, True)])
source(s, f"Source: Ionic redeployment plan. Proceeds of {fcr(proceeds)} from the 19 Sells, the trim and the 3 fund actions (matches the order sheet, p. 29). Amounts indicative.")

# PAGE 28 — BEFORE & AFTER
s = content(4, "The Plan", "28", "Before & after", "The book, before and after the plan",
            "The improvement in one glance: less concentration, no sub-scale or Regular-plan clutter, a step toward the house view.")
lg = STK[0]['wt_book']
tor = [
    ("Largest single name", lg, f"{lg:.1f}%", 8.0, "8.0% (guideline)"),
    ("Sell-rated weight", sell_pct_book, f"{sell_pct_book:.1f}%", 0.0, "0.0%"),
    ("Regular-plan fund", 100*reg_val/grand, f"{100*reg_val/grand:.1f}%", 0.0, "0.0%"),
    ("Foreign equity share", intl_pct, f"{intl_pct:.1f}%", 25.0, "toward 25%"),
]
tornado(s, tor, spine=6.7, y0=2.35, pitch=0.78)
txt(s, ML, 5.85, UW, 0.5, [
    [("Fund schemes reduce from 17 to 15 (exit Bandhan, redeem the Regular-plan duplicate; the Nippon switch is net neutral).", SERIF, 11, INK, False)],
    [("Proposed values are guidelines or house-view targets, not forecasts of return.", SERIF, 10.5, SLATE, False, True)]], ls=1.2)
source(s, "Source: Ionic Quant Scorecard; redeployment plan. Targets are house guidelines.")

# PAGE 29 — ORDER SHEET
s = content(4, "The Plan", "29", "Priority actions", "Priority actions",
            "Everything in this review, on one page the client can execute from, sequenced tax-aware.")
lines = [
    ("Trim", STK[0]['name'][:20], "Toward single-name guideline", STK[0]['value']*0.3, "Low-vol / value", "p.15"),
    ("Sell", STK[1]['name'][:20], "Full exit, scored 27", STK[1]['value'], "Low-vol / value", "p.14"),
    ("Sell", f"{T['n_eq_sell']-1} further Sell names", "Scored below 40", sell_val-STK[1]['value'], "Redeploy per plan", "p.16"),
    ("Switch", "Nippon Multi-Cap", "Mandate shape to flexi-cap", next((f['value'] for f in MF if 'Nippon' in f['scheme']),0), "Flexi-cap fund", "p.23"),
    ("Redeem", "ICICI Multi-Asset (Regular)", "Duplicate, costlier plan", reg_val, "ICICI Direct plan", "p.23"),
    ("Exit", "Bandhan Small-Cap", "Sub-scale duplicate", next((f['value'] for f in MF if 'Bandhan' in f['scheme']),0), "Primary small-cap", "p.23"),
]
rows = []
tot = 0
for i,(act,ins,instr,amt,dest,ref) in enumerate(lines):
    tot += amt
    rows.append([("n",str(i+1)), ("pill",act), ("b",ins), instr, ("n",fl(amt)), dest, ("n",ref)])
utbl(s, ML, 1.95, UW, [("#",0.3,"l"),("Action",0.95,"c"),("Instrument",2.5,"l"),("Instruction",2.9,"l"),("₹ L",0.85,"r"),("Proceeds to",2.0,"l"),("Ref",0.6,"r")],
     rows, rowh=0.42, fs=9.5,
     totals=[None,None,"Total redeployed", None, ("n",fl(tot)), None, None])
txt(s, ML, 6.35, UW, 0.3, [("Reviewed with client on  ____________________", SANS, 9.5, SLATE, False)])
rule(s, ML, 6.62, 4.0, HAIR, 0.01)
source(s, "Source: this review. Sequence sells tax-aware per p. 26; amounts indicative pending dealing.")

# =====================================================================
# QUANT LENS (new)
# =====================================================================
# PAGE 30 — EFFICIENT FRONTIER
s = content(0, "The opportunity set", "30", "The opportunity set", "Where the book sits against an efficient frontier")
_debt = sum(f["value"] for f in MF if f["category"] and ("Hybrid" in f["category"] or "Arbitrage" in f["category"]) and f["value"])
cur_intl = intl_pct; cur_debt = round(100*_debt/grand, 1); cur_dom = round(100 - cur_intl - cur_debt, 1)
_ef = cl.efficient_frontier(["Domestic equity", "International", "Debt & cash", "Gold"],
        [11.5, 9.5, 7.0, 8.0], [16, 15, 4, 14],
        [[1, .6, .1, 0], [.6, 1, .1, .1], [.1, .1, 1, -.1], [0, .1, -.1, 1]],
        [("Today", [cur_dom, cur_intl, cur_debt, 0.1], cl.SELL), ("House-view target", [60, 20, 12, 8], cl.HOLD)],
        "p30_ef", figsize=(8.4, 5.2))
pic(s, _ef, ML, 1.78, 7.7, 4.6, valign="top")
txt(s, 9.0, 1.9, 3.4, 0.24, [("READING THE MAP", SANS, 9, SLATE, True, False, 150)])
txt(s, 9.0, 2.25, 3.4, 3.8, [
    [("Each dot is a feasible asset-class mix; colour rises with reward-per-unit-risk. The upper-left edge is the efficient frontier.", SERIF, 11, INK, False)],
    [("Today (red) sits inside the frontier: too much single-basket equity risk for its return. The house-view target (green) steps left toward the frontier by adding international, gold and a debt cushion.", SERIF, 11, INK, False)]], ls=1.22)
source(s, "Illustrative opportunity set from long-run capital-market assumptions, not a forecast. Client return history was not used; this maps asset-class mixes only.")

# PAGE 31 — VALUE MAP
s = content(0, "Quality vs price", "31", "Quality versus price", "Are we paying fair prices for the quality we own?")
_vm_s = [st for st in STK if st["roe"] is not None and st["pe"] is not None and 0 < st["pe"] < 120 and -30 < st["roe"] < 120]
_vm = cl.value_map([st["pe"] for st in _vm_s], [st["roe"] for st in _vm_s], [st["value"] for st in _vm_s],
                   [callcol(st) for st in _vm_s], [sym_of(st) for st in _vm_s], "p31_vm", figsize=(8.4, 5.2))
pic(s, _vm, ML, 1.78, 7.7, 4.65, valign="top")
txt(s, 9.0, 1.9, 3.4, 0.24, [("THE FOUR QUADRANTS", SANS, 9, SLATE, True, False, 150)])
txt(s, 9.0, 2.25, 3.4, 3.8, [
    [("Every holding on quality (ROE) against valuation (P/E); bubble size is position value, colour is the call.", SERIF, 11, INK, False)],
    [("The names to question sit bottom-right: expensive for the quality they deliver. Green Holds cluster top-left, high quality at a fair price; red Sells lean expensive-and-mediocre.", SERIF, 11, INK, False)]], ls=1.22)
source(s, f"Source: Ionic Quant Scorecard (ROE, P/E on TTM). {len(STK)-len(_vm_s)} names with no or extreme P/E omitted.")

# PAGE 32 — PORTFOLIO FACTOR PROFILE
s = content(0, "Factor profile", "32", "The book's factor profile", "How the equity book scores across the six quant pillars")
_keys = ["quality", "growth3y", "growth1y", "value", "stage3y", "macro"]
_acc = {k: 0.0 for k in _keys}; _w = 0.0
for st in STK:
    if st["binding"] is None: continue
    ww = st["wt_book"] or 0
    for k in _keys:
        v = st["pillars"].get(k)
        if v is not None: _acc[k] += v * ww
    _w += ww
_vals = [round(_acc[k]/_w) if _w else 0 for k in _keys]
_fr = cl.radar(["Quality", "Growth 3Y", "Growth 1Y", "Value", "Trend", "Macro"], _vals, "p32_factor",
               color=cl.NAVY, figsize=(4.7, 4.7))
pic(s, _fr, ML+0.3, 1.95, 4.7, 4.4)
_pairs = sorted(zip(["Quality", "Growth (3Y)", "Growth (1Y)", "Value", "Price trend", "Sector / macro"], _vals), key=lambda t: t[1])
txt(s, 7.4, 2.0, 4.9, 0.24, [("WHAT THE PROFILE SAYS", SANS, 9, SLATE, True, False, 150)])
txt(s, 7.4, 2.35, 4.9, 3.6, [
    [(f"Weighted across the equity book, the strongest pillar is {_pairs[-1][0].lower()} ({_pairs[-1][1]:.0f}) and the weakest is {_pairs[0][0].lower()} ({_pairs[0][1]:.0f}).", SERIF, 11.5, INK, False)],
    [("A value-light, trend-led profile is consistent with an aggressive, large-cap growth book; it also explains why a stretch in valuation is the most common reason names fall below the Sell line.", SERIF, 11.5, INK, False)]], ls=1.24)
source(s, "Source: Ionic Quant Scorecard pillar scores, weighted by position value across the scored equity book (0-100 each).")

# PAGE 33 — WEALTH PROJECTION (goal planning)
s = content(0, "Growth projection", "33", "Growth projection", "Where this book could grow, and the goals it can fund")
_proj = cl.projection_cone(grand, 15, 11.0, 13.0, "p33_proj",
        goals=[(10, "Illustrative goal corpus", grand*2.2), (15, "Illustrative long-term target", grand*4.0)],
        figsize=(9.0, 5.0))
pic(s, _proj, ML, 1.78, 7.8, 4.7, valign="top")
txt(s, 9.05, 1.95, 3.35, 0.24, [("READING THE CONE", SANS, 9, SLATE, True, False, 150)])
txt(s, 9.05, 2.3, 3.35, 3.9, [
    [("At an assumed 11% return and 13% volatility, the book's central path roughly doubles in about seven years.", SERIF, 11, INK, False)],
    [("The shaded band is the 10th to 90th percentile of outcomes; real experience will vary. The dashed lines are placeholder goal corpora.", SERIF, 11, INK, False)],
    [("Share your actual goals and target dates and we will size each funding gap and set a glidepath.", SERIF, 11, INK, False)]], ls=1.22)
source(s, "Illustrative lognormal projection at long-run capital-market assumptions, not a forecast. Goal lines are placeholders pending your goal inputs.")

# =====================================================================
# HOLDINGS IN DETAIL (all 68, with 2-8 line commentary)
# =====================================================================
divider(5, "Holdings in detail", "Every direct-equity holding, with the desk note and the scorecard's call.", [], "34")

def note_for(st):
    for key in (sym_of(st), st.get("symbol") or "", st.get("name") or ""):
        n = CST.get(key)
        if isinstance(n, dict) and n.get("text"): return n["text"]
    return "Retained and monitored; no separate scorecard view for this holding."

def holding_card(s, x, y, w, h, st):
    rect(s, x, y, w, h, fill=PANEL)
    rect(s, x, y, 0.05, h, fill=(SELL if st["call"] == "Sell" else HOLD))
    txt(s, x+0.2, y+0.13, w-1.25, 0.3, [(sym_of(st)[:24], SANS, 12.5, NAVY, True)])
    pill(s, x+w-1.12, y+0.16, st["call"] if st["call"] in ("Sell", "Hold") else "Hold", 0.94)
    sub = f"{st['wt_book']:.1f}% of book"
    if st["binding"] is not None: sub += f"   ·   score {st['binding']:.0f}"
    if st["sector"]: sub += f"   ·   {st['sector'][:20]}"
    txt(s, x+0.2, y+0.46, w-0.4, 0.22, [(sub, SANS, 8.5, SLATE, False)])
    txt(s, x+0.2, y+0.73, w-0.4, h-0.83, [[(note_for(st), SANS, 9, INK, False)]], ls=1.12)

_STKsort = sorted(STK, key=lambda z: -(z["wt_book"] or 0))
_pos = [(ML, 1.72), (ML+5.8, 1.72), (ML, 4.12), (ML+5.8, 4.12)]
_npg = math.ceil(len(_STKsort)/4)
for _pi in range(_npg):
    _chunk = _STKsort[_pi*4:_pi*4+4]
    _lo, _hi = _pi*4+1, min(_pi*4+4, len(_STKsort))
    s = content(5, "Holdings in detail", f"H{_pi+1}", "Holdings in detail", f"Direct equity, {_lo} to {_hi} of {len(_STKsort)} by weight")
    for st, (cx, cy) in zip(_chunk, _pos):
        holding_card(s, cx, cy, 5.6, 2.25, st)
    source(s, "Source: Ionic Quant Scorecard and Ionic equity desk. Weight = share of the total portfolio; score = binding (lower of 3Y and 1Y).")

# =====================================================================
# APPENDIX
# =====================================================================
divider(0, "Appendix", "Every holding, every score, every source.", [], "A-1")
prs.slides[-1].shapes  # divider used section 0 -> ghost "00"; acceptable

def appendix_equity(page_rows, folio, part):
    s = content(0, "Appendix", folio, "Appendix", f"Direct equity register  ·  {part}", gold_eyebrow=True)
    half = (len(page_rows)+1)//2
    L = page_rows[:half]; R = page_rows[half:]
    def block(rows, x):
        rr = [[("b",st["name"][:20]), ("n",f"{st['wt_book']:.1f}%"), ("bar", gov(st)), ("pill", st["call"] if st["call"] in ("Sell","Hold") else "Hold")] for st in rows]
        utbl(s, x, 1.95, 5.36, [("Name",2.3,"l"),("Wt",0.8,"r"),("Score",1.6,"l"),("Call",0.9,"c")], rr, rowh=0.245, fs=8.5)
    block(L, ML); block(R, 6.85)
    source(s, "Source: Ionic Quant Scorecard, 21 Jul 2026. Sorted by weight.")
    return s

appendix_equity(STK[:34], "A-2", "1 of 2")
appendix_equity(STK[34:68], "A-3", "2 of 2")

# A-4 fund register
s = content(0, "Appendix", "A-4", "Appendix", "Mutual fund register")
frows = []
for f in MF:
    plan = "Regular" if (f["scheme"]=="ICICI Prudential Multi Asset Fund") else "Direct"
    perf = (f"+{f['alpha']:.1f} pp vs bench" if f["alpha"] is not None else "under 3Y")
    frows.append([("b",f["scheme"].replace(" - Direct Plan","")[:30]), (f["category"] or "")[:16], plan, ("n",fl(f["value"]) if f["value"] else "-"), perf, ("pill", f["action"])])
utbl(s, ML, 1.95, UW, [("Scheme",3.3,"l"),("Category",1.7,"l"),("Plan",0.9,"l"),("₹ L",1.0,"r"),("3Y vs benchmark",2.0,"l"),("Action",1.0,"c")], frows, rowh=0.28, fs=8.5)
source(s, "Source: client statement, Jul 2026; QFRA fund NAV feed. Alpha window to Dec 2024.")

# A-5 methodology
s = content(0, "Appendix", "A-5", "Appendix", "Methodology & basis of preparation")
txt(s, ML, 1.95, 5.6, 0.24, [("THE IONIC SCORE", SANS, 9, SLATE, True, False, 150)])
paras(s, ML, 2.30, 5.6, [
    "A 0 to 100 cross-sectional composite of quality, valuation and momentum pillars.",
    "Blended 60% on a three-year horizon and 40% on a one-year horizon.",
    "A Sell triggers below 40; at or above 40 the call is Hold. The model never issues Buy.",
    "TTM v7: one-year growth and P/E reflect the latest reported quarter where available.",
    "Refreshed each quarter on point-in-time data; a score can be wrong or stale.",
], fn=SERIF, sz=11, ls=1.25)
txt(s, 6.85, 1.95, 5.5, 0.24, [("BASIS OF PREPARATION", SANS, 9, SLATE, True, False, 150)])
paras(s, 6.85, 2.30, 5.5, [
    "† Fund alpha is measured to 31 December 2024, the most recent date with audited NAV histories common to the seasoned schemes. It will be refreshed in the Q3 review.",
    "‡ Acquisition dates and costs for the direct-equity holdings were not part of the statement provided; tax estimates cover the fund book only until the demat trade file is shared.",
    "Illustrative broad-market sector weights are used for over/underweight context only.",
    "All figures trace to the client statement (Jul 2026), the Ionic Quant Scorecard, or the QFRA fund NAV feed.",
], fn=SERIF, sz=11, ls=1.25)
source(s, "Source: Ionic Wealth research standards.")

# A-6 disclaimer / back
s = slide(WHITE)
txt(s, ML, 0.9, UW, 0.24, [("DISCLAIMER", SANS, 9, SLATE, True, False, 200)])
txt(s, ML, 1.3, UW, 3.3, [
    ("This document is a review of the existing holdings of the named client, prepared by Ionic Wealth on a non-discretionary basis. "
     "It is confidential and for the named client only; it must not be redistributed. It is not investment advice, not a recommendation to buy any security, and not a solicitation. "
     "The recommendations are Sell (including a full Exit or, for a fund, a Redeem), Switch, Trim or Hold on holdings already owned; the model does not issue Buy calls. "
     "Scores are model output on point-in-time data and can be wrong or stale. Past performance is not indicative of future results. "
     "Mutual fund investments are subject to market risk; read all scheme related documents carefully. "
     "Tax treatment depends on individual circumstances and may change; consult a tax adviser before dealing. "
     "Prepared exclusively for the named client.", SERIF, 11.5, INK, False)], ls=1.3)
rect(s, 0, 5.0, CW, 2.5, fill=NAVYD)
txt(s, 0, 5.9, CW, 0.5, [("IONIC WEALTH", SANS, 16, WHITE, True, False, 250)], align=PP_ALIGN.CENTER)
rule(s, (CW-2.4)/2, 6.5, 2.4, GOLD, 0.014)
txt(s, 0, 6.7, CW, 0.3, [("Private & Confidential  ·  Prepared exclusively for the named client", SANS, 8.5, ONNAVY, False)], align=PP_ALIGN.CENTER)

prs.save(OUT)
print("SAVED", OUT, "| slides", len(prs.slides._sldIdLst))
