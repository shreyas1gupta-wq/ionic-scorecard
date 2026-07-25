# -*- coding: utf-8 -*-
"""build_pr_v3.py -- PREMIUM Ionic Wealth Portfolio Review deck (v3, ground-up rebuild).
Rejects the old sample-template look. Built to the Fable art-direction spec v1.0:
private-bank register, one gold accent per page, tinted enamel pills, navy-ramp charts,
negatives in ink (never red). This build renders the 4 HERO pages for sign-off
(Cover, Snapshot, Direct Equity, Mutual Funds). Full expanded template follows approval.
Reads compact showcase_data.json. No client PII is committed (deck path is gitignored).
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# ---------------- palette (Fable spec 2.2) ----------------
NAVYD  = RGBColor(0x0F, 0x1E, 0x3D)   # cover bg
NAVY   = RGBColor(0x1F, 0x38, 0x64)   # primary
NT1    = RGBColor(0x4A, 0x65, 0x91)   # navy tint 1
NT2    = RGBColor(0x8C, 0xA3, 0xC4)   # navy tint 2
NT3    = RGBColor(0xC9, 0xD4, 0xE4)   # navy tint 3
GOLD   = RGBColor(0xC9, 0xA2, 0x27)
INK    = RGBColor(0x16, 0x23, 0x3B)
SLATE  = RGBColor(0x6B, 0x72, 0x80)
HAIR   = RGBColor(0xE5, 0xE7, 0xEB)
PANEL  = RGBColor(0xF7, 0xF8, 0xFA)
ZEBRA  = RGBColor(0xFA, 0xFB, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ONNAVY = RGBColor(0xA7, 0xB3, 0xC9)   # muted text on cover
SELL   = RGBColor(0xB2, 0x3A, 0x48)   # terracotta text/hairline
SELLBG = RGBColor(0xF6, 0xE7, 0xE9)   # sell pill fill
HOLD   = RGBColor(0x2E, 0x7D, 0x6F)   # teal text/hairline
HOLDBG = RGBColor(0xE7, 0xF1, 0xEF)   # hold pill fill
TRIM   = RGBColor(0x8A, 0x6E, 0x1B)
TRIMBG = RGBColor(0xF5, 0xEF, 0xDD)

SERIF = "Georgia"
SANS  = "Bahnschrift"

CW, CH = 13.333, 7.5
ML = 0.90
CWIDTH = 11.53

DIGEST = os.environ.get("SHOWCASE_DIGEST",
    r"C:\Users\SHREYA~1.1GU\AppData\Local\Temp\claude\c--Users-Shreyas-1Gupta-OneDrive---Angel-Broking-Limited-Desktop-Backup-NIFTY-500\5ec2bf16-8c38-4f40-9e4f-8e07be6545fd\scratchpad\showcase_data.json")
OUT = os.environ.get("SHOWCASE_OUT",
    r"C:\Users\SHREYA~1.1GU\AppData\Local\Temp\claude\c--Users-Shreyas-1Gupta-OneDrive---Angel-Broking-Limited-Desktop-Backup-NIFTY-500\5ec2bf16-8c38-4f40-9e4f-8e07be6545fd\scratchpad\PR_v3_hero.pptx")

D = json.load(open(DIGEST, encoding="utf-8"))
T = D["totals"]

prs = Presentation()
prs.slide_width  = Inches(CW)
prs.slide_height = Inches(CH)
BLANK = prs.slide_layouts[6]


# ---------------- primitives ----------------
def no_shadow(sh):
    try: sh.shadow.inherit = False
    except Exception: pass

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); no_shadow(r)
    return s

def rect(s, x, y, w, h, fill=None, line=None, lw=0.75, round_=0.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    no_shadow(shp)
    if round_:
        try: shp.adjustments[0] = round_
        except Exception: pass
    return shp

def rule(s, x, y, w, color=HAIR, h=0.01):   # thin rectangle rule (exact height, Fable 2.7)
    return rect(s, x, y, w, h, fill=color)

def txt(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if isinstance(paras[0], tuple): paras = [paras]
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(0)
        if line_spacing: p.line_spacing = line_spacing
        for run in para:
            t, fn, sz, col, bold = run[0], run[1], run[2], run[3], run[4]
            italic = run[5] if len(run) > 5 else False
            spc = run[6] if len(run) > 6 else None
            r = p.add_run(); r.text = t
            r.font.name = fn; r.font.size = Pt(sz); r.font.bold = bold
            r.font.italic = italic; r.font.color.rgb = col
            if spc is not None:
                r.font._rPr.set('spc', str(int(spc)))
    return tb

def eyebrow_title(s, eyebrow, title, tag=None):
    txt(s, ML, 0.52, 8.0, 0.25, [(eyebrow.upper(), SANS, 10, GOLD, True, False, 280)])
    rule(s, ML, 0.82, 0.55, GOLD, h=0.018)
    txt(s, ML, 0.90, 9.5, 0.55, [(title, SERIF, 26, NAVY, False)])
    if tag:
        txt(s, 9.43, 0.60, 3.0, 0.25, [(tag, SANS, 9, SLATE, False)], align=PP_ALIGN.RIGHT)

def footer(s, page):
    rule(s, ML, 7.02, CWIDTH, HAIR, h=0.012)
    txt(s, ML, 7.10, 3.0, 0.25, [("IONIC WEALTH", SANS, 8, NAVY, True, False, 250)])
    txt(s, 5.17, 7.10, 3.0, 0.25, [("Private & Confidential", SANS, 8, SLATE, False)], align=PP_ALIGN.CENTER)
    txt(s, 9.43, 7.10, 3.0, 0.25, [(f"Portfolio Review  ·  July 2026  ·  {str(page).zfill(2)}",
                                    SANS, 8, SLATE, False)], align=PP_ALIGN.RIGHT)

def pill(s, x, y, text, w=0.72):
    fill, tc = {"Sell": (SELLBG, SELL), "Hold": (HOLDBG, HOLD),
                "Switch": (SELLBG, SELL), "Trim": (TRIMBG, TRIM)}.get(text, (PANEL, SLATE))
    r = rect(s, x, y, w, 0.22, fill=fill, line=tc, lw=0.75, round_=0.5)
    txt(s, x, y - 0.006, w, 0.23, [(text.upper(), SANS, 8.5, tc, True, False, 150)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return r


def cr(v):  return f"₹ {v/1e7:.2f} Cr"
def L(v):   return f"₹ {v/1e5:,.1f} L"


# ---------------- table ----------------
def table(s, x, y, w, cols, rows, rowh=0.42, fs=10.5, zebra=False):
    """cols: (label, width, align). rows: list of cells; cell = str | ('num',t[,col]) |
    ('b',t[,col]) | ('pill',text)."""
    tot = sum(c[1] for c in cols)
    xs, ws, cx = [], [], x
    for (_, cw, _a) in cols:
        xs.append(cx); ws.append(w * cw / tot); cx += w * cw / tot
    for i, (label, _cw, al) in enumerate(cols):
        a = PP_ALIGN.RIGHT if al == "r" else (PP_ALIGN.CENTER if al == "c" else PP_ALIGN.LEFT)
        txt(s, xs[i], y, ws[i], 0.26, [(label.upper(), SANS, 9.5, SLATE, True, False, 200)],
            align=a, anchor=MSO_ANCHOR.MIDDLE)
    rule(s, x, y + 0.30, w, NAVY, h=0.015)
    ry = y + 0.34
    for ri, row in enumerate(rows):
        if zebra and ri % 2 == 1:
            rect(s, x, ry, w, rowh, fill=ZEBRA)
        for i, cell in enumerate(row):
            al = cols[i][2]
            a = PP_ALIGN.RIGHT if al == "r" else (PP_ALIGN.CENTER if al == "c" else PP_ALIGN.LEFT)
            if isinstance(cell, tuple) and cell[0] == "pill":
                pw = min(ws[i] - 0.06, 0.82)
                pill(s, xs[i] + (ws[i] - pw) / 2, ry + (rowh - 0.22) / 2, cell[1], pw)
            elif isinstance(cell, tuple):
                tag, tval = cell[0], cell[1]
                col = cell[2] if len(cell) > 2 else INK
                txt(s, xs[i], ry, ws[i], rowh, [(tval, SANS, fs, col, tag == "b")],
                    align=a, anchor=MSO_ANCHOR.MIDDLE)
            else:
                txt(s, xs[i], ry, ws[i], rowh, [(str(cell), SANS, fs, INK, False)],
                    align=a, anchor=MSO_ANCHOR.MIDDLE)
        rule(s, x, ry + rowh, w, HAIR, h=0.008)
        ry += rowh
    return ry


def donut(s, x, y, wh, pairs, colors, hole=65):
    cd = CategoryChartData()
    cd.categories = [p[0] for p in pairs]
    cd.add_series("a", [p[1] for p in pairs])
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(wh), Inches(wh), cd)
    ch = gf.chart
    ch.has_title = False; ch.has_legend = False
    ser = ch.plots[0].series[0]
    for i, pt in enumerate(ser.points):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.color.rgb = WHITE; pt.format.line.width = Pt(1.5)
    csp = ch._chartSpace
    for dc in csp.iter(qn('c:doughnutChart')):
        hs = dc.find(qn('c:holeSize'))
        if hs is None:
            hs = dc.makeelement(qn('c:holeSize'), {}); dc.append(hs)
        hs.set('val', str(hole))
    return gf

def stacked_bar(s, x, y, w, h, segs):
    tot = sum(p[1] for p in segs) or 1
    cx = x
    for i, (label, pct, col) in enumerate(segs):
        sw = w * pct / tot
        if sw <= 0.002: continue
        rect(s, cx + (0.02 if i else 0), y, sw - (0.02 if i else 0), h, fill=col)
        cx += sw
    caption = "   ·   ".join(f"{lb} {pc:.1f}%" for lb, pc, _c in segs)
    txt(s, x, y + h + 0.06, w, 0.22, [(caption, SANS, 9, SLATE, False)], align=PP_ALIGN.LEFT)

def kpi(s, x, y, w, h, label, value, sub):
    rect(s, x, y, w, h, fill=PANEL)
    rule(s, x, y, w, NAVY, h=0.02)
    txt(s, x + 0.15, y + 0.12, w - 0.3, 0.22, [(label.upper(), SANS, 9, SLATE, True, False, 120)])
    txt(s, x + 0.15, y + 0.34, w - 0.3, 0.4, [(value, SANS, 20, NAVY, True)])
    txt(s, x + 0.15, y + 0.74, w - 0.3, 0.24, [(sub, SANS, 9, SLATE, False)])

def card(s, x, y, w, h, name, call, body):
    rect(s, x, y, w, h, fill=PANEL)
    rule(s, x, y, 0.03, (SELL if call in ("Sell", "Switch") else NAVY), h=h)  # left rule
    txt(s, x + 0.18, y + 0.12, w - 1.05, 0.32, [(name, SANS, 10.5, NAVY, True)], anchor=MSO_ANCHOR.MIDDLE)
    pill(s, x + w - 0.86, y + 0.14, call, 0.72)
    txt(s, x + 0.18, y + 0.44, w - 0.36, h - 0.54, [[(body, SANS, 9.5, INK, False)]], line_spacing=1.18)


# =================================================================
# SLIDE 1 -- COVER
# =================================================================
s = slide(NAVYD)
txt(s, ML, 0.62, 6.0, 0.30, [("IONIC WEALTH", SANS, 13, WHITE, True, False, 300)])
rule(s, ML, 1.02, 0.35, GOLD, h=0.018)
txt(s, ML, 2.55, 8.0, 0.28, [("PRIVATE WEALTH  ·  PORTFOLIO REVIEW", SANS, 11, GOLD, True, False, 300)])
txt(s, ML, 2.90, 11.0, 1.05, [("Portfolio Review", SERIF, 54, WHITE, False)])
txt(s, ML, 4.05, 10.0, 0.4, [("Holdings, allocation and recommendations for the family portfolio.",
                              SERIF, 15, ONNAVY, False, True)])
txt(s, ML, 5.35, 4.0, 0.22, [("PREPARED FOR", SANS, 9, ONNAVY, True, False, 280)])
txt(s, ML, 5.62, 7.0, 0.4, [(D["client_label"], SERIF, 20, WHITE, False)])
txt(s, ML, 6.12, 6.0, 0.28, [("21 July 2026  ·  Mumbai", SANS, 10, ONNAVY, False)])
rule(s, ML, 6.78, CWIDTH, RGBColor(0x2A, 0x3B, 0x5E), h=0.012)
txt(s, ML, 6.92, 6.0, 0.25, [("Ionic Wealth  ·  Internal review of existing holdings", SANS, 8, ONNAVY, False)])
txt(s, 8.93, 6.92, 3.5, 0.25, [("Private & Confidential", SANS, 8, ONNAVY, False)], align=PP_ALIGN.RIGHT)

# =================================================================
# SLIDE 2 -- SNAPSHOT (signature)
# =================================================================
s = slide()
eyebrow_title(s, "Portfolio Snapshot", "Where the portfolio stands", tag="As of 21 July 2026")
# hero (single gold number = the one gold accent besides chrome)
txt(s, ML, 1.65, 5.0, 0.22, [("TOTAL PORTFOLIO VALUE", SANS, 9.5, SLATE, True, False, 200)])
txt(s, ML - 0.02, 1.90, 6.2, 0.85, [(cr(T["grand"]), SANS, 46, GOLD, False)])
txt(s, ML, 2.86, 6.2, 0.28, [(f"{T['n_stocks']} direct stocks and {T['n_mf']} fund schemes across two sleeves",
                              SANS, 11, SLATE, False)])
rule(s, ML, 3.34, 6.20, HAIR, h=0.008)
# 4 KPI tiles (2x2)
kpi(s, ML,   3.55, 2.95, 1.05, "Direct Equity", cr(T["equity"]), f"{T['eq_pct']}% of portfolio")
kpi(s, 4.15, 3.55, 2.95, 1.05, "Mutual Funds", cr(T["mf"]),     f"{T['mf_pct']}% of portfolio")
kpi(s, ML,   4.85, 2.95, 1.05, "Top-10 Weight", f"{T['top10_book_pct']}%", "of the whole book")
kpi(s, 4.15, 4.85, 2.95, 1.05, "Flagged To Act", f"{T['n_eq_sell']} + 3", "stocks  ·  fund schemes")
# right zone: donut + legend + mcap bar
txt(s, 7.55, 1.65, 4.88, 0.22, [("ASSET ALLOCATION", SANS, 9.5, SLATE, True, False, 200)])
donut(s, 7.55, 1.95, 2.55, [("Direct Equity", T["eq_pct"]), ("Mutual Funds", T["mf_pct"])], [NAVY, NT2])
txt(s, 7.72, 2.86, 2.20, 0.6, [[(cr(T["grand"]), SANS, 12, NAVY, True)], [("TOTAL", SANS, 8, SLATE, True, False, 150)]],
    align=PP_ALIGN.CENTER, line_spacing=1.0)
rect(s, 10.35, 2.52, 0.13, 0.13, fill=NAVY)
txt(s, 10.55, 2.44, 2.0, 0.22, [("Direct Equity", SANS, 10.5, INK, True)])
txt(s, 10.55, 2.66, 2.0, 0.22, [(f"{cr(T['equity'])}  ·  {T['eq_pct']}%", SANS, 9.5, SLATE, False)])
rect(s, 10.35, 3.22, 0.13, 0.13, fill=NT2)
txt(s, 10.55, 3.14, 2.0, 0.22, [("Mutual Funds", SANS, 10.5, INK, True)])
txt(s, 10.55, 3.36, 2.0, 0.22, [(f"{cr(T['mf'])}  ·  {T['mf_pct']}%", SANS, 9.5, SLATE, False)])
mp = D["mcap_split_pct"]
txt(s, 7.55, 4.85, 4.88, 0.22, [("MARKET-CAP MIX, DIRECT EQUITY", SANS, 9.5, SLATE, True, False, 200)])
segcol = {"Large": NAVY, "Mid": NT1, "Small": NT2, "Unknown": NT3}
segs = [(k, mp[k], segcol.get(k, NT3)) for k in ["Large", "Mid", "Small", "Unknown"] if k in mp]
stacked_bar(s, 7.55, 5.18, 4.88, 0.30, segs)
# in-brief editorial strip
rule(s, ML, 6.22, CWIDTH, HAIR, h=0.008)
txt(s, ML, 6.34, CWIDTH, 0.55,
    [("In brief:  ", SERIF, 12, NAVY, True, True),
     (f"a large-cap-heavy, aggressive core-satellite book worth {cr(T['grand'])}, split roughly evenly "
      f"between direct equity and funds. We flag {T['n_eq_sell']} direct stocks and three fund schemes to act on; "
      "the pages that follow set out each call.", SERIF, 12, INK, False, True)], line_spacing=1.15)
footer(s, 3)

# =================================================================
# SLIDE 3 -- DIRECT EQUITY
# =================================================================
s = slide()
eyebrow_title(s, "Direct Equity", "Review & recommendations", tag=f"Top 10 of {T['n_stocks']} holdings  ·  {cr(T['equity'])}")
stocks = D["top_stocks"]
by_sym = {(st.get("symbol") or st["name"]): st for st in stocks}
rows = []
for st in stocks[:10]:
    nm = (st.get("symbol") or st["name"])[:18]
    sec = (st.get("sector") or "-")[:16]
    sc = st.get("score3y")
    rows.append([("b", nm), ("num", sec) if False else sec, ("num", L(st["value"])),
                 ("num", f"{st['wt_eq']:.1f}%"),
                 ("b", (f"{sc:.0f}" if sc is not None else "-"), NAVY),
                 ("pill", st["call"] if st["call"] in ("Sell", "Hold", "Trim") else "Hold")])
cols = [("Company", 2.15, "l"), ("Sector", 1.55, "l"), ("Value", 1.05, "r"),
        ("Weight", 0.85, "r"), ("Ionic Score", 1.10, "r"), ("Call", 1.05, "c")]
endy = table(s, ML, 1.60, 7.60, cols, rows, rowh=0.42, fs=10.5)
txt(s, ML, endy + 0.08, 7.60, 0.24,
    [("Full 68-stock scorecard in Appendix A. Ionic Score is a 0 to 100 composite of quality, valuation and "
      "momentum; a Sell triggers below 40.", SANS, 8, SLATE, False)])

txt(s, 8.80, 1.60, 3.63, 0.22, [("ANALYST SPOTLIGHT", SANS, 9.5, SLATE, True, False, 200)])
def g(sym, key, default=None):
    st = by_sym.get(sym); return (st.get(key) if st else default)
def sc_str(sym):
    v = g(sym, "score3y"); return f"{v:.0f}" if v is not None else "n/a"
def wt_str(sym):
    v = g(sym, "wt_eq"); return f"{v:.1f}%" if v is not None else ""
feat = [
    ("RELIANCE", "Sell",
     f"Reliance screens Sell at {sc_str('RELIANCE')} on the Ionic scorecard, where a full valuation and soft "
     f"relative momentum outweigh the franchise quality. It is also a top single-stock position at {wt_str('RELIANCE')}. "
     "We act on the call by trimming hard toward a market weight and rotating the freed capital into higher-scoring names."),
    ("TITAN", "Hold",
     f"Titan is the largest holding at {wt_str('TITAN')} and carries a Hold at {sc_str('TITAN')}. The jewellery "
     "franchise keeps compounding and the quant profile stays above our Sell line. We retain it, while "
     "noting that a single stock at this weight is a concentration worth watching."),
    ("TCS", "Hold",
     f"TCS holds a Hold at {sc_str('TCS')}. Strong cash generation and a defensive services mix keep it in the core; "
     "near-term caution on IT spending shows up as a moderated growth reading, not a broken thesis. Retain as a core large-cap."),
]
cy = 1.90
for sym, call, body in feat:
    card(s, 8.80, cy, 3.63, 1.55, sym.title() if sym != "TCS" else "TCS", call, body)
    cy += 1.70
footer(s, 12)

# =================================================================
# SLIDE 4 -- MUTUAL FUNDS
# =================================================================
s = slide()
eyebrow_title(s, "Mutual Funds", "Scheme review", tag=f"{T['n_mf']} schemes  ·  {cr(T['mf'])}")
mf = D["mf"]
def catshort(c):
    c = (c or "").replace("Equity: ", "").replace("Hybrid: ", "")
    return {"Thematic-Infrastructure": "Thematic", "Multi Asset Allocation": "Multi Asset",
            "Large & MidCap": "Large & Mid", "Value Oriented": "Value"}.get(c, c)[:14]
def mf_action(f):
    sc = f["scheme"] or ""
    if sc == "ICICI Prudential Multi Asset Fund": return "Sell"      # costlier regular plan
    if "Nippon India Multi Cap" in sc: return "Switch"
    if "Bandhan Small Cap" in sc: return "Sell"
    return "Hold"
mf_sorted = sorted(mf, key=lambda f: -(f["value"] or 0))
rows = []
for f in mf_sorted:
    nm = (f["scheme"] or "").replace(" - Direct Plan", "").replace(" Fund", "")
    nm = nm.replace("Motilal Oswal", "MO").replace("ICICI Prudential", "ICICI Pru")[:26]
    al = f.get("alpha")
    alc = ("num", (f"{al:+.1f}%" if al is not None else "—"), INK)
    act = mf_action(f)
    rows.append([("b", nm), catshort(f["category"]), ("num", L(f["value"])), alc, ("pill", act)])
cols = [("Scheme", 2.95, "l"), ("Category", 1.35, "l"), ("Value", 1.05, "r"),
        ("3Y Alpha", 0.95, "r"), ("Call", 0.90, "c")]
endy = table(s, ML, 1.60, 7.20, cols, rows, rowh=0.27, fs=9, zebra=True)
txt(s, ML, endy + 0.07, 7.20, 0.24,
    [("3Y alpha vs each scheme's category benchmark, direct-plan NAVs to Dec 2024. Only four schemes carry a 3-year record.",
      SANS, 8, SLATE, False)])

txt(s, 8.65, 1.60, 3.78, 0.22, [("WHY WE ACT ON THESE THREE", SANS, 9.5, SLATE, True, False, 200)])
rule(s, 8.65, 1.86, 0.35, GOLD, h=0.015)
sell_notes = [
    ("Nippon India Multi Cap", "Switch",
     "A multi-cap mandate forces a fixed small and mid cap quota we would rather size ourselves. Its three-year "
     "record is genuinely strong, so this is a switch, not a criticism: move into the flexi-cap sleeve."),
    ("ICICI Pru Multi Asset (Regular)", "Sell",
     "The same multi-asset fund the family already owns in the direct plan, held here in the costlier regular "
     "plan. No reason to pay the extra trail: redeem and consolidate into the existing direct-plan holding."),
    ("Bandhan Small Cap", "Sell",
     "A sub-scale small-cap position of about ₹ 3 L sitting beside a ₹ 73 L small-cap fund. It adds cost "
     "and duplication without moving the needle. Exit and let the primary small-cap holding carry the sleeve."),
]
cy = 2.05
for nm, call, body in sell_notes:
    card(s, 8.65, cy, 3.78, 1.45, nm, call, body)
    cy += 1.57
footer(s, 14)

prs.save(OUT)
print("SAVED", OUT, "| slides", len(prs.slides._sldIdLst))
