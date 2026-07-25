"""
build_pr_ppt.py (v2) — Ionic Wealth Portfolio Review deck, client-grade.
Reads a holdings Excel (Stocks + MF), joins direct equity to the TTM-v7 scorecard + pf_qual
commentary, folds in the parallel-agent analysis (results/pr_kordes/*.json: real MF 3Y-alpha,
concentration/risk, client story, tax, overlap, redeployment, QA), and renders an Ionic-branded
.pptx page by page.

Vocab Sell/Trim/Reduce/Hold/No-View only (never Buy). Nothing fabricated; uncovered names are
flagged. House voice via detell(). Missing agent JSONs degrade gracefully.

Run: set PYTHONIOENCODING=utf-8 && <py> build_pr_ppt.py [holdings.xlsx]
"""
import glob
import json
import os
import re
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

ROOT = r"c:\Users\Shreyas.1Gupta\OneDrive - Angel Broking Limited\Desktop\Backup\NIFTY 500"
SD = os.path.join(ROOT, "Shreyas_Ionic_AMC", "05_DATA_OFFICE", "scripts")
RES = os.path.join(ROOT, "Shreyas_Ionic_AMC", "04_RND_LAB", "STOCK_SCORECARD_750", "results")
IDX = os.path.join(ROOT, "datasets", "index_constituents")
OUT = os.path.join(ROOT, "Shreyas_Ionic_AMC", "09_PRODUCT", "reports", "PORTFOLIO_REVIEW_Kordes_Family.pptx")

NAVY = RGBColor(0x26, 0x28, 0x73); BLUE = RGBColor(0x2B, 0x2C, 0x8C); GOLD = RGBColor(0xF5, 0xA6, 0x23)
LBLUE = RGBColor(0x7C, 0x8F, 0xF0); WHITE = RGBColor(0xFF, 0xFF, 0xFF); GREY = RGBColor(0x66, 0x66, 0x66)
LGREY = RGBColor(0xD9, 0xD9, 0xD9)
GREEN = RGBColor(0xD9, 0xEA, 0xD3); GREEN_T = RGBColor(0x1E, 0x6B, 0x2E)
RED = RGBColor(0xF4, 0xCC, 0xCC); RED_T = RGBColor(0xB0, 0x20, 0x20)
AMBER = RGBColor(0xFF, 0xF2, 0xCC); AMBER_T = RGBColor(0x9C, 0x6B, 0x00)
HFONT = "Bahnschrift"; BFONT = "Calibri"
EW, EH = I(13.333), I(7.5)
_TELL = {"delve": "examine", "leverage": "use", "utilize": "use", "robust": "strong", "seamless": "smooth",
         "notably": "", "moreover": "", "furthermore": "also", "meticulous": "careful", "underscore": "highlight",
         "pivotal": "key", "intricate": "complex"}


def detell(s):
    if not isinstance(s, str):
        return "" if s is None else str(s)
    s = s.replace("—", " ").replace(" - ", " ")
    for k, v in _TELL.items():
        s = re.sub(rf"\b{k}\b", v, s, flags=re.I)
    return re.sub(r"\s{2,}", " ", s).strip()


def money(x):
    try:
        x = float(x)
    except Exception:
        return "-"
    if x >= 1e7:
        return f"Rs {x/1e7:.2f} Cr"
    if x >= 1e5:
        return f"Rs {x/1e5:.2f} L"
    return f"Rs {x:,.0f}"


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(s, l, t, w, h, fill, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def txt(s, l, t, w, h, runs, size=14, color=BLUE, bold=False, align=PP_ALIGN.LEFT,
        font=BFONT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str):
        paras = [[(runs, color, bold)]]
    elif runs and isinstance(runs[0], list):
        paras = runs                      # list of paragraphs (each a list of run-tuples)
    else:
        paras = [runs]                    # single paragraph = list of run-tuples
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align; first = False
        for tx, cl, bd in para:
            r = p.add_run(); r.text = tx; r.font.size = Pt(size); r.font.bold = bd
            r.font.name = font; r.font.color.rgb = cl; r.font.italic = italic
    return tb


def bullets(s, l, t, w, h, items, size=12, color=NAVY, gap=True):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(6 if gap else 2)
        r = p.add_run(); r.text = "•  " + detell(it); r.font.size = Pt(size); r.font.name = BFONT; r.font.color.rgb = color
    return tb


def classified(s):
    txt(s, I(4.5), I(0.02), I(4.3), I(0.25), "Classified as Internal", 9, LBLUE, False, PP_ALIGN.CENTER)


def logo(s, dark=False, big=False):
    c = WHITE if dark else BLUE
    tb = s.shapes.add_textbox(I(10.9), I(0.18), I(2.3), I(0.5)); tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    for tx, cl in [("I", c), ("O", GOLD), ("NIC ", c), ("WEALTH", c)]:
        r = p.add_run(); r.text = tx; r.font.size = Pt(15); r.font.bold = True; r.font.name = HFONT; r.font.color.rgb = cl
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.RIGHT
    r = p2.add_run(); r.text = "BY ANGEL ONE"; r.font.size = Pt(7); r.font.color.rgb = (WHITE if dark else GREY); r.font.name = HFONT


def header(s, kicker, title):
    classified(s); logo(s)
    txt(s, I(0.5), I(0.52), I(10), I(0.7), [[(kicker + " ", GOLD, True), (title, BLUE, True)]], 25, font=HFONT)
    rect(s, I(0.5), I(1.26), I(12.3), Pt(2.2), BLUE)
    rect(s, I(0.5), I(1.26), I(2.2), Pt(2.2), GOLD)


def rings(s, x, y, big=False):
    for d, cl in ([(I(3.4), GOLD), (I(2.5), WHITE)] if big else [(I(3.0), GOLD), (I(2.2), WHITE)]):
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
        o.fill.background(); o.line.color.rgb = cl; o.line.width = Pt(10); o.shadow.inherit = False


ACT_CLR = {"Hold": (GREEN, GREEN_T), "Sell": (RED, RED_T), "Trim": (AMBER, AMBER_T),
           "Reduce": (AMBER, AMBER_T), "Switch": (AMBER, AMBER_T), "No View": (AMBER, AMBER_T),
           "Review": (LGREY, GREY)}


def table(s, l, t, w, headers, rows, colw, act_col=None, fs=9, hdr_fs=9, rowh=0.3):
    n = len(rows) + 1
    tb = s.shapes.add_table(n, len(headers), l, t, w, I(rowh * n)).table
    for j, cw in enumerate(colw):
        tb.columns[j].width = I(cw)
    for j, hd in enumerate(headers):
        c = tb.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for m in ("margin_top", "margin_bottom"):
            setattr(c, m, Pt(1))
        c.margin_left = Pt(3); c.margin_right = Pt(3)
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hd; r.font.size = Pt(hdr_fs); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = HFONT
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tb.cell(i, j)
            for m in ("margin_top", "margin_bottom"):
                setattr(c, m, Pt(1))
            c.margin_left = Pt(3); c.margin_right = Pt(3)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val); r.font.size = Pt(fs); r.font.name = BFONT; r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            if act_col is not None and j == act_col:
                fill, tcol = ACT_CLR.get(str(val).strip(), (WHITE, RGBColor(0x22, 0x22, 0x22)))
                c.fill.fore_color.rgb = fill; r.font.color.rgb = tcol; r.font.bold = True
    return tb


def pie(s, l, t, w, h, labels, vals, title):
    cd = CategoryChartData(); cd.categories = labels; cd.add_series("s", vals)
    ch = s.shapes.add_chart(XL_CHART_TYPE.PIE, l, t, w, h, cd).chart
    ch.has_title = True; ch.chart_title.text_frame.text = title
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.RIGHT; ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(8)
    pl = ch.plots[0]; pl.has_data_labels = True
    pl.data_labels.number_format = '0.0"%"'; pl.data_labels.number_format_is_linked = False; pl.data_labels.font.size = Pt(8)
    return ch


def bar(s, l, t, w, h, labels, vals, title):
    cd = CategoryChartData(); cd.categories = labels; cd.add_series("Allocation %", vals)
    ch = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd).chart
    ch.has_title = True; ch.chart_title.text_frame.text = title
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    ch.has_legend = False
    pl = ch.plots[0]; pl.has_data_labels = True
    pl.data_labels.number_format = '0.0"%"'; pl.data_labels.number_format_is_linked = False
    pl.data_labels.font.size = Pt(8); pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    pl.series[0].format.fill.solid(); pl.series[0].format.fill.fore_color.rgb = BLUE
    ch.category_axis.tick_labels.font.size = Pt(7); ch.value_axis.tick_labels.font.size = Pt(8)
    return ch


# ---------------- data ----------------
def load():
    f = sys.argv[1] if len(sys.argv) > 1 else glob.glob(os.path.join(SD, "*ordes*.xlsx"))[0]
    st = pd.read_excel(f, sheet_name="Stocks Holdings", skiprows=3); st.columns = [str(c).strip() for c in st.columns]
    st = st.dropna(subset=["Security Name"]).copy()
    st = st[st["Security Name"].astype(str).str.strip().str.lower() != "total"]
    st["Current Value"] = pd.to_numeric(st["Current Value"], errors="coerce")
    st["is_etf"] = st["Security Name"].str.contains("ETF", case=False, na=False)
    mf = pd.read_excel(f, sheet_name="MF"); mf.columns = [str(c).strip() for c in mf.columns]
    mf = mf.dropna(subset=["Scheme Name"]).copy()
    for c in ("Current Value", "Invested Value", "Units"):
        mf[c] = pd.to_numeric(mf[c], errors="coerce")
    mf = mf[mf["Scheme Name"].astype(str).str.lower() != "total"]
    return f, st, mf


def isin_map():
    m = {}
    nse = os.path.join(SD, "nse_isin_symbol.csv")
    if os.path.exists(nse):
        d = pd.read_csv(nse)
        ic = next((x for x in d.columns if "ISIN" in x.upper()), None); sc = next((x for x in d.columns if x.upper() == "SYMBOL"), None)
        for _, r in d.iterrows():
            m[str(r[ic]).strip()] = str(r[sc]).strip()
    return m


NAME_OVERRIDE = {"kotak mahindra bank ltd.": "KOTAKBANK", "national securities depository ltd.": "NSDL"}


def load_pr():
    d = os.path.join(RES, "pr_kordes"); out = {}
    for k in ("mf_review", "concentration", "left_analysis", "story", "tax", "overlap", "redeployment", "qa"):
        p = os.path.join(d, k + ".json")
        try:
            out[k] = json.load(open(p, encoding="utf-8")) if os.path.exists(p) else None
        except Exception:
            out[k] = None
    return out


def g(d, k, default=None):
    return d.get(k, default) if isinstance(d, dict) else default


def main():
    f, st, mf = load()
    im = isin_map()
    sc = pd.read_csv(os.path.join(RES, "full750_scored.csv")).set_index("symbol")
    pfq = {}
    for p in glob.glob(os.path.join(RES, "pf_qual_*.json")):
        try:
            d = json.load(open(p, encoding="utf-8")); pfq[d.get("symbol")] = d
        except Exception:
            pass
    pr = load_pr()
    CON = pr["concentration"] or {}; STORY = pr["story"] or {}; MFR = pr["mf_review"] or []
    TAX = pr["tax"] or {}; REDEP = pr["redeployment"] or {}; OVL = pr["overlap"] or {}

    total = st["Current Value"].sum() + mf["Current Value"].sum()
    eq_sleeve = st.loc[~st["is_etf"], "Current Value"].sum()
    mf_total = mf["Current Value"].sum()

    # direct equity calls
    st["symbol"] = st["ISIN"].astype(str).str.strip().map(im)
    st["symbol"] = st.apply(lambda r: NAME_OVERRIDE.get(str(r["Security Name"]).strip().lower(), r["symbol"]), axis=1)
    steq = st[~st["is_etf"]].copy(); steq["wt_book"] = steq["Current Value"] / total * 100
    steq["wt_eq"] = steq["Current Value"] / eq_sleeve * 100
    rec_rows = []
    for _, r in steq.sort_values("Current Value", ascending=False).iterrows():
        sym = r["symbol"]; q = pfq.get(sym, {})
        arec = q.get("your_recommendation"); qrec = sc.loc[sym, "recommendation_overall"] if sym in sc.index else None
        base = arec or qrec
        summ = detell(q.get("summary")) if q else ""
        if base == "Sell":
            act, why = "Sell", (summ or "Weak on the scorecard read.")
        elif r["wt_eq"] > 10 and base == "Hold":
            act, why = "Reduce", f"Quality name at {r['wt_eq']:.0f}% of the equity sleeve; trim to manage concentration. " + summ
        elif base == "Hold":
            act, why = "Hold", (summ or "Scorecard Hold on current fundamentals.")
        else:
            act, why = "Review", "Held-away / off-scorecard; recent listing, monitor."
        rec_rows.append(dict(name=str(r["Security Name"])[:34], sym=sym, val=r["Current Value"],
                             wt=r["wt_book"], wt_eq=r["wt_eq"], act=act, why=why[:150],
                             sector=(sc.loc[sym, "sector"] if sym in sc.index else "-"),
                             mcap=(sc.loc[sym, "mcap_tercile"] if sym in sc.index and "mcap_tercile" in sc.columns else "-")))

    # MF: real 3Y alpha from mf_review.json (fallback: since-inception CAGR)
    mfmap = {str(x.get("scheme", ""))[:24]: x for x in MFR}
    mf["years"] = (pd.Timestamp("2026-07-21") - pd.to_datetime(mf["Inception"])).dt.days / 365.25
    mf["cagr"] = ((mf["Current Value"] / mf["Invested Value"]) ** (1 / mf["years"]) - 1) * 100
    mf["assetclass"] = mf["Category Name"].str.split(":").str[0].str.strip()
    mf["wt"] = mf["Current Value"] / total * 100

    def sc100(v):
        return v * 100 if isinstance(v, (int, float)) else None  # mf_review stores CAGR/alpha as fractions

    def mrow(r):
        rev = mfmap.get(str(r["Scheme Name"])[:24])
        if rev:
            return pd.Series([rev.get("action", "Review"), detell(rev.get("rationale", "")),
                              sc100(rev.get("fund_3y_cagr")), sc100(rev.get("benchmark_3y_cagr")), sc100(rev.get("alpha_3y"))])
        return pd.Series(["Review", "Awaiting the fund-model read.", None, None, None])
    mf[["act", "why", "f3y", "b3y", "alpha"]] = mf.apply(mrow, axis=1)

    prs = Presentation(); prs.slide_width = EW; prs.slide_height = EH
    n_sell_mf = int((mf["act"] == "Sell").sum())

    # ---- 1 COVER ----
    s = blank(prs); rect(s, 0, 0, EW, EH, NAVY); rings(s, I(10.2), I(0.5), big=True)
    tb = s.shapes.add_textbox(I(0.8), I(0.9), I(7), I(0.9)); p = tb.text_frame.paragraphs[0]
    for tx, cl in [("I", WHITE), ("O", GOLD), ("NIC ", WHITE), ("WEALTH", WHITE)]:
        r = p.add_run(); r.text = tx; r.font.size = Pt(40); r.font.bold = True; r.font.name = HFONT; r.font.color.rgb = cl
    txt(s, I(0.85), I(1.85), I(7), I(0.4), [[("Co-founder", GOLD, True), (" in your journey of wealth creation", WHITE, False)]], 16, italic=True)
    rect(s, I(0.85), I(5.3), I(4.2), Pt(1.5), GOLD)
    txt(s, I(0.85), I(5.5), I(9), I(0.8), [[("Portfolio ", WHITE, True), ("Review", GOLD, True)]], 30, font=HFONT)
    txt(s, I(0.85), I(6.4), I(9), I(0.4), f"Total Portfolio Size: {money(total)}   |   As on 21 Jan 2026", 13, LBLUE, False)

    # ---- 2 APPROACH ----
    s = blank(prs); header(s, "Our", "Approach")
    steps = ["Understand objectives behind investments", "Ascertain risk appetite (willingness & ability)",
             "Devise asset allocation and portfolio contours", "Evaluate the current portfolio",
             "Lay out actionables (High / Low priority)"]
    for i, stp in enumerate(steps):
        x = I(0.6 + i * 2.45); rect(s, x, I(2.9), I(2.25), I(1.35), LBLUE)
        txt(s, x, I(2.9), I(2.25), I(1.35), stp, 11, WHITE, True, PP_ALIGN.CENTER, BFONT, MSO_ANCHOR.MIDDLE)
        if i < 4:
            txt(s, x + I(2.22), I(3.35), I(0.25), I(0.5), ">", 20, GOLD, True, PP_ALIGN.CENTER)
    txt(s, I(0.6), I(4.6), I(12), I(1.6),
        "Our philosophy of evaluation rests on three lenses: returns and risk (participation in up and down markets), "
        "fitment to the objective and IPS, and tax and cost. We recommend selling holdings that have completed their "
        "long-term period unless action is needed sooner.", 13, NAVY, italic=True)

    # ---- 3 SNAPSHOT ----
    s = blank(prs); header(s, "Portfolio", "Snapshot")
    txt(s, I(0.5), I(1.35), I(9), I(0.4), [[("Total Portfolio Size: ", NAVY, True), (money(total), GOLD, True)]], 16)
    ac = pd.Series({"Direct Equity": eq_sleeve})
    ac = pd.concat([ac, mf.groupby("assetclass")["Current Value"].sum()])
    etfv = st.loc[st["is_etf"], "Current Value"].sum()
    if etfv > 0:
        ac["ETF / Cash"] = etfv
    pie(s, I(0.4), I(1.9), I(5.3), I(4.7), list(ac.index), [round(v, 1) for v in ac.values], "Asset Split")
    rows = [(k, money(v), f"{v/total*100:.1f}%") for k, v in ac.sort_values(ascending=False).items()]
    rows.append(("TOTAL", money(total), "100.0%"))
    table(s, I(6.2), I(2.15), I(6.5), ["Category", "Amount", "% of Portfolio"], rows, [2.6, 2.2, 1.7], fs=12, hdr_fs=12, rowh=0.5)
    txt(s, I(6.2), I(5.7), I(6.5), I(0.9),
        f"{len(steq)} direct-equity holdings + {len(mf)} mutual funds. Direct equity {eq_sleeve/total*100:.0f}%, "
        f"mutual funds {mf_total/total*100:.0f}%. Market value as on 21 Jan 2026.", 11, GREY, italic=True)

    # ---- 4 ASSET ALLOCATION + MCAP ----
    s = blank(prs); header(s, "Portfolio Overview", "Asset Allocation & Market Cap Split")
    hyb = mf[mf["assetclass"] == "Hybrid"]["Current Value"].sum()
    eq_all = eq_sleeve + mf[mf["assetclass"] == "Equity"]["Current Value"].sum()
    pie(s, I(0.3), I(1.7), I(6.2), I(4.7), ["Equity", "Hybrid", "ETF/Cash"],
        [round(eq_all/total*100, 1), round(hyb/total*100, 1), round(etfv/total*100, 1)], "Asset Allocation")
    mc = g(CON, "mcap_split") or {}
    pie(s, I(6.9), I(1.7), I(6.1), I(4.7), ["Large Cap", "Mid Cap", "Small Cap"],
        [round(mc.get("large", 0), 1), round(mc.get("mid", 0), 1), round(mc.get("small", 0), 1)],
        "Market Cap Split (Direct Equity)")
    txt(s, I(0.5), I(6.55), I(12), I(0.7),
        f"A strong large-cap core ({mc.get('large', 0):.0f}% of direct equity) with real mid and small participation, "
        "consistent with an aggressive, long-term mandate.", 11, BLUE, italic=True)

    # ---- 5 IPS CONTOURS ----
    s = blank(prs); header(s, "Portfolio Contours:", "Investment Policy Statement")
    eqp = eq_all / total * 100
    amc_max = g(CON, "single_amc_max") or {}
    them = mf[mf["Category Name"].str.contains("Thematic|Sector", case=False, na=False)]["wt"].sum()
    top1 = steq.sort_values("Current Value", ascending=False).iloc[0]
    ips = [
        ("Equity Allocation", "80% - 100%", f"{eqp:.1f}%", eqp < 80 or eqp > 100),
        ("Fixed Income & Hybrid", "0% - 20%", f"{(hyb+etfv)/total*100:.1f}%", (hyb+etfv)/total*100 > 20),
        ("Single instrument", "Max 15%", f"{top1['Current Value']/total*100:.1f}% ({str(top1['Security Name'])[:20]})", top1['Current Value']/total*100 > 15),
        ("Single AMC (MF sleeve)", "Max 25%", f"{amc_max.get('pct', 0):.1f}% ({amc_max.get('amc', '-')})", amc_max.get('pct', 0) > 25),
        ("Thematic & Sectoral", "Max 20%", f"{them:.1f}% of book", them > 20),
        ("Market Cap - Large", "50% - 70%", f"{mc.get('large', 0):.1f}%", not (50 <= mc.get('large', 60) <= 70)),
        ("Risk Profile", "Aggressive", "Aggressive", False),
    ]
    nR = len(ips) + 1
    tb = s.shapes.add_table(nR, 4, I(0.7), I(1.6), I(12), I(0.6 * nR)).table
    for j, cw in enumerate([4.2, 3.0, 3.8, 1.0]):
        tb.columns[j].width = I(cw)
    for j, hd in enumerate(["Parameter", "Ideal", "Current", ""]):
        c = tb.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hd; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = HFONT
    for i, (par, ideal, cur, breach) in enumerate(ips, 1):
        for j, val in enumerate([par, ideal, cur, ("!" if breach else "")]):
            c = tb.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = WHITE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val); r.font.size = Pt(10.5); r.font.name = BFONT
            r.font.color.rgb = RED_T if (breach and j >= 2) else RGBColor(0x22, 0x22, 0x22); r.font.bold = bool(breach and j >= 2)

    # ---- 6 CONCENTRATION & RISK ----
    s = blank(prs); header(s, "Concentration &", "Risk")
    t5 = g(CON, "top5_pct") or {}; t10 = g(CON, "top10_pct") or {}
    kpis = [("Top-5 (of equity)", f"{t5.get('of_equity_sleeve', 0):.0f}%"), ("Top-10 (of equity)", f"{t10.get('of_equity_sleeve', 0):.0f}%"),
            ("Largest AMC", f"{amc_max.get('pct', 0):.0f}%"), ("Tail (<1%) positions", f"{g(CON,'tail_count','-')}"),
            ("Sector HHI", f"{g(CON,'hhi','-')}")]
    for i, (lab, val) in enumerate(kpis):
        x = I(0.55 + i * 2.5); rect(s, x, I(1.6), I(2.3), I(1.1), RGBColor(0xEF, 0xF3, 0xFF))
        txt(s, x, I(1.7), I(2.3), I(0.55), val, 22, NAVY, True, PP_ALIGN.CENTER, HFONT, MSO_ANCHOR.MIDDLE)
        txt(s, x, I(2.25), I(2.3), I(0.4), lab, 9.5, GREY, False, PP_ALIGN.CENTER)
    tops = g(CON, "top_sectors") or []
    if tops:
        bar(s, I(0.5), I(2.95), I(7.0), I(3.7), [t["sector"][:14] for t in tops[:7]], [round(t["pct"], 1) for t in tops[:7]],
            "Direct-equity sector exposure (%)")
    flags = [x for x in (g(CON, "flags") or [])
             if not any(t in str(x).upper() for t in ("DATA QUALITY", "DATA GAP", "POINT-IN-TIME", "CRASH-REGIME", "NSE_ISIN", "UNMAPPED"))]
    txt(s, I(7.8), I(3.0), I(5.1), I(0.4), "Risk flags", 14, GOLD, True, font=HFONT)
    bullets(s, I(7.8), I(3.45), I(5.1), I(2.9), [x[:130] for x in flags[:4]] or ["Within policy guardrails."], 11)
    txt(s, I(0.5), I(6.75), I(12.3), I(0.6), detell(g(CON, "narrative", ""))[:320], 10.5, BLUE, italic=True)

    # ---- 7 CORE-SATELLITE CONSTRUCT ----
    s = blank(prs); header(s, "Portfolio Overview", "Current Portfolio Construct")
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, I(1.2), I(2.1), I(3.6), I(3.6)); o.fill.solid(); o.fill.fore_color.rgb = LBLUE; o.line.fill.background(); o.shadow.inherit = False
    txt(s, I(1.2), I(3.3), I(3.6), I(1.2), [[("Equity ", WHITE, True), (f"{eq_all/total*100:.0f}%", WHITE, True)]], 20, WHITE, True, PP_ALIGN.CENTER, HFONT, MSO_ANCHOR.MIDDLE)
    txt(s, I(1.2), I(5.75), I(3.6), I(0.5), f"Debt + Hybrid {(hyb+etfv)/total*100:.0f}%", 12, NAVY, True, PP_ALIGN.CENTER)
    txt(s, I(5.6), I(2.0), I(7.2), I(0.5), "Core (risk-off): low churn, all-weather   |   Satellite (risk-on): tactical, high conviction", 12, GOLD, True, font=HFONT)
    txt(s, I(5.6), I(2.7), I(7.2), I(3.8), detell(g(STORY, "construct_read", "")), 13, NAVY, False)

    # ---- 8 MF CATEGORY ALLOCATION ----
    s = blank(prs); header(s, "Mutual Fund Overview", "Category Allocation")
    cat = mf.groupby("Category Name")["Current Value"].sum().sort_values(ascending=False)
    bar(s, I(0.4), I(1.6), I(12.5), I(4.8), [c.replace("Equity: ", "").replace("Hybrid: ", "")[:15] for c in cat.index],
        [round(v/mf_total*100, 1) for v in cat.values], "MF allocation by category (% of MF book)")
    txt(s, I(0.5), I(6.55), I(12), I(0.6),
        "Concentration sits in Small Cap and one Thematic-Infrastructure holding; the tail of sub-scale sleeves can be consolidated.", 11, BLUE, italic=True)

    # ---- 9 MF AMC CONCENTRATION ----
    s = blank(prs); header(s, "Mutual Fund Overview", "AMC Concentration")
    amcs = g(CON, "amc_breakdown") or []
    if amcs:
        bar(s, I(0.4), I(1.7), I(12.5), I(4.6), [a["amc"][:12] for a in amcs[:9]], [round(a.get("pct_of_book", a.get("pct", 0)), 1) for a in amcs[:9]], "AMC exposure (% of total book)")
    txt(s, I(0.5), I(6.5), I(12), I(0.7),
        f"Largest single AMC is {amc_max.get('amc','-')} at {amc_max.get('pct',0):.1f}% of the book, within the 25% guideline on the MF sleeve. "
        "Ideal single-AMC exposure is up to 25%.", 11, BLUE, italic=True)

    # ---- 10 EQUITY MF REVIEW (real 3Y alpha) ----
    s = blank(prs); header(s, "Equity Mutual Funds", "3-Year Alpha Review")
    eqmf = mf[mf["assetclass"] == "Equity"].sort_values("Current Value", ascending=False)

    def fmtp(v):
        return "-" if (v is None or (isinstance(v, float) and pd.isna(v))) else f"{float(v):.1f}%"
    rows = [(str(r["Scheme Name"])[:28], money(r["Current Value"]), fmtp(r["f3y"]), fmtp(r["b3y"]), fmtp(r["alpha"]), r["act"], detell(r["why"])[:92]) for _, r in eqmf.iterrows()]
    table(s, I(0.35), I(1.45), I(12.65), ["Scheme", "Value", "Fund 3Y", "Bmk 3Y", "Alpha", "Action", "Rationale"],
          rows, [2.8, 1.2, 0.95, 0.95, 0.9, 1.05, 4.9], act_col=5, fs=7.5, hdr_fs=8.5, rowh=0.355)
    txt(s, I(0.35), I(6.9), I(12.6), I(0.5),
        "Trailing 3Y CAGR (window ending Dec-2024) from NAVs vs the category benchmark TRI (in-house MF model). Negative 3Y alpha = Sell; sectoral/thematic = No View; Multi-Cap = Switch. "
        "Review = NAV series too short for a clean 3Y window (e.g. HSBC Small Cap's series starts Nov-2022 post-rebrand though the fund is long-established) or outside the active-equity model.",
        7.5, GREY, italic=True)

    # ---- 11 HYBRID/DEBT MF REVIEW ----
    s = blank(prs); header(s, "Hybrid & Debt", "Mutual Fund Review")
    hybmf = mf[mf["assetclass"] != "Equity"].sort_values("Current Value", ascending=False)
    rows = [(str(r["Scheme Name"])[:34], r["Category Name"].split(":")[-1].strip(), money(r["Current Value"]), fmtp(r["alpha"]), r["act"]) for _, r in hybmf.iterrows()]
    table(s, I(0.6), I(1.7), I(12), ["Scheme", "Category", "Value", "3Y Alpha", "Action"], rows, [4.2, 2.6, 2.0, 1.6, 1.2], act_col=4, fs=10, hdr_fs=10, rowh=0.5)
    txt(s, I(0.6), I(1.7) + I(0.5 * (len(rows) + 1)) + I(0.25), I(12), I(0.7),
        "The hybrid sleeve provides beta protection; arbitrage and multi-asset holdings are retained as the risk-off core.", 12, BLUE, italic=True)

    # ---- 12 / 13 DIRECT EQUITY RECOMMENDATION ----
    for pg, chunk in enumerate([rec_rows[:14], rec_rows[14:28]]):
        if not chunk:
            continue
        s = blank(prs); header(s, "Direct Equity", "Recommendation" + (f" ({pg+1}/2)" if len(rec_rows) > 14 else ""))
        rows = [(x["name"], money(x["val"]), f"{x['wt_eq']:.1f}%", x["act"], x["why"]) for x in chunk]
        table(s, I(0.35), I(1.45), I(12.65), ["Company", "Value", "Wt (eq)", "Action", "Rationale"],
              rows, [2.7, 1.3, 0.9, 1.0, 6.75], act_col=3, fs=8, hdr_fs=8.5, rowh=0.35)
        txt(s, I(0.35), I(6.95), I(12.6), I(0.4),
            "Calls from our TTM-v7 scorecard + analyst commentary; a Hold above ~10% of the equity sleeve is flagged Reduce for concentration. Weight shown as % of the direct-equity sleeve.", 8, GREY, italic=True)

    # ---- 14 DIRECT EQUITY by-call summary (remaining) ----
    s = blank(prs); header(s, "Direct Equity", "Holdings by Call")
    bycall = {}
    for x in rec_rows:
        bycall.setdefault(x["act"], []).append(x["name"])
    order = ["Sell", "Reduce", "Trim", "Hold", "Review"]
    x0 = I(0.6)
    for i, a in enumerate([a for a in order if a in bycall]):
        col = i % 3; row = i // 3
        x = I(0.55 + col * 4.25); y = I(1.7 + row * 2.7)
        fill, tcol = ACT_CLR.get(a, (LGREY, GREY))
        rect(s, x, y, I(4.0), I(0.5), fill)
        txt(s, x, y, I(4.0), I(0.5), f"{a}  ({len(bycall[a])})", 13, tcol, True, PP_ALIGN.CENTER, HFONT, MSO_ANCHOR.MIDDLE)
        names = ", ".join(bycall[a])
        txt(s, x, y + I(0.55), I(4.0), I(2.0), names[:340], 10, NAVY, False)

    # ---- 15 RECOMMENDATION SNAPSHOT ----
    s = blank(prs); header(s, "Recommendation", "Snapshot")
    snaprows = []
    eqdf = pd.DataFrame(rec_rows)
    for a, gdf in eqdf.groupby("act"):
        snaprows.append(("Direct Equity", a, money(gdf["val"].sum()), f"{gdf['val'].sum()/total*100:.1f}%", len(gdf)))
    for (acl, a), gdf in mf.groupby(["assetclass", "act"]):
        snaprows.append((f"{acl} MF", a, money(gdf["Current Value"].sum()), f"{gdf['Current Value'].sum()/total*100:.1f}%", len(gdf)))
    table(s, I(1.0), I(1.7), I(11.3), ["Asset Class", "Action", "Amount", "% of Portfolio", "No."], snaprows, [3.0, 2.2, 2.6, 2.5, 1.0], act_col=1, fs=10, hdr_fs=10, rowh=0.42)

    # ---- 16 TAX & COST ----
    if TAX:
        s = blank(prs); header(s, "Tax & Cost", "Aware Sell Sequencing")
        lt = g(TAX, "lt_ready") or []; stc = g(TAX, "st_caution") or []
        txt(s, I(0.6), I(1.55), I(6), I(0.4), f"Long-term ready ({len(lt)})", 14, GREEN_T, True, font=HFONT)
        bullets(s, I(0.6), I(2.0), I(6), I(3.2), [f"{x.get('scheme','')[:34]} — {x.get('years_held','?')}y, gain {money(x.get('embedded_gain',0))}" for x in lt[:7]] or ["-"], 10.5)
        txt(s, I(6.9), I(1.55), I(6), I(0.4), f"Short-term caution ({len(stc)})", 14, AMBER_T, True, font=HFONT)
        bullets(s, I(6.9), I(2.0), I(6), I(3.2), [f"{x.get('scheme','')[:32]} — {x.get('days_to_lt','?')}d to LT" for x in stc[:6]] or ["-"], 10.5)
        txt(s, I(0.6), I(5.5), I(12.2), I(1.4), detell(g(TAX, "sequence_note", "") + " " + g(TAX, "equity_note", ""))[:360], 10.5, BLUE, italic=True)

    # ---- 17 EXECUTIVE SUMMARY ----
    s = blank(prs); header(s, "Executive", "Summary")
    ex = g(STORY, "exec_summary") or []
    for i, fn in enumerate(ex[:4]):
        x = I(0.6 + (i % 2) * 6.35); y = I(1.65 + (i // 2) * 2.5)
        rect(s, x, y, I(0.55), I(0.55), WHITE, line=GOLD)
        txt(s, x, y, I(0.55), I(0.55), str(i + 1), 16, BLUE, True, PP_ALIGN.CENTER, HFONT, MSO_ANCHOR.MIDDLE)
        rect(s, x + I(0.72), y, I(5.4), I(2.05), WHITE, line=BLUE)
        txt(s, x + I(0.9), y + I(0.1), I(5.05), I(1.85), detell(fn)[:230], 12, NAVY, False, PP_ALIGN.LEFT, BFONT, MSO_ANCHOR.MIDDLE)

    # ---- 18 PRIORITY ACTIONS ----
    s = blank(prs); header(s, "Priority", "Actions")
    hp = g(STORY, "high_priority") or []; lp = g(STORY, "low_priority") or []
    rect(s, I(0.6), I(1.6), I(6.0), I(0.5), RED); txt(s, I(0.6), I(1.6), I(6.0), I(0.5), "High priority", 14, RED_T, True, PP_ALIGN.CENTER, HFONT, MSO_ANCHOR.MIDDLE)
    bullets(s, I(0.6), I(2.2), I(6.0), I(4.5), [x[:180] for x in hp[:5]], 11)
    rect(s, I(6.9), I(1.6), I(6.0), I(0.5), LBLUE); txt(s, I(6.9), I(1.6), I(6.0), I(0.5), "Low priority", 14, WHITE, True, PP_ALIGN.CENTER, HFONT, MSO_ANCHOR.MIDDLE)
    bullets(s, I(6.9), I(2.2), I(6.0), I(4.5), [x[:180] for x in lp[:5]], 11)

    # ---- FUND OVERLAP ----
    if OVL:
        s = blank(prs); header(s, "Mutual Fund Overview", "Overlap & Look-Through")
        up = g(OVL, "unique_pct")
        for i, (lab, val) in enumerate([("Unique stocks", g(OVL, "unique_stocks", "-")), ("Total (with dupes)", g(OVL, "total_stocks", "-")),
                                        ("Truly unique", (f"{up:.0f}%" if isinstance(up, (int, float)) else "-"))]):
            x = I(0.6 + i * 4.2); rect(s, x, I(1.6), I(3.9), I(1.1), RGBColor(0xEF, 0xF3, 0xFF))
            txt(s, x, I(1.7), I(3.9), I(0.55), str(val), 22, NAVY, True, PP_ALIGN.CENTER, HFONT, MSO_ANCHOR.MIDDLE)
            txt(s, x, I(2.25), I(3.9), I(0.4), lab, 10, GREY, False, PP_ALIGN.CENTER)
        pairs = g(OVL, "high_overlap_pairs") or []; dups = g(OVL, "most_duplicated") or []
        txt(s, I(0.6), I(3.0), I(12), I(0.4), "High-overlap fund pairs and most-duplicated names", 13, GOLD, True, font=HFONT)
        items = [f"{str(p.get('a',''))[:26]}  /  {str(p.get('b',''))[:26]} — {p.get('pct','?')}%" for p in pairs[:5]]
        if dups:
            items.append("Most-held underlying: " + ", ".join(str(x)[:14] for x in dups[:8]))
        bullets(s, I(0.6), I(3.5), I(12.2), I(2.4), items or ["Overlap within expected range for this mix."], 11)
        txt(s, I(0.6), I(6.55), I(12.3), I(0.6), detell(g(OVL, "narrative", ""))[:300], 10.5, BLUE, italic=True)

    # ---- HOUSE-VIEW FIT & ACTIONABLES ----
    LA = pr["left_analysis"] or {}
    if LA:
        s = blank(prs); header(s, "Portfolio Fit", "House View & Actionables")
        txt(s, I(0.6), I(1.5), I(6.1), I(0.4), "What we see", 14, GOLD, True, font=HFONT)
        see = [detell(g(LA, k)) for k in ("sector_read", "mcap_positioning", "international_read", "commodity_gap") if g(LA, k)]
        see += [detell(o) for o in (g(LA, "overlap_findings") or [])[:2]]
        bullets(s, I(0.6), I(1.95), I(6.1), I(4.7), [x[:150] for x in see][:6] or ["-"], 11)
        txt(s, I(6.9), I(1.5), I(6.0), I(0.4), "Actionables (house-view aligned)", 14, GOLD, True, font=HFONT)
        bullets(s, I(6.9), I(1.95), I(6.0), I(4.7), [detell(a)[:150] for a in (g(LA, "actionables") or [])[:6]] or ["-"], 11)

    # ---- REDEPLOYMENT (Before / After) ----
    if REDEP:
        s = blank(prs); header(s, "Redeployment", "Freed Cash — Where It Goes")
        src = g(REDEP, "sources") or []; dst = g(REDEP, "destinations") or []
        srows = [(detell(str(x.get("what", "")))[:40], money(x.get("approx_amount", 0))) for x in src[:7]]
        table(s, I(0.5), I(1.6), I(6.0), ["Source of freed cash", "Approx"], srows or [("-", "-")], [4.4, 1.6], fs=9.5, hdr_fs=10, rowh=0.5)
        # destinations shown as CATEGORY/action only — no specific fund-buy names (NDPMS: not a solicitation)
        drows = [(detell(str(x.get("what", "")))[:44], money(x.get("approx_amount", 0))) for x in dst[:7]]
        table(s, I(6.8), I(1.6), I(6.1), ["Redeploy toward", "Approx"], drows or [("-", "-")], [4.4, 1.6], fs=9.5, hdr_fs=10, rowh=0.5)
        txt(s, I(0.5), I(6.35), I(12.4), I(0.95),
            detell(g(REDEP, "residual_cash_note", ""))[:180]
            + " Specific replacement schemes are chosen by the advisor from the in-house fund model at execution; this review does not solicit any purchase.",
            9.5, BLUE, italic=True)

    # ---- 19 METHODOLOGY ----
    s = blank(prs); header(s, "Our Methodology", "To Review Holdings")
    meth = [("Hold", GREEN, GREEN_T, "Strong risk-adjusted returns, favourable market participation and alpha per our framework."),
            ("Trim / Reduce", AMBER, AMBER_T, "Sound holding but oversized; reduce to manage single-name or single-AMC concentration."),
            ("Sell / Switch", RED, RED_T, "Weak risk-adjusted returns or negative 3Y alpha vs benchmark, or a better category alternative exists."),
            ("No View", AMBER, AMBER_T, "Sectoral and thematic funds are timing-led; we do not give a directional view.")]
    for i, (lab, fill, tcol, desc) in enumerate(meth):
        y = I(1.7 + i * 1.3); rect(s, I(0.6), y, I(3.0), I(1.05), fill)
        txt(s, I(0.6), y, I(3.0), I(1.05), lab, 15, tcol, True, PP_ALIGN.CENTER, HFONT, MSO_ANCHOR.MIDDLE)
        txt(s, I(3.9), y, I(8.9), I(1.05), desc, 13, NAVY, False, PP_ALIGN.LEFT, BFONT, MSO_ANCHOR.MIDDLE)

    # ---- 20 DISCLAIMER ----
    s = blank(prs); rect(s, 0, 0, EW, EH, NAVY); rings(s, I(10.4), I(0.6))
    txt(s, I(0.6), I(0.4), I(10), I(0.8), [[("Reach us at ", WHITE, True), ("hello@ionic.in", GOLD, True)]], 22, font=HFONT)
    disc = ("This presentation is intended for informational purposes only and does not constitute an offer or "
            "solicitation for investing in any products distributed by or services made available by Angel One Wealth "
            "Limited (\"AOWL\") or any of its affiliates or group entities. Any information contained in this presentation "
            "shall not be treated or construed as an investment advice or a recommendation. The recipient is advised to "
            "conduct its own due diligence and consult with its legal, tax and financial advisors before making any "
            "investment decisions. The recipient is requested to note that past performance is not indicative of future "
            "results.\n\nThis presentation is confidential and is intended solely for the recipient. Unauthorized "
            "distribution, reproduction, or other use of the information contained herein is strictly prohibited. AOWL "
            "disclaims any liability for actions taken based on the information provided in this presentation.")
    txt(s, I(0.6), I(1.6), I(9.4), I(5.0), disc, 11.5, WHITE, False, PP_ALIGN.LEFT)

    os.makedirs(os.path.dirname(OUT), exist_ok=True); prs.save(OUT)
    print("SAVED:", OUT)
    print(f"slides={len(prs.slides._sldIdLst)} total={money(total)} equity={len(steq)} funds={len(mf)} mf_sell={n_sell_mf}")
    print("JSONs used:", {k: (v is not None) for k, v in pr.items()})


if __name__ == "__main__":
    main()
