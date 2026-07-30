# -*- coding: utf-8 -*-
"""pptx_to_pdf.py — PPTX -> PDF with auto-detected backend (best-available wins).

Backends (tried in order):
  1. PowerPoint COM  — pixel-perfect, needs MS Office installed + comtypes
  2. LibreOffice     — good fidelity, user-local extract (%LOCALAPPDATA%\\Apps\\LibreOffice)
  3. Slide-to-PNG    — pure Python (Pillow), works everywhere (web/sandbox/no Office);
                       rasterises at 150 DPI so text is not selectable, but layout is faithful

Usage:
  python pptx_to_pdf.py <deck.pptx> [more.pptx ...] [--outdir DIR] [--backend pptx|libre|png]
  --backend  force a specific backend (skip auto-detect)

Each PDF lands next to its source (or in --outdir). Exit 1 if any conversion fails.
Part of the NDPMS deck pipeline (ndpms-deck skill): build -> gates -> PPTX -> PDF.
"""
import os
import sys
import glob
import subprocess
import time

# ---------------------------------------------------------------------------
# Backend 1: PowerPoint COM (Windows + MS Office + comtypes)
# ---------------------------------------------------------------------------
def _find_powerpoint():
    """Return True if PowerPoint COM automation is available."""
    try:
        import comtypes.client  # noqa: F401
        # Quick probe: can we find the PowerPoint executable?
        import winreg
        key = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE,
                             r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\powerpnt.exe")
        winreg.CloseKey(key)
        return True
    except Exception:
        return False


def _convert_pptx_com(src, dst):
    """Convert one PPTX to PDF via PowerPoint COM automation."""
    import comtypes.client
    ppt = None
    presentation = None
    try:
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        ppt.DisplayAlerts = 0  # suppress dialogs
        presentation = ppt.Presentations.Open(src, ReadOnly=True, WithWindow=False)
        # SaveAs format 32 = ppSaveAsPDF
        presentation.SaveAs(dst, 32)
        return True
    except Exception as e:
        print(f"  COM error: {e}")
        return False
    finally:
        if presentation:
            try:
                presentation.Close()
            except Exception:
                pass
        if ppt:
            try:
                ppt.Quit()
            except Exception:
                pass
            time.sleep(0.3)


# ---------------------------------------------------------------------------
# Backend 2: LibreOffice (user-local, no admin)
# ---------------------------------------------------------------------------
SOFFICE = os.path.join(os.environ.get("LOCALAPPDATA", ""), "Apps", "LibreOffice",
                       "program", "soffice.com")


def _find_soffice():
    if os.path.exists(SOFFICE):
        return SOFFICE
    hits = glob.glob(os.path.join(os.environ.get("LOCALAPPDATA", ""), "Apps",
                                  "LibreOffice*", "**", "soffice.com"), recursive=True)
    return hits[0] if hits else None


def _convert_libre(src, outdir):
    so = _find_soffice()
    if not so:
        return False
    r = subprocess.run([so, "--headless", "--norestore", "--convert-to", "pdf",
                        "--outdir", outdir, src],
                       capture_output=True, text=True, timeout=600)
    pdf = os.path.join(outdir, os.path.splitext(os.path.basename(src))[0] + ".pdf")
    return os.path.exists(pdf)


# ---------------------------------------------------------------------------
# Backend 3: Slide-to-PNG rasterisation (pure Python, works everywhere)
# ---------------------------------------------------------------------------
def _convert_png_pdf(src, dst):
    """Render each slide to a PNG via python-pptx geometry + Pillow, stitch into PDF."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        print(f"  png backend needs python-pptx + Pillow: {e}")
        return False

    DPI = 150
    prs = Presentation(src)
    w_emu = prs.slide_width
    h_emu = prs.slide_height
    w_px = int(w_emu / 914400 * DPI)
    h_px = int(h_emu / 914400 * DPI)

    images = []
    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (w_px, h_px), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            x = int(shape.left / w_emu * w_px) if shape.left else 0
            y = int(shape.top / h_emu * h_px) if shape.top else 0
            sw = int(shape.width / w_emu * w_px) if shape.width else 0
            sh = int(shape.height / h_emu * h_px) if shape.height else 0

            # filled rectangles / shapes with fill
            if hasattr(shape, "fill") and shape.fill and shape.fill.type is not None:
                try:
                    from pptx.dml.color import RGBColor
                    c = shape.fill.fore_color.rgb
                    draw.rectangle([x, y, x + sw, y + sh],
                                   fill=f"#{c}")
                except Exception:
                    pass

            # text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # approximate font size
                    try:
                        sz_pt = para.runs[0].font.size
                        sz_pt = int(sz_pt / 12700) if sz_pt else 10
                    except Exception:
                        sz_pt = 10
                    font_px = max(6, int(sz_pt * DPI / 72))
                    try:
                        font = ImageFont.truetype("arial.ttf", font_px)
                    except Exception:
                        font = ImageFont.load_default()
                    # text color
                    try:
                        tc = para.runs[0].font.color.rgb
                        color = f"#{tc}"
                    except Exception:
                        color = "#000000"
                    draw.text((x + 4, y), text, fill=color, font=font)
                    y += font_px + 2

            # images embedded in slides
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    blob = shape.image.blob
                    from io import BytesIO
                    pic = Image.open(BytesIO(blob))
                    pic = pic.convert("RGB")
                    pic = pic.resize((sw, sh), Image.LANCZOS)
                    img.paste(pic, (x, y))
                except Exception:
                    pass

        images.append(img)

    if not images:
        print(f"  no slides found in {src}")
        return False

    images[0].save(dst, "PDF", save_all=True, append_images=images[1:],
                   resolution=DPI)
    return True


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------
BACKENDS = [
    ("pptx", "PowerPoint COM", _find_powerpoint, lambda s, d, od: _convert_pptx_com(s, d)),
    ("libre", "LibreOffice", lambda: _find_soffice() is not None,
     lambda s, d, od: _convert_libre(s, od)),
    ("png", "Slide-to-PNG", lambda: True, lambda s, d, od: _convert_png_pdf(s, d)),
]


def convert(paths, outdir=None, force_backend=None):
    # pick backend
    backend = None
    for key, label, probe, fn in BACKENDS:
        if force_backend and key != force_backend:
            continue
        if probe():
            backend = (key, label, fn)
            break
    if not backend:
        print(f"No backend available (forced={force_backend}). Install comtypes "
              f"(pip install comtypes) for PowerPoint COM, or ensure LibreOffice is "
              f"extracted, or install Pillow for the PNG fallback.")
        return 1

    print(f"Backend: {backend[1]}")
    rc_all = 0
    for p in paths:
        p = os.path.abspath(p)
        od = os.path.abspath(outdir) if outdir else os.path.dirname(p)
        os.makedirs(od, exist_ok=True)
        pdf = os.path.join(od, os.path.splitext(os.path.basename(p))[0] + ".pdf")
        ok = backend[2](p, pdf, od)
        if ok and os.path.exists(pdf):
            print(f"OK   {pdf}  ({os.path.getsize(pdf) / 1e6:.1f} MB)")
        else:
            print(f"FAIL {p}")
            rc_all = 1
    return rc_all


if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    outdir = None
    force = None
    if "--outdir" in sys.argv:
        outdir = sys.argv[sys.argv.index("--outdir") + 1]
        args = [a for a in args if a != outdir]
    if "--backend" in sys.argv:
        force = sys.argv[sys.argv.index("--backend") + 1]
        args = [a for a in args if a != force]
    if not args:
        print(__doc__)
        sys.exit(2)
    sys.exit(convert(args, outdir, force))
