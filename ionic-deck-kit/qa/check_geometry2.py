# -*- coding: utf-8 -*-
"""check_geometry2.py — layout QA v2. Unlike v1 (box rectangles), this simulates the RENDERED
text extent: per-font char widths, wrap simulation per paragraph, effective bbox = box grown to
the simulated text height (and width when word-wrap is off). Finds the overlaps PowerPoint shows
that box-only checks miss.
Flags: (1) effective-text vs effective-text overlap, (2) text spilling past slide/footer bounds,
(3) text under a picture's ink area, (4) severe in-box overflow (clipping risk).
Usage: python check_geometry2.py <deck.pptx> [--json out.json]
"""
import sys, json, math
from pptx import Presentation

IN = 914400.0
SW, SH = 13.333, 7.5
# rough average char width in inches per pt of font size
CHAR_W = {"Georgia": 0.0102, "Bahnschrift": 0.0075}
DEF_W = 0.0088


def eff_box(sh):
    """(x0,y0,x1,y1) grown to simulated text extent."""
    try:
        x0, y0 = sh.left / IN, sh.top / IN
        w, h = sh.width / IN, sh.height / IN
    except Exception:
        return None
    tf = sh.text_frame
    wrap = tf.word_wrap is not False
    total_h, max_w = 0.0, 0.0
    for p in tf.paragraphs:
        chars, size, line_w = 0, 0.0, 0.0
        for r in p.runs:
            t = r.text or ""
            chars += len(t)
            sz = r.font.size.pt if r.font.size else 10
            size = max(size, sz)
            cw = (CHAR_W.get(r.font.name, DEF_W) * sz
                  * (1.05 if r.font.bold else 1.0) * (1.03 if r.font.italic else 1.0))
            line_w += len(t) * cw          # per-run width sum — mixed serif/sans measured correctly
        if chars == 0:
            total_h += 0.04
            continue
        line_h = size / 72.0 * 1.22
        if wrap and w > 0.05:
            lines = max(1, math.ceil(line_w / w))
            total_h += lines * line_h
            max_w = max(max_w, min(line_w, w))
        else:
            total_h += line_h
            max_w = max(max_w, line_w)
    return (x0, y0, x0 + max(w, max_w if not wrap else w), y0 + max(h, total_h))


def raw_box(sh):
    try:
        return (sh.left / IN, sh.top / IN, (sh.left + sh.width) / IN, (sh.top + sh.height) / IN)
    except Exception:
        return None


def inter(a, b):
    x0, y0 = max(a[0], b[0]), max(a[1], b[1])
    x1, y1 = min(a[2], b[2]), min(a[3], b[3])
    return max(0.0, x1 - x0) * max(0.0, y1 - y0)


def area(a):
    return max(0.0, (a[2] - a[0]) * (a[3] - a[1]))


def txt_of(sh):
    try:
        return sh.text_frame.text.strip() if sh.has_text_frame else ""
    except Exception:
        return ""


def check(path):
    prs = Presentation(path)
    F = []
    for i, s in enumerate(prs.slides, 1):
        texts = []
        pics = []
        for sh in s.shapes:
            if sh.shape_type == 13:
                b = raw_box(sh)
                if b:  # trim ~7% padding around chart PNGs (transparent margins)
                    dx, dy = (b[2] - b[0]) * 0.07, (b[3] - b[1]) * 0.07
                    pics.append((b[0] + dx, b[1] + dy, b[2] - dx, b[3] - dy))
            elif sh.has_text_frame:
                t = txt_of(sh)
                if t:
                    eb = eff_box(sh)
                    rb = raw_box(sh)
                    if eb and rb:
                        texts.append((t, eb, rb))
        for a in range(len(texts)):
            ta, ea, ra = texts[a]
            # 2: spill past slide bounds / into the footer band.
            #
            # The exemption here used to be the literal string "portfolio review", which coupled
            # this gate to one product's footer text. Renaming the footer for the QFRA-2 deck made
            # the gate emit 18 phantom spill-bounds findings on a deck with no defect - and the
            # danger is the reverse: someone "fixes" the non-defect by moving real chrome.
            # Exempt by ROLE instead: page chrome is deliberately PLACED in the footer band, so its
            # RAW top already sits at or below the band line. Content that spills INTO the band
            # starts above it and grows down, which is the case we actually want to catch.
            in_footer_band_by_design = ra[1] >= 7.10
            if (ea[2] > SW - 0.02 or ea[3] > 7.10) and not in_footer_band_by_design:
                F.append({"slide": i, "kind": "spill-bounds", "text": ta[:55],
                          "eff": [round(v, 2) for v in ea]})
            # 4: severe in-box vertical overflow (clipping)
            if (ea[3] - ra[3]) > 0.22 and len(ta) > 40:
                F.append({"slide": i, "kind": "clip-risk", "text": ta[:55],
                          "extra_in": round(ea[3] - ra[3], 2)})
            # 1: effective text-vs-text
            for b_ in range(a + 1, len(texts)):
                tb, ebb, rbb = texts[b_]
                if len(ta) < 12 or len(tb) < 12:
                    continue
                ov = inter(ea, ebb)
                small = min(area(ea), area(ebb))
                if small > 0 and ov / small > 0.22 and ov > 0.02:
                    F.append({"slide": i, "kind": "text-overlap", "pct": round(ov / small * 100),
                              "a": ta[:40], "b": tb[:40]})
            # 3: text under picture ink
            for pb in pics:
                ov = inter(ea, pb)
                if area(ea) > 0 and ov / area(ea) > 0.30 and len(ta) > 12:
                    F.append({"slide": i, "kind": "under-image", "text": ta[:45],
                              "pct": round(ov / area(ea) * 100)})
    return F


if __name__ == "__main__":
    path = sys.argv[1]
    fs = check(path)
    by = {}
    for f in fs:
        by.setdefault(f["kind"], []).append(f)
    print(f"{path}: {len(fs)} findings")
    for k, v in sorted(by.items()):
        print(f"  {k}: {len(v)}")
        for f in v[:15]:
            print("   ", {kk: vv for kk, vv in f.items() if kk != "kind"})
    if "--json" in sys.argv:
        json.dump(fs, open(sys.argv[sys.argv.index("--json") + 1], "w", encoding="utf-8"), indent=1)
